"""PPT 生成插件：让 AI 用 python-pptx 正确生成 PPT（避免小模型手写 XML 导致中文乱码）。
模型只需要提供标题和每页要点，服务器生成合法 .pptx 并返回下载链接。"""

import os
import re
import time

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FILES_DIR = os.path.join(BASE, "files")
os.makedirs(FILES_DIR, exist_ok=True)


def _safe(s, n=60):
    s = re.sub(r"[\\/:*?\"<>|]", "_", str(s or "演示文稿"))
    return s[:n] or "演示文稿"


def create_pptx(args):
    """用 python-pptx 生成演示文稿（编码可靠；支持多布局/主题色/副标题/要点层级/图片嵌入）。"""
    title = str(args.get("title", "")).strip() or "演示文稿"
    subtitle = str(args.get("subtitle", "")).strip() or ""
    theme = str(args.get("theme", "")).strip().lower() or "blue"
    layout = str(args.get("layout", "")).strip().lower() or "auto"   # auto/corporate/clean/dark
    if layout == "auto":
        layout = "corporate"
    slides = args.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return {"error": "slides 参数不能为空：请提供每页 {title, bullets:[...]} 列表"}
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return {"error": "服务器缺少 python-pptx，无法生成 PPT"}

    def _resolve_image(src):
        """图片来源：data URL（base64）/ http(s) URL / 本地路径 / files 相对路径 → 本地文件路径"""
        s = str(src or "").strip().strip('"')
        if not s:
            return None
        # data URL（截图工具常返回 base64）
        if s.startswith("data:image/"):
            import base64 as _b64
            m = re.match(r"^data:image/(\w+);base64,(.+)$", s, re.S)
            if not m:
                return None
            ext = "png" if m.group(1) in ("png", "jpeg", "jpg", "gif", "bmp") else "png"
            tmp = os.path.join(FILES_DIR, f"_img_{int(time.time() * 1000)}.png")
            try:
                raw = _b64.b64decode(m.group(2))
                if len(raw) < 1024 or len(raw) > 20_000_000:
                    return None
                with open(tmp, "wb") as f:
                    f.write(raw)
                return tmp
            except Exception:
                return None
        if s.startswith("http://") or s.startswith("https://"):
            import urllib.request
            tmp = os.path.join(FILES_DIR, f"_img_{int(time.time() * 1000)}_{abs(hash(s)) % 100000}.png")
            try:
                req = urllib.request.Request(s, headers={"User-Agent": "Mozilla/5.0"})
                with urllib.request.urlopen(req, timeout=20) as resp:
                    data = resp.read()
                if len(data) < 1024 or len(data) > 20_000_000:
                    return None
                with open(tmp, "wb") as f:
                    f.write(data)
                return tmp
            except Exception:
                return None
        p = s
        if not os.path.isfile(p):
            alt = os.path.join(FILES_DIR, s)
            if os.path.isfile(alt):
                p = alt
        return p if os.path.isfile(p) else None
    # 主题色板（设计感）
    THEMES = {
        "blue":   (RGBColor(0x1F, 0x4E, 0x79), RGBColor(0xD6, 0xE4, 0xF0)),
        "green":  (RGBColor(0x1B, 0x5E, 0x20), RGBColor(0xE2, 0xEF, 0xDA)),
        "red":    (RGBColor(0x8B, 0x1A, 0x1A), RGBColor(0xF4, 0xD9, 0xD9)),
        "purple": (RGBColor(0x4B, 0x1A, 0x6E), RGBColor(0xE9, 0xD5, 0xF5)),
        "gold":   (RGBColor(0x7A, 0x5C, 0x0A), RGBColor(0xF5, 0xEA, 0xC8)),
        "teal":   (RGBColor(0x0E, 0x4D, 0x4D), RGBColor(0xD2, 0xEF, 0xEF)),
        "gray":   (RGBColor(0x33, 0x33, 0x33), RGBColor(0xE8, 0xE8, 0xE8)),
    }
    dark, light = THEMES.get(theme, THEMES["blue"])
    try:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches as _In, Pt as _Pt
        prs = Presentation()

        def _add_rect(slide, x, y, w, h, color, line=False):
            """添加色块/色带（装饰元素）"""
            sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _In(x), _In(y), _In(w), _In(h))
            sh.fill.solid()
            sh.fill.fore_color.rgb = color
            if not line:
                sh.line.fill.background()
            sh.shadow.inherit = False
            return sh

        # ===== 封面页 =====
        cover = prs.slides.add_slide(prs.slide_layouts[6])   # 空白布局，自由排版
        if layout == "clean":
            # 极简：白底 + 深色标题 + 细线 + 副标题
            tb = cover.shapes.add_textbox(_In(1.2), _In(2.6), _In(11), _In(1.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.size = _Pt(44)
            r.font.bold = True
            r.font.color.rgb = dark
            _add_rect(cover, 1.3, 4.4, 2.0, 0.05, dark)
            if subtitle:
                sb = cover.shapes.add_textbox(_In(1.3), _In(4.7), _In(10), _In(0.8))
                sr_ = sb.text_frame.paragraphs[0].add_run()
                sr_.text = subtitle
                sr_.font.size = _Pt(18)
                sr_.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        elif layout == "dark":
            # 深色主题页
            _add_rect(cover, 0, 0, 13.33, 7.5, RGBColor(0x14, 0x16, 0x1A))
            _add_rect(cover, 0, 3.2, 13.33, 0.03, light)
            tb = cover.shapes.add_textbox(_In(1.2), _In(2.4), _In(11), _In(1.6))
            p = tb.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.size = _Pt(42)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            if subtitle:
                sb = cover.shapes.add_textbox(_In(1.2), _In(4.3), _In(10), _In(0.8))
                sr_ = sb.text_frame.paragraphs[0].add_run()
                sr_.text = subtitle
                sr_.font.size = _Pt(18)
                sr_.font.color.rgb = light
        else:
            # corporate：主题色背景 + 白字 + 装饰
            _add_rect(cover, 0, 0, 13.33, 7.5, dark)
            _add_rect(cover, 0, 0, 13.33, 0.18, light)
            _add_rect(cover, 0, 7.32, 13.33, 0.18, light)
            _add_rect(cover, 0, 2.2, 0.35, 3.1, RGBColor(0xFF, 0xFF, 0xFF))
            tb = cover.shapes.add_textbox(_In(1.1), _In(2.5), _In(11), _In(1.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.size = _Pt(40)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            _add_rect(cover, 1.15, 4.45, 2.6, 0.07, light)
            if subtitle:
                sb = cover.shapes.add_textbox(_In(1.15), _In(4.75), _In(10.5), _In(0.9))
                sr_ = sb.text_frame.paragraphs[0].add_run()
                sr_.text = subtitle
                sr_.font.size = _Pt(20)
                sr_.font.color.rgb = light
            _add_rect(cover, 12.4, 7.0, 0.55, 0.2, light)

        # ===== 内容页（浅色背景 + 左色条 + 页码 + 顶部色带）=====
        total = len([s for s in slides if isinstance(s, dict)])
        for i, s in enumerate(slides):
            if not isinstance(s, dict):
                continue
            stitle = str(s.get("title", "")).strip() or f"第 {i + 1} 页"
            bullets = s.get("bullets") or s.get("points") or []
            image = s.get("image") or s.get("img") or s.get("screenshot") or ""
            slide = prs.slides.add_slide(prs.slide_layouts[6])   # 空白布局，自由排版
            is_dark = layout == "dark"
            if is_dark:
                _add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x14, 0x16, 0x1A))
            # 顶部主题色带（corporate）/ 细线（clean/dark）
            if layout == "clean":
                _add_rect(slide, 0, 0, 13.33, 0.04, dark)
            else:
                _add_rect(slide, 0, 0, 13.33, 0.12, dark)
            # 左侧装饰（corporate 竖条；clean/dark 细线）
            if layout == "corporate":
                _add_rect(slide, 0, 0.9, 0.14, 1.0, dark)
            else:
                _add_rect(slide, 0.62, 2.0, 1.6, 0.04, dark)
            # 标题（深色粗体，靠近色条）
            tb = slide.shapes.add_textbox(_In(0.6), _In(0.75), _In(12), _In(1.1))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = stitle
            r.font.size = _Pt(30)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if is_dark else dark
            # 标题下装饰线
            _add_rect(slide, 0.62, 1.85, 3.2, 0.05, dark if not is_dark else light)
            # 内容区（文本）
            body_box = slide.shapes.add_textbox(_In(0.62), _In(2.2), _In(11.8), _In(4.5))
            body = body_box.text_frame
            body.word_wrap = True
            body_text_color = RGBColor(0xE8, 0xEA, 0xF0) if is_dark else RGBColor(0x33, 0x33, 0x33)
            if isinstance(bullets, list):
                first = True
                for j, b in enumerate(bullets):
                    if isinstance(b, dict):
                        btext = str(b.get("text", "")).strip()
                        blevel = int(b.get("level", 0) or 0)
                    else:
                        btext = str(b).strip()
                        blevel = 0
                    if not btext:
                        continue
                    p = body.paragraphs[0] if first else body.add_paragraph()
                    first = False
                    p.text = btext
                    p.level = min(2, max(0, blevel))
                    p.space_after = _Pt(8)
                    for rr in p.runs:
                        rr.font.size = _Pt(20 if p.level == 0 else 16)
                        rr.font.color.rgb = body_text_color
            else:
                body.paragraphs[0].text = str(bullets or "")
            # 图片嵌入（右侧）
            if image:
                img_path = _resolve_image(image)
                if img_path:
                    try:
                        slide.shapes.add_picture(img_path, _In(7.6), _In(2.3), width=_In(5.0))
                        if img_path.startswith(os.path.join(FILES_DIR, "_img_")):
                            try:
                                os.remove(img_path)
                            except Exception:
                                pass
                    except Exception as e:
                        return {"error": f"图片嵌入失败（格式可能不受支持：{str(e)[:80]}）。"
                                         "请改用 PNG/JPG 格式图片，或去掉该页图片后重试。"}
                else:
                    return {"error": f"第 {i + 1} 页图片无法解析：{str(image)[:80]}。"
                                     "请提供有效的图片 URL / 本地路径 / base64 图片。"}
            # 页脚：页码 + 底部装饰线
            _add_rect(slide, 0, 7.32, 13.33, 0.05, dark)
            pg = slide.shapes.add_textbox(_In(12.2), _In(7.05), _In(0.9), _In(0.4))
            pf = pg.text_frame
            pp = pf.paragraphs[0]
            pr_ = pp.add_run()
            pr_.text = f"{i + 1} / {max(total, i + 1)}"
            pr_.font.size = _Pt(12)
            pr_.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        fname = f"{_safe(title)}_{int(time.time())}.pptx"
        fpath = os.path.join(FILES_DIR, fname)
        prs.save(fpath)
        url = f"http://127.0.0.1:8000/files/{fname}"
        return {
            "ok": True,
            "file": fname,
            "url": url,
            "note": f"PPT 已生成：{fname}（封面 + {len(slides)} 页，主题色 {theme}）。请在回答里引用链接，用户点击即可预览/下载。",
        }
    except Exception as e:
        return {"error": f"PPT 生成失败: {e}"}


PLUGIN_TOOLS = [
    {
        "name": "create_pptx",
        "description": "创建 PPT 演示文稿（服务器 python-pptx 生成，中文可靠；支持多种布局/主题色/图片嵌入）。"
                       "当用户要求『做一个 PPT/演示文稿/幻灯片/汇报』时，必须使用本工具，不要自己手写 .pptx。"
                       "布局选择（按内容气质选，不要死用一套）：corporate=商务主题色（默认，适合汇报）；"
                       "clean=极简白底（适合学术/清新）；dark=深色科技风（适合发布/炫酷）。"
                       "图片获取：优先用 playwright-mcp 浏览器截图（browser_take_screenshot），或图片 URL；"
                       "通过每页 image 参数嵌入。"
                       "参数：title=标题；subtitle=副标题；theme=blue/green/red/purple/gold/teal/gray；"
                       "layout=corporate/clean/dark；slides=[{title, bullets:[...], image}]。生成后返回链接，请在回答中引用。",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "演示文稿标题"},
                "subtitle": {"type": "string", "description": "副标题（封面页）"},
                "theme": {"type": "string", "enum": ["blue", "green", "red", "purple", "gold", "teal", "gray"], "description": "主题色"},
                "layout": {"type": "string", "enum": ["auto", "corporate", "clean", "dark"], "description": "布局风格：corporate 商务 / clean 极简 / dark 深色科技"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "本页标题"},
                            "bullets": {"type": "array", "items": {"type": "object"}, "description": "本页要点：[字符串 或 {text, level}]"},
                            "image": {"type": "string", "description": "本页图片：网页截图文件路径 / 图片 URL / files 目录文件名"}
                        },
                        "required": ["title"]
                    },
                    "description": "每页内容，最多 20 页"
                }
            },
            "required": ["title", "slides"],
        },
        "handler": create_pptx,
    }
]
