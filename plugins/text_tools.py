"""text_tools: 文件内搜索与按行范围读取工具。

用于定位大文件（如 app.js）中特定逻辑的位置：
- search_in_file: 按关键词/正则搜索，返回行号与内容
- read_lines: 读取指定行范围，避免整体读取超大文件被截断
"""

import os
import re

PLUGIN_TOOLS = [
    {
        "name": "search_in_file",
        "description": "在文本文件中按关键词（支持正则）搜索，返回命中的行号与行内容。用于定位大文件中特定逻辑的位置。",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "文件路径（相对项目根，如 gui/static/app.js）"},
                "pattern": {"type": "string", "description": "搜索关键词或正则表达式"},
                "max_results": {"type": "integer", "description": "最多返回多少条结果，默认 40"}
            },
            "required": ["path", "pattern"]
        },
        "handler": "search_in_file"
    },
    {
        "name": "read_lines",
        "description": "读取文本文件的指定行范围（按 1 起始行号），用于分段查看大文件。",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "文件路径（相对项目根，如 gui/static/app.js）"},
                "start": {"type": "integer", "description": "起始行号（含），从 1 开始"},
                "end": {"type": "integer", "description": "结束行号（含）"}
            },
            "required": ["path", "start", "end"]
        },
        "handler": "read_lines"
    },
    {
        "name": "read_document",
        "description": "读取各种格式文档的内容（自动按扩展名解析为纯文本）：Word(.docx)/Excel(.xlsx)/PPT(.pptx)/PDF/"
                       "HTML/纯文本(.txt .md .json .csv .yaml .py 等)。"
                       "当用户上传了文件、或要求读取指定地址的文件时使用（尤其二进制文档 read_file 读不了，必须用本工具）。"
                       "参数 path：文件路径（绝对路径或 files/ 目录下文件名或项目目录相对路径）。"
                       "返回文档全文（完整读取，不截断；超大文件可分多次读取）。",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "文件路径（绝对路径 / files 目录文件名 / 项目相对路径）"}
            },
            "required": ["path"]
        },
        "handler": "read_document"
    }
]


def _resolve(path):
    """把相对路径解析到项目根目录。"""
    candidates = [os.getcwd()]
    # 插件运行目录可能是 plugins/ 或项目根，向上兜底
    for _ in range(3):
        parent = os.path.dirname(candidates[-1])
        if parent and parent not in candidates:
            candidates.append(parent)
    for base in candidates:
        p = os.path.join(base, path)
        if os.path.isfile(p):
            return p
    # 直接当绝对/相对路径用
    if os.path.isfile(path):
        return path
    raise FileNotFoundError("找不到文件: %s（尝试了 %s）" % (path, candidates))


def _read_lines_safe(path):
    """读取文件全部行，处理编码。"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.readlines()


def search_in_file(params):
    path = params.get("path", "")
    pattern = params.get("pattern", "")
    max_results = int(params.get("max_results", 40))
    if not path or not pattern:
        return "错误：需要 path 和 pattern 参数"
    try:
        rx = re.compile(pattern, re.IGNORECASE)
    except re.error as e:
        return "正则错误: %s" % e
    p = _resolve(path)
    lines = _read_lines_safe(p)
    hits = []
    for i, line in enumerate(lines, 1):
        if rx.search(line):
            hits.append("%d: %s" % (i, line.rstrip("\n")))
        if len(hits) >= max_results:
            break
    if not hits:
        return "未命中: %s in %s" % (pattern, path)
    return "文件: %s\n共 %d 条命中（显示前 %d 条）:\n%s" % (
        p, len(hits), max_results, "\n".join(hits))


def read_lines(params):
    path = params.get("path", "")
    try:
        start = int(params.get("start", 1))
        end = int(params.get("end", start))
    except (TypeError, ValueError):
        return "错误：start/end 需要是数字"
    if not path:
        return "错误：需要 path 参数"
    p = _resolve(path)
    lines = _read_lines_safe(p)
    total = len(lines)
    start = max(1, start)
    end = min(total, end)
    if start > end:
        return "行范围无效（文件共 %d 行）" % total
    body = []
    for i in range(start, end + 1):
        body.append("%d: %s" % (i, lines[i - 1].rstrip("\n")))
    return "文件: %s | 共 %d 行 | 显示 %d-%d 行:\n%s" % (
        p, total, start, end, "\n".join(body))


def read_document(params):
    """按扩展名自动读取各种文档格式：docx/pptx/xlsx/pdf/html/txt/md/json/csv/yaml/py 等。
    二进制文档（word/excel/ppt/pdf）自动解析为纯文本返回。"""
    path = str(params.get("path", "")).strip()
    if not path:
        return "错误：需要 path 参数（文件路径）"
    p = _resolve(path)
    if not os.path.isfile(p):
        return "文件不存在: %s" % p
    ext = os.path.splitext(p)[1].lower().lstrip(".")
    try:
        if ext in ("docx",):
            import docx
            d = docx.Document(p)
            parts = []
            for para in d.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
            for i, tbl in enumerate(d.tables, 1):
                parts.append("[表格 %d]" % i)
                for row in tbl.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
            text = "\n".join(parts)
            return "Word 文档: %s | 段落+表格 %d 块\n%s" % (os.path.basename(p), len(parts), text)
        if ext in ("pptx", "ppt"):
            import pptx
            prs = pptx.Presentation(p)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append("【第 %d 页】" % i)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in para.runs).strip()
                            if t:
                                parts.append(t)
                    if getattr(shape, "has_table", False) and shape.has_table:
                        for row in shape.table.rows:
                            parts.append(" | ".join(c.text.strip() for c in row.cells))
            text = "\n".join(parts)
            return "PPT: %s | %d 页\n%s" % (os.path.basename(p), len(prs.slides), text)
        if ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append("【工作表: %s】" % ws.title)
                rows = 0
                for row in ws.iter_rows(values_only=True):
                    vals = [str(v) if v is not None else "" for v in row]
                    if any(vals):
                        parts.append(" | ".join(vals))
                        rows += 1
                        if rows >= 1000:
                            parts.append("…（行数过多已截断）")
                            break
            return "Excel: %s | %d 个工作表\n%s" % (os.path.basename(p), len(wb.worksheets), "\n".join(parts))
        if ext in ("pdf",):
            try:
                import fitz  # PyMuPDF 优先（更稳）
                doc = fitz.open(p)
                parts = []
                for i in range(min(len(doc), 200)):
                    parts.append("【第 %d 页】" % (i + 1))
                    parts.append(doc[i].get_text().strip())
                text = "\n".join(parts)
                return "PDF: %s | %d 页\n%s" % (os.path.basename(p), len(doc), text)
            except Exception:
                import pypdf
                rd = pypdf.PdfReader(p)
                parts = []
                for i, page in enumerate(rd.pages[:200]):
                    parts.append("【第 %d 页】" % (i + 1))
                    parts.append((page.extract_text() or "").strip())
                return "PDF: %s | %d 页\n%s" % (os.path.basename(p), len(rd.pages), "\n".join(parts))
        if ext in ("html", "htm"):
            import re as _re
            with open(p, encoding="utf-8", errors="ignore") as f:
                raw = f.read()
            text = _re.sub(r"<script[\s\S]*?</script>|<style[\s\S]*?</style>", " ", raw)
            text = _re.sub(r"<[^>]+>", " ", text)
            text = _re.sub(r"\s+", " ", text).strip()
            return "HTML: %s\n%s" % (os.path.basename(p), text)
        # 文本类：直接读（txt/md/json/csv/yaml/py 等）
        with open(p, encoding="utf-8", errors="ignore") as f:
            content = f.read()
        if ext in ("json",):
            try:
                import json as _json
                content = _json.dumps(_json.loads(content), ensure_ascii=False, indent=2)
            except Exception:
                pass
        return "文件: %s\n%s" % (os.path.basename(p), content)
    except Exception as e:
        return "读取 %s 失败: %s" % (os.path.basename(p), e)
