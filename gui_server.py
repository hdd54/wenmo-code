# --- PyInstaller windowed stdio fix ---
# 当以 PyInstaller windowed（无控制台）模式打包运行时，sys.stdout/stderr/stdin
# 为 None；uvicorn 初始化日志格式器时调用 sys.stdout.isatty() 会崩溃（AttributeError）。
# 此处将 None 重定向到空设备，保证程序正常启动。
import sys as _sys, os as _os
if _sys.stdout is None:
    _sys.stdout = open(_os.devnull, "w", encoding="utf-8")
if _sys.stderr is None:
    _sys.stderr = open(_os.devnull, "w", encoding="utf-8")
if _sys.stdin is None:
    _sys.stdin = open(_os.devnull, "r", encoding="utf-8")
# --- end PyInstaller windowed stdio fix ---

"""
GUI 后端（第二课上半场）：
把第一课 chat.py 的聊天能力变成 HTTP 服务，前端页面通过 SSE 流式接收回复。

运行：python gui_server.py  然后浏览器打开 http://127.0.0.1:8000

架构上你只需要懂三件事：
1. 前端（浏览器）发 HTTP 请求给这个服务 —— 前后端分离
2. 本服务代替前端去调模型 API（key 只藏在服务端，不暴露给浏览器）
3. SSE 流式：模型每吐一个字，就用 "data: {...}" 推给前端 —— 所以打字效果是"真的在打字"
"""
import asyncio
import base64
import copy
import hashlib
import io
import json
import os
import re
import shlex
import subprocess
import sys
import tempfile
import threading
import time
import uuid
import urllib.request
import urllib.parse
import zipfile

# ---- 打包适配（必须在任何项目模块 import 之前）：PyInstaller 打包版数据目录 → %APPDATA%/问墨 ----
# history/skills_loader/plugins_loader/mcp_client 等读取 WENMO_DATA_DIR 决定数据位置，
# 否则都写进临时解包目录（每次启动变化）导致数据丢失。
if getattr(sys, "frozen", False):
    _appdata = os.environ.get("APPDATA") or os.path.expanduser("~")
    _data_root = os.path.join(_appdata, "问墨")
    os.makedirs(_data_root, exist_ok=True)
    os.environ["WENMO_DATA_DIR"] = _data_root          # 数据目录（history/workspace/files）
    _res_root = getattr(sys, "_MEIPASS", os.path.dirname(os.path.abspath(__file__)))
    os.environ["WENMO_RES_DIR"] = _res_root            # 资源目录（gui/plugins/skills 只读）

# latex2mathml 延迟导入（见 math_to_omml，启动提速）
import uvicorn
from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import Response, StreamingResponse, FileResponse, JSONResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.gzip import GZipMiddleware
from lxml import etree
from openai import OpenAI, AsyncOpenAI
from pydantic import BaseModel

from chat import load_providers, save_providers
import history as history_store
from mcp_client import mcp_call, mcp_openai_tools, mcp_refresh, mcp_snapshot
import plugins_loader
import cluster
from execution_context import current_tenant, current_workspace
from conversation_workspace import resolve_conversation_workspace
from permission_engine import evaluate_permission, load_runtime_policy, normalize_policy
from sandbox_runner import get_sandbox_diagnostics, get_sandbox_status
from tenant_state import atomic_write_json as write_tenant_json
from tenant_state import load_json as load_tenant_json
from tenant_state import tenant_file
from tenant_state import files_dir as tenant_files_dir
from tenant_state import apps_dir as tenant_apps_dir
from tenant_state import tenant_name
from usage_accounting import from_openai_usage, reported_tool_usage
from billing_api import router as billing_router


def _files_dir():
    return tenant_files_dir(create=True)


def _user_apps_dir():
    return tenant_apps_dir(create=True)


def _atexit_cleanup():
    """优雅退出：清理残留的临时文件（.tmp*），避免下次启动残留脏文件"""
    try:
        history_store.cleanup_tmp_files()
    except Exception:
        pass
    try:
        for qid in list(PENDING_ASKS.keys()):
            try:
                if not PENDING_ASKS[qid].done():
                    PENDING_ASKS[qid].set_result('（服务退出）')
            except Exception:
                pass
    except Exception:
        pass


import atexit
atexit.register(_atexit_cleanup)   # 分布式任务集群 P1（任务队列 + worker 池）
import skills_loader
from skills_loader import load_skills, match_skills, generic_skills

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
STATIC_DIR = os.path.join(BASE_DIR, "gui", "static")

# ---- 打包适配（顶部 import 前已设 WENMO_DATA_DIR/WENMO_RES_DIR）：这里只重定向 STATIC_DIR ----
# 首次启动引导：把打包自带的默认配置（providers/settings/mcp.json）复制到数据目录，
# 之后用户保存的 key/设置都写数据目录（可写、持久、不随安装目录变化）
if getattr(sys, "frozen", False) and os.environ.get("WENMO_DATA_DIR"):
    _cfg_boot_done = False
    try:
        _data_root = os.environ["WENMO_DATA_DIR"]
        _src_candidates = [os.path.dirname(sys.executable), getattr(sys, "_MEIPASS", "")]
        for _cfg_name in ("providers.json", "settings.json", "mcp.json"):
            _dst = os.path.join(_data_root, _cfg_name)
            if os.path.isfile(_dst):
                continue  # 已有用户配置，不动
            for _src_dir in _src_candidates:
                _src = os.path.join(_src_dir, _cfg_name)
                if os.path.isfile(_src):
                    try:
                        shutil.copy2(_src, _dst)
                        _cfg_boot_done = True
                    except Exception:
                        pass
                    break
    except Exception:
        pass

# ---- 内容层分离（程序与内容分离方案）：gui/plugins/skills/apps 数据目录优先 + seed 兜底 + 首启种子化 ----
CONTENT_DIR = ""
if getattr(sys, "frozen", False) and os.environ.get("WENMO_DATA_DIR"):
    _data_root = os.environ["WENMO_DATA_DIR"]
    _res_root = os.environ.get("WENMO_RES_DIR") or getattr(sys, "_MEIPASS", "")
    _seed_dir = os.path.join(_res_root, "seed") if _res_root else ""
    _content_dir = os.path.join(_data_root, "content")
    try:
        if _seed_dir and os.path.isdir(_seed_dir) and not os.path.isdir(_content_dir):
            shutil.copytree(_seed_dir, _content_dir)   # 首启种子化：seed → 数据目录
        os.makedirs(_content_dir, exist_ok=True)
        CONTENT_DIR = _content_dir
    except Exception:
        CONTENT_DIR = _content_dir if os.path.isdir(_content_dir) else ""

if getattr(sys, "frozen", False):
    _res_root = os.environ.get("WENMO_RES_DIR") or BASE_DIR
    # 前端静态：数据目录 content/gui/static 优先（用户可自定义主题/前端），seed 兜底
    _c_static = os.path.join(CONTENT_DIR, "gui", "static") if CONTENT_DIR else ""
    if _c_static and os.path.isdir(_c_static):
        STATIC_DIR = _c_static
    else:
        _seed_static = os.path.join(_res_root, "seed", "gui", "static")
        if os.path.isdir(_seed_static):
            STATIC_DIR = _seed_static
        else:
            STATIC_DIR = os.path.join(_res_root, "gui", "static")

# ---- MCP server 模式（内置集成）：问墨.exe --mcp-server <脚本名> → 以 MCP server 运行 ----
# 复用主程序已打包的 mcp 库与全部依赖，无需独立 exe；开发版直接 python gui_server.py --mcp-server xxx
if "--mcp-server" in sys.argv:
    _mi = sys.argv.index("--mcp-server")
    _mcp_name = sys.argv[_mi + 1] if _mi + 1 < len(sys.argv) else ""
    try:
        import importlib.util
        _exe_dir = os.path.dirname(os.path.abspath(sys.executable)) if getattr(sys, "frozen", False) else ""
        _mcp_cands = [
            os.path.join(BASE_DIR, _mcp_name + ".py"),
            os.path.join(BASE_DIR, "seed", "mcp-servers", _mcp_name + ".py"),
            os.path.join(BASE_DIR, "mcp-servers", _mcp_name + ".py"),
            os.path.join(_exe_dir, _mcp_name + ".py") if _exe_dir else "",
            os.path.join(os.environ.get("WENMO_DATA_DIR", ""), "content", "mcp-servers", _mcp_name + ".py"),
        ]
        _mcp_file = next((c for c in _mcp_cands if c and os.path.isfile(c)), "")
        if not _mcp_file:
            print("MCP server not found: %s" % _mcp_name)
            sys.exit(1)
        _spec = importlib.util.spec_from_file_location("wenmo_mcp_" + _mcp_name, _mcp_file)
        _mod = importlib.util.module_from_spec(_spec)
        _spec.loader.exec_module(_mod)
        if hasattr(_mod, "main"):
            import asyncio
            asyncio.run(_mod.main())
        sys.exit(0)
    except SystemExit:
        raise
    except Exception as _e:
        print("MCP server error: %s" % _e)
        sys.exit(1)


app = FastAPI(title="问墨·code Chat")
app.add_middleware(GZipMiddleware, minimum_size=1024)
app.include_router(billing_router)


@app.middleware("http")
async def no_cache_static(request, call_next):
    """Protect the loopback control plane and set deterministic static/file headers."""
    try:
        host = (urllib.parse.urlparse("//" + (request.headers.get("host") or "")).hostname or "").lower()
    except ValueError:
        host = ""
    if host not in ("127.0.0.1", "localhost", "::1"):
        return JSONResponse({"detail": "仅允许本机访问"}, status_code=403)
    origin = request.headers.get("origin")
    if origin:
        origin_scheme = ""
        try:
            parsed_origin = urllib.parse.urlparse(origin)
            origin_host = (parsed_origin.hostname or "").lower()
            origin_scheme = parsed_origin.scheme
            origin_port = parsed_origin.port or (443 if parsed_origin.scheme == "https" else 80)
            request_port = request.url.port or (443 if request.url.scheme == "https" else 80)
        except Exception:
            origin_host = ""
            origin_port = -1
            request_port = -2
        if (origin_host != host or origin_port != request_port
                or origin_scheme != request.url.scheme):
            return JSONResponse({"detail": "已拒绝跨站请求"}, status_code=403)
    tenant_token = None
    try:
        import auth as _request_auth
        session_token = (request.cookies.get("wenmo_session") or
                         request.headers.get("X-Wenmo-Token", "")).strip()
        login = _request_auth.validate_session(session_token) if session_token else None
        tenant_token = current_tenant.set("github-" + login if login else "local")
    except Exception:
        tenant_token = current_tenant.set("local")
    try:
        response = await call_next(request)
    finally:
        current_tenant.reset(tenant_token)
    p = request.url.path
    if p.startswith("/files/"):
        # Files may contain user/model supplied HTML or SVG. A sandboxed response cannot
        # access this app's localStorage or privileged loopback APIs even when opened.
        response.headers["Content-Security-Policy"] = "sandbox; default-src 'none'"
        response.headers["X-Content-Type-Options"] = "nosniff"
        response.headers["Cache-Control"] = "no-store"
    elif p == "/" or p.endswith((".js", ".css", ".html", ".png", ".svg", ".ico")):
        response.headers["Cache-Control"] = ("public, max-age=3600" if getattr(sys, "frozen", False)
                                             else "no-store")
    return response


def resolve_key(cfg):
    """获取 API key（服务端版）。
    注意：chat.py 里的 get_api_key() 在缺 key 时会 sys.exit 整个进程，
    服务器进程不能退出，所以这里重新实现一个"只返回 None"的版本。
    这教给你一个工程经验：命令行工具的报错方式和服务器完全不同。"""
    key = cfg.get("api_key", "").strip()
    if key:
        return key
    env_name = cfg.get("api_key_env", "").strip()
    if env_name and not getattr(sys, "frozen", False):
        return os.environ.get(env_name, "").strip()
    return None


# 各远端供应商的推荐模型列表（仅列 /chat/completions 兼容的模型；前端可自由输入自定义模型名）
MODELS_BY_PROVIDER = {
    "openai": [],  # 通过官方 /models 动态发现，也允许用户直接填写模型 ID
    "deepseek": ["deepseek-v4-flash", "deepseek-v4-pro", "deepseek-chat", "deepseek-reasoner"],
    "zen": [
        # ---- 完整 61 个模型（2026-08 从 zen /v1/models 实时拉取）----
        "claude-fable-5", "claude-opus-5", "claude-opus-4-8", "claude-opus-4-7",
        "claude-opus-4-6", "claude-opus-4-5", "claude-sonnet-5", "claude-sonnet-4-6",
        "claude-sonnet-4-5", "claude-sonnet-4", "claude-haiku-4-5",
        "gemini-3.6-flash", "gemini-3.5-flash-lite", "gemini-3.5-flash",
        "gemini-3.1-pro", "gemini-3-flash",
        "gpt-5.6-sol", "gpt-5.6-terra", "gpt-5.6-luna", "gpt-5.5", "gpt-5.5-pro",
        "gpt-5.4", "gpt-5.4-pro", "gpt-5.4-mini", "gpt-5.4-nano",
        "gpt-5.3-codex-spark", "gpt-5.3-codex", "gpt-5.2", "gpt-5.2-codex", "gpt-5.1",
        "gpt-5.1-codex-max", "gpt-5.1-codex", "gpt-5.1-codex-mini", "gpt-5", "gpt-5-codex",
        "gpt-5-nano",
        "grok-build-0.1", "grok-4.5",
        "deepseek-v4-pro", "deepseek-v4-flash",
        "glm-5.2", "glm-5.1", "glm-5",
        "minimax-m3", "minimax-m2.7", "minimax-m2.5",
        "kimi-k3", "kimi-k2.7-code", "kimi-k2.6", "kimi-k2.5",
        "qwen3.6-plus", "qwen3.5-plus", "big-pickle",
        "deepseek-v4-flash-free", "mimo-v2.5-free", "ling-3.0-flash-free",
        "ling-3.0-tiny-free", "nemotron-3-ultra-free", "north-mini-code-free",
        "laguna-s-2.1-free", "longcat-2.0-free",
    ],
    "free": [
        # ---- 聚合各平台免费档（官方确认免费；需在设置里给对应平台配免费 key）----
        # DeepSeek（已有 key 立即可用）
        "deepseek-v4-flash-free", "deepseek-chat", "deepseek-reasoner",
        # 智谱官方免费（glm-4.7-flash 永久免费）
        "glm-4.7-flash",
        # 硅基流动限免模型
        "deepseek-ai/DeepSeek-OCR", "tencent/Hunyuan-MT-7B",
        # Kimi 免费档
        "kimi-latest", "kimi-k3",
        # 通义免费档
        "qwen-turbo", "qwen-plus",
        # 豆包免费档
        "doubao-lite-32k",
    ],
    "opencode_go": [
        "deepseek-v4-flash", "deepseek-v4-pro",
        "kimi-k3", "kimi-k2.7-code", "kimi-k2.6", "kimi-k2.5",
        "glm-5.2", "glm-5.1", "grok-4.5", "grok-build-0.1",
        "mimo-v2.5", "mimo-v2.5-pro",
        "claude-fable-5", "claude-sonnet-5", "claude-opus-5", "claude-haiku-4-5",
        "gemini-3.6-flash", "gemini-3.5-flash", "gpt-5.6-sol", "gpt-5.6-terra",
        "gpt-5.6-luna", "gpt-5.5", "gpt-5.4-mini", "gpt-5-nano",
        "qwen3.6-plus", "qwen3.5-plus",
    ],
    "ollama": [],  # ����ģ�ͣ��� /api/ollama/models ʵʱ����
    "qianwen": ["qwen-max", "qwen-plus", "qwen-turbo", "qwen-long", "qwen3-max", "qwen3-plus"],
    "zhipu": ["glm-4-plus", "glm-4-air", "glm-4-flash", "glm-5-flash"],
    "siliconflow": [
        "deepseek-ai/DeepSeek-V3", "deepseek-ai/DeepSeek-R1", "deepseek-ai/DeepSeek-V3.1",
        "Qwen/Qwen2.5-72B-Instruct", "Qwen/QwQ-32B", "Qwen/Qwen3-235B-A22B",
        "THUDM/glm-4-9b-chat",
    ],
    "kimi": ["kimi-latest", "kimi-k3", "moonshot-v1-8k", "moonshot-v1-32k", "moonshot-v1-128k"],
    "doubao": ["doubao-pro-32k", "doubao-lite-32k"],
}


def ollama_alive() -> bool:
    """Ollama 守护进程是否在运行"""
    try:
        with urllib.request.urlopen("http://localhost:11434/api/tags", timeout=3) as resp:
            return resp.status == 200
    except Exception:
        return False


def _gpu_free_mem_gb():
    """查询 NVIDIA 可用显存（GB）；无 GPU / 失败返回 0（llama.cpp 走纯 CPU）"""
    try:
        out = subprocess.run(
            ["nvidia-smi", "--query-gpu=memory.free", "--format=csv,noheader,nounits"],
            capture_output=True, text=True, timeout=5, creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        if out.returncode != 0:
            return 0
        vals = [float(v.strip()) for v in out.stdout.strip().splitlines() if v.strip()]
        if not vals:
            return 0
        return max(vals) / 1024.0   # MiB → GiB
    except Exception:
        return 0


@app.get("/api/providers")
def list_providers():
    """给前端展示用的供应商列表 —— 注意：绝不在接口里暴露 api_key 本身。
    本地 GGUF 排在最前：它是离线主力，用户应第一时间看到。"""
    providers = load_providers()
    # 动态加入"本地 GGUF"供应商（模型 = 当前加载的 GGUF 文件）
    local_entry = {
        "name": "本地 GGUF（llama.cpp）",
        "model": LOCAL_STATE["name"] or "未加载",
        "has_key": True,  # 本地模型不需要 key
        "models": [LOCAL_STATE["name"]] if LOCAL_STATE["name"] else [],
    }
    entries = [("local", local_entry)] + [(k, v) for k, v in providers.items() if k != "local"]
    settings = load_settings()
    # 供应商价格表（元/百万 tokens，前端计费用）。pricing 模块缺失时返回 None → 前端不显示计费
    try:
        from pricing import get_billing_mode, get_price, get_price_source
    except Exception:
        get_price = None
        get_billing_mode = None
        get_price_source = None

    def _price_payload(provider_key, model):
        quote = get_price(provider_key, model) if get_price else (None, None, None, True)
        price = list(quote[:3]) if quote[0] is not None else None
        mode = get_billing_mode(provider_key) if get_billing_mode else "unknown"
        source = get_price_source(provider_key, model) if get_price_source else {
            "kind": "unavailable", "invoice_reconciled": False}
        return {"price": price, "price_est": bool(quote[3]), "billing_mode": mode,
                "price_source": source}

    return {
        "providers": [
            {
                "key": k,
                "name": v.get("name", k),
                "model": v["model"],
                "has_key": bool(v.get("api_key", "").strip() or (
                    not getattr(sys, "frozen", False)
                    and v.get("api_key_env", "").strip()
                    and os.environ.get(v["api_key_env"], "").strip()
                )),
                "models": MODELS_BY_PROVIDER.get(k, []),
                "ctx": settings["local_ctx"] if k == "local" else (8192 if k == "ollama" else settings["remote_ctx"]),
                **_price_payload(k, v.get("model", "")),
            }
            for k, v in entries
        ]
    }


@app.get("/api/ollama/models")
def list_ollama_models():
    """实时列出 Ollama 已安装的本地模型（绕过前端直连 11434 的跨域问题）"""
    if not ollama_alive():
        raise HTTPException(status_code=400, detail="Ollama 未运行，请先启动 ollama serve")
    try:
        with urllib.request.urlopen("http://localhost:11434/api/tags", timeout=5) as resp:
            data = json.loads(resp.read().decode("utf-8"))
        models = [
            {"name": m.get("name"), "size": m.get("size", 0)}
            for m in data.get("models", [])
        ]
        return {"models": models, "running": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"读取 Ollama 模型失败: {e}")


class OllamaImportRequest(BaseModel):
    path: str  # 本地 GGUF 文件的绝对路径


@app.post("/api/ollama/import")
def import_gguf(req: OllamaImportRequest):
    """把本地 GGUF 文件注册进 Ollama：写 Modelfile -> ollama create
    教学点：浏览器不能直接执行 shell，所以"文件导入"这种能力必须由服务器代劳。
    """
    path = req.path.strip().strip('"')
    if not path.lower().endswith(".gguf"):
        raise HTTPException(status_code=400, detail="文件必须是 .gguf 格式")
    if not os.path.isfile(path):
        raise HTTPException(status_code=400, detail=f"文件不存在: {path}")
    if os.path.getsize(path) == 0:
        raise HTTPException(status_code=400, detail="文件是空的")
    if not ollama_alive():
        raise HTTPException(status_code=400, detail="Ollama 未运行，请先启动 ollama serve")

    name = os.path.splitext(os.path.basename(path))[0].lower()
    name = "".join(c if (c.isalnum() or c in "._-") else "-" for c in name).strip("-.")
    if not name:
        raise HTTPException(status_code=400, detail="无法从文件名生成模型名")

    # 已存在同名模型 -> 拒绝（避免误覆盖）
    try:
        with urllib.request.urlopen("http://localhost:11434/api/tags", timeout=5) as resp:
            existing = {m.get("name") for m in json.loads(resp.read().decode()).get("models", [])}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"读取 Ollama 模型失败: {e}")
    if name in existing:
        raise HTTPException(status_code=409, detail=f"模型 {name} 已存在，请先删除再导入（或换个文件名）")

    # 写 Modelfile（Windows 反斜杠路径转成正斜杠，ollama 更稳）
    modelfile = f"FROM {path.replace(os.sep, '/')}\n"
    fd, modelfile_path = tempfile.mkstemp(suffix=".Modelfile", prefix="agent-tutorial-", text=True)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            f.write(modelfile)
        proc = subprocess.run(
            ["ollama", "create", name, "-f", modelfile_path],
            capture_output=True, text=True, timeout=120,
        )
        if proc.returncode != 0:
            # 去掉 ollama 进度条的 ANSI 转义码，让错误消息干净可读
            raw = proc.stderr.strip() or proc.stdout.strip()
            clean = re.sub(r"\x1b\[[0-9;?]*[a-zA-Z]", "", raw)
            raise HTTPException(status_code=500, detail=f"ollama create 失败: {clean}")
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="ollama create 超时（模型过大？）")
    finally:
        try:
            os.remove(modelfile_path)
        except OSError:
            pass

    return {"ok": True, "name": name}


# ============================================================================
# 本地 GGUF 推理服务：直接加载任意 .gguf 文件，不依赖 Ollama
# 原理：用 llama-cpp-python 的服务器模式，把 GGUF 文件变成一个
#       OpenAI 兼容的本地 API 服务 —— "模型文件即服务"。
# 同一时间只跑一个本地模型（加载新文件会先停掉旧的，省内存）。
# ============================================================================
LOCAL_STATE = {
    "proc": None,      # subprocess.Popen
    "path": None,      # 原始文件路径
    "name": None,      # 模型名（文件名）
    "port": None,      # 分配的端口
    "mmproj": None,    # 自动配对的图像投影文件（若有）
    "status": "idle",  # idle | loading | ready | error
    "error": None,
}
_LOCAL_LOG = os.path.join(BASE_DIR, "local_server.log")


def _find_free_port(start=8090, end=8999):
    import socket
    for port in range(start, end):
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            try:
                s.bind(("127.0.0.1", port))
                return port
            except OSError:
                continue
    raise RuntimeError("没有可用端口")


def _kill_llama_processes():
    """扫描并杀掉所有 llama-server.exe 进程（含进程树），返回被杀 PID 列表。
    软件设计是"同一时间只跑一个本地模型"，因此所有残留 llama-server 均可安全清理。"""
    killed = []
    try:
        out = subprocess.run(
            ["tasklist", "/FI", "IMAGENAME eq llama-server.exe", "/FO", "CSV", "/NH"],
            capture_output=True, text=True, timeout=10,
        ).stdout
        for line in out.splitlines():
            parts = [p.strip().strip('"') for p in line.strip().split('","')]
            if len(parts) < 2 or "llama-server.exe" not in parts[0]:
                continue
            pid = parts[1]
            if not pid.isdigit():
                continue
            try:
                subprocess.run(["taskkill", "/F", "/T", "/PID", pid],
                               capture_output=True, text=True, timeout=10)
                killed.append(int(pid))
            except Exception as e:
                print("[local] 清理 llama-server PID %s 失败: %s" % (pid, e), flush=True)
    except Exception as e:
        print("[local] 扫描 llama-server 进程失败: %s" % e, flush=True)
    return killed


def _local_stop():
    """停止本地推理服务：先终止受管进程（terminate→kill 递进），
    再扫描清理全部残留 llama-server 进程（含 proc 丢失的孤儿进程），
    杀完验证并留痕；失败不静默吞掉。"""
    proc = LOCAL_STATE["proc"]
    if proc is not None and proc.poll() is None:
        try:
            proc.terminate()
            proc.wait(timeout=5)
            print("[local] 已终止受管进程 PID %s" % proc.pid, flush=True)
        except Exception:
            try:
                proc.kill()
                proc.wait(timeout=5)
                print("[local] 已强制终止受管进程 PID %s" % proc.pid, flush=True)
            except Exception as e:
                print("[local] 警告：受管进程 PID %s 终止失败: %s" % (proc.pid, e), flush=True)
    # 扫描清理所有残留 llama-server（含历史孤儿进程，proc 丢失也能清掉）
    killed = _kill_llama_processes()
    if killed:
        print("[local] 清理残留 llama-server 进程: %s" % killed, flush=True)
    # 停止后验证：受管进程若仍存活则强杀兜底
    if proc is not None and proc.poll() is None:
        print("[local] 警告：进程 PID %s 仍在运行，taskkill /F /T 兜底" % proc.pid, flush=True)
        try:
            subprocess.run(["taskkill", "/F", "/T", "/PID", str(proc.pid)],
                           capture_output=True, text=True, timeout=10)
        except Exception as e:
            print("[local] 警告：强杀 PID %s 失败: %s" % (proc.pid, e), flush=True)
    LOCAL_STATE.update(proc=None, path=None, name=None, port=None, mmproj=None, status="idle", error=None)


@app.get("/api/local/status")
def local_status():
    return {k: LOCAL_STATE[k] for k in ("status", "name", "path", "port", "mmproj", "error")}


class LocalLoadRequest(BaseModel):
    path: str


class LocalScanRequest(BaseModel):
    path: str


def _match_mmproj(model_path, mmprojs):
    """为文本模型找配套的 mmproj（图像投影）：
    优先"同目录 + 名称前缀匹配"，否则"同目录唯一 mmproj"兜底。"""
    d = os.path.dirname(model_path)
    dir_mm = [m for m in mmprojs if os.path.dirname(m["path"]) == d]
    if not dir_mm:
        return None
    base = os.path.splitext(os.path.basename(model_path))[0].lower()
    for m in dir_mm:
        pfx = os.path.splitext(os.path.basename(m["path"]))[0]
        if pfx.lower().startswith("mmproj"):
            pfx = pfx[len("mmproj"):].lstrip("-")
        pfx = pfx.lower()
        if pfx and base.startswith(pfx):
            return m["path"]
    if len(dir_mm) == 1:
        return dir_mm[0]["path"]
    return None


@app.post("/api/local/scan")
def local_scan(req: LocalScanRequest):
    """递归扫描文件夹里的所有 .gguf 文件，按大小倒序返回。
    mmproj-* 是多模态投影文件（图像部分），必须搭配同目录的文本模型使用；
    扫描会给每个文本模型自动配对好 mmproj（如果有）。"""
    path = req.path.strip().strip('"')
    if not os.path.isdir(path):
        raise HTTPException(status_code=400, detail=f"文件夹不存在: {path}")
    models, mmprojs = [], []
    for root, _dirs, files in os.walk(path):
        for f in files:
            if not f.lower().endswith(".gguf"):
                continue
            full = os.path.join(root, f)
            try:
                size = os.path.getsize(full)
            except OSError:
                continue
            if f.lower().startswith("mmproj"):
                mmprojs.append({"path": full, "name": os.path.splitext(f)[0], "size": size, "kind": "mmproj"})
            else:
                models.append({"path": full, "name": os.path.splitext(f)[0], "size": size, "kind": "model"})
    for m in models:
        m["mmproj"] = _match_mmproj(m["path"], mmprojs)
    found = sorted(models + mmprojs, key=lambda x: x["size"], reverse=True)
    return {"path": path, "count": len(found), "files": found}


@app.post("/api/local/load")
def local_load(req: LocalLoadRequest):
    """加载一个 GGUF 文件并起本地推理服务（新文件会替换旧模型）"""
    path = req.path.strip().strip('"')
    if not path.lower().endswith(".gguf"):
        raise HTTPException(status_code=400, detail="文件必须是 .gguf 格式")
    if not os.path.isfile(path):
        raise HTTPException(status_code=400, detail=f"文件不存在: {path}")
    if os.path.getsize(path) == 0:
        raise HTTPException(status_code=400, detail="文件是空的")

    _local_stop()
    name = os.path.splitext(os.path.basename(path))[0]
    port = _find_free_port()

    # 用官方 llama.cpp 的 llama-server.exe 加载（若不在项目 llama/ 目录则报错）
    # 多路径查找 llama-server.exe：exe 旁/llama → _internal/llama → 数据目录/llama
    llama_exe = None
    _llama_cands = []
    if getattr(sys, "frozen", False):
        _llama_cands.append(os.path.join(os.path.dirname(os.path.abspath(sys.executable)), "llama", "llama-server.exe"))
    _llama_cands += [
        os.path.join(BASE_DIR, "llama", "llama-server.exe"),
        os.path.join(os.environ.get("WENMO_DATA_DIR", ""), "llama", "llama-server.exe"),
    ]
    for _c in _llama_cands:
        if _c and os.path.isfile(_c):
            llama_exe = _c
            break
    if not llama_exe:
        raise HTTPException(status_code=500, detail="未找到 llama-server.exe（请在程序目录或数据目录 llama/ 放置）")
    cmd = [
        llama_exe, "-m", path,
        "--host", "127.0.0.1", "--port", str(port),
        "-c", str(load_settings()["local_ctx"]),   # 上下文上限可在设置里调整（16K~1M）
    ]
    # GPU 加速（llama.cpp）：按可用显存自动分配 GPU 层数 + KV 缓存量化
    # 显存充足 → 全部层上 GPU；显存不足 → 按模型大小估算可容纳的层数
    try:
        vram = _gpu_free_mem_gb()   # 可用显存（GB），无 GPU 返回 0
    except Exception:
        vram = 0
    model_gb = os.path.getsize(path) / (1024 ** 3)
    if vram > 1.5:
        if vram >= model_gb * 1.3 + 1.0:
            ngl = 999      # 显存够：全部层上 GPU（999 = 全部）
        else:
            # 估算：模型每层约 model_gb/层数；留 1.2GB 给 KV/计算
            ngl = max(1, int((vram - 1.2) / max(model_gb / 80, 0.01)) * 10)
        cmd += ["-ngl", str(min(ngl, 999))]
        # KV 缓存 8bit 量化：显存占用约减半，可换更多 GPU 层（llama.cpp ≥ b4xxx 支持）
        cmd += ["-ctk", "q8_0", "-ctv", "q8_0"]
        # 线程：物理核心数（默认即可，交给 llama 自检），这里不覆盖
    # 自动配对同目录的 mmproj（多模态图像投影），让图像理解能力直接可用
    mmproj_path = None
    mmprojs = []
    model_dir = os.path.dirname(path)
    try:
        for f in os.listdir(model_dir):
            if f.lower().startswith("mmproj") and f.lower().endswith(".gguf"):
                mmprojs.append({"path": os.path.join(model_dir, f), "name": f, "size": 0})
    except OSError:
        pass
    if mmprojs:
        mmproj_path = _match_mmproj(path, mmprojs)
    if mmproj_path:
        cmd += ["--mmproj", mmproj_path]

    logf = open(_LOCAL_LOG, "a", encoding="utf-8")
    proc = subprocess.Popen(
        cmd, stdout=logf, stderr=subprocess.STDOUT,
        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
    )
    LOCAL_STATE.update(proc=proc, path=path, name=name, port=port, mmproj=mmproj_path, status="loading", error=None)

    def _wait_ready():
        import time
        deadline = time.time() + 240
        while time.time() < deadline:
            if proc.poll() is not None:
                LOCAL_STATE["status"] = "error"
                LOCAL_STATE["error"] = f"推理服务进程退出（代码 {proc.returncode}），详见 local_server.log"
                return
            try:
                with urllib.request.urlopen(f"http://127.0.0.1:{port}/health", timeout=2) as resp:
                    if resp.status == 200:
                        LOCAL_STATE["status"] = "ready"
                        return
            except Exception:
                pass
            time.sleep(1.5)
        LOCAL_STATE["status"] = "error"
        LOCAL_STATE["error"] = "加载超时（240 秒），详见 local_server.log"

    threading.Thread(target=_wait_ready, daemon=True).start()
    return {"ok": True, "name": name, "port": port, "status": "loading", "mmproj": mmproj_path}


@app.post("/api/local/stop")
def local_stop():
    _local_stop()
    return {"ok": True, "status": "idle"}


@app.get("/api/skills")
def list_skills():
    """列出 opencode 技能库里的技能（供前端展示）；标记通用型/特点型分类"""
    skills = load_skills()
    return {
        "total": len(skills),
        "generic": [n for n in skills_loader.GENERIC_SKILL_NAMES if any(s["name"] == n for s in skills)],
        "skills": [
            {
                "name": s["name"],
                "description": s["description"],
                "source": s["source"],
                "is_generic": skills_loader.is_generic(s["name"]),
            }
            for s in skills
        ],
    }


# ============================================================================
# 上下文设置（16K ~ 1M 可调）+ 系统资源检测（内存/显存）
# ============================================================================
SETTINGS_FILE = os.path.join(os.environ.get("WENMO_DATA_DIR") or BASE_DIR, "settings.json")  # 打包版：读数据目录
_CONFIG_WRITE_LOCK = threading.Lock()
DEFAULT_SETTINGS = {"local_ctx": 16384, "remote_ctx": 32768, "vision_provider": "", "vision_model": "",
"default_provider": "", "default_model": "", "agent_provider": "", "agent_model": "",
"agent_vision_provider": "", "agent_vision_model": "",
"permissions": {"write_files": "allow", "run_command": "ask"},
"permission_policy": {}, "sandbox_mode": "required", "sandbox_image": "wenmo-agent:locked"}


def _atomic_json_write(path, data):
    """Serialize config writes so crashes/concurrent requests cannot leave truncated JSON."""
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
    tmp = path + ".tmp." + uuid.uuid4().hex[:8]
    try:
        with _CONFIG_WRITE_LOCK:
            with open(tmp, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
                f.flush()
                os.fsync(f.fileno())
            os.replace(tmp, path)
    finally:
        try:
            if os.path.isfile(tmp):
                os.remove(tmp)
        except OSError:
            pass


def load_settings():
    try:
        with open(SETTINGS_FILE if tenant_name() == "local" else tenant_file("settings.json"),
                  encoding="utf-8") as f:
            s = json.load(f)
        return {**copy.deepcopy(DEFAULT_SETTINGS), **s}
    except Exception:
        return copy.deepcopy(DEFAULT_SETTINGS)


def _settings_file():
    return SETTINGS_FILE if tenant_name() == "local" else tenant_file("settings.json")


def _ram_info():
    """Windows 内存使用（ctypes 标准库）"""
    import ctypes

    class MEMORYSTATUSEX(ctypes.Structure):
        _fields_ = [
            ("dwLength", ctypes.c_ulong), ("dwMemoryLoad", ctypes.c_ulong),
            ("ullTotalPhys", ctypes.c_ulonglong), ("ullAvailPhys", ctypes.c_ulonglong),
            ("ullTotalPageFile", ctypes.c_ulonglong), ("ullAvailPageFile", ctypes.c_ulonglong),
            ("ullTotalVirtual", ctypes.c_ulonglong), ("ullAvailVirtual", ctypes.c_ulonglong),
            ("ullAvailExtendedVirtual", ctypes.c_ulonglong),
        ]

    m = MEMORYSTATUSEX()
    m.dwLength = ctypes.sizeof(MEMORYSTATUSEX)
    ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(m))
    return {"total": m.ullTotalPhys, "used": m.ullTotalPhys - m.ullAvailPhys}


def _vram_info():
    """NVIDIA 显存使用（nvidia-smi）；无 NVIDIA 显卡返回 None"""
    try:
        out = subprocess.run(
            ["nvidia-smi", "--query-gpu=memory.total,memory.used", "--format=csv,noheader,nounits"],
            capture_output=True, text=True, timeout=5,
        ).stdout.strip()
        parts = out.split(",")
        if len(parts) == 2:
            return {"total": int(parts[0].strip()) * 2**20, "used": int(parts[1].strip()) * 2**20}
    except Exception:
        pass
    return None


@app.get("/api/settings/context")
def get_context_settings():
    """上下文设置 + 图像识别模型 + 系统资源（内存/显存）+ 本地模型状态"""
    s = load_settings()
    return {
        "local_ctx": s["local_ctx"],
        "remote_ctx": s["remote_ctx"],
        "vision_provider": s.get("vision_provider", ""),
        "vision_model": s.get("vision_model", ""),
        "default_provider": s.get("default_provider", ""),
        "default_model": s.get("default_model", ""),
        "agent_provider": s.get("agent_provider", ""),
        "agent_model": s.get("agent_model", ""),
        "agent_vision_provider": s.get("agent_vision_provider", ""),
        "agent_vision_model": s.get("agent_vision_model", ""),
        "permissions": s.get("permissions") or {"write_files": "allow", "run_command": "ask"},
        "resources": {
            "ram": _ram_info(),
            "vram": _vram_info(),
            "local_model": {"name": LOCAL_STATE["name"], "status": LOCAL_STATE["status"], "mmproj": LOCAL_STATE["mmproj"]},
        },
    }


class ContextSettingsRequest(BaseModel):
    local_ctx: int = 0
    remote_ctx: int = 0
    vision_provider: str = ""
    vision_model: str = ""
    default_provider: str = ""
    default_model: str = ""
    agent_provider: str = ""
    agent_model: str = ""
    agent_vision_provider: str = ""
    agent_vision_model: str = ""
    permissions: dict = {}


@app.post("/api/settings/context")
def set_context_settings(req: ContextSettingsRequest):
    """保存上下文上限、图像模型、默认文本模型与智能体模型；本地上下文改动且模型在跑 → 自动用新值重载"""
    s = load_settings()
    changed_local = False
    if 16384 <= req.local_ctx <= 1048576:
        changed_local = req.local_ctx != s["local_ctx"]
        s["local_ctx"] = req.local_ctx
    if 16384 <= req.remote_ctx <= 1048576:
        s["remote_ctx"] = req.remote_ctx
    s["vision_provider"] = req.vision_provider.strip()
    s["vision_model"] = req.vision_model.strip()
    s["default_provider"] = req.default_provider.strip()
    s["default_model"] = req.default_model.strip()
    s["agent_provider"] = req.agent_provider.strip()
    s["agent_model"] = req.agent_model.strip()
    s["agent_vision_provider"] = req.agent_vision_provider.strip()
    s["agent_vision_model"] = req.agent_vision_model.strip()
    # 权限（对标 opencode）：write_files=allow/ask/deny；run_command=allow/ask/deny
    perms = s.get("permissions") or {}
    if isinstance(req.permissions, dict):
        for k in ("write_files", "run_command"):
            if req.permissions.get(k) in ("allow", "ask", "deny"):
                perms[k] = req.permissions[k]
    s["permissions"] = perms
    _atomic_json_write(_settings_file(), s)
    if changed_local and LOCAL_STATE["status"] in ("ready", "loading") and LOCAL_STATE["path"]:
        local_load(LocalLoadRequest(path=LOCAL_STATE["path"]))  # 内部先停旧再启新
    return {"ok": True, "local_ctx": s["local_ctx"], "remote_ctx": s["remote_ctx"],
            "vision_provider": s["vision_provider"], "vision_model": s["vision_model"],
            "default_provider": s["default_provider"], "default_model": s["default_model"],
            "agent_provider": s["agent_provider"], "agent_model": s["agent_model"],
            "agent_vision_provider": s["agent_vision_provider"], "agent_vision_model": s["agent_vision_model"]}


class PermissionRequest(BaseModel):
    write_files: str = ""   # allow / ask / deny
    run_command: str = ""   # allow / ask / deny
    policy: dict = {}
    sandbox_mode: str = ""
    sandbox_image: str = ""


@app.get("/api/settings/permissions")
def get_permissions():
    """读取兼容开关 + 按工具/路径/命令的完整权限矩阵。"""
    s = load_settings()
    return {
        **(s.get("permissions") or {"write_files": "allow", "run_command": "ask"}),
        "policy": normalize_policy(s.get("permission_policy")),
        "sandbox_mode": s.get("sandbox_mode", "required"),
        "sandbox_image": s.get("sandbox_image", "wenmo-agent:locked"),
        "sandbox": get_sandbox_status().__dict__,
    }


@app.post("/api/settings/permissions")
def set_permissions(req: PermissionRequest):
    """保存权限设置"""
    s = load_settings()
    perms = s.get("permissions") or {}
    if req.write_files in ("allow", "ask", "deny"):
        perms["write_files"] = req.write_files
    if req.run_command in ("allow", "ask", "deny"):
        perms["run_command"] = req.run_command
    s["permissions"] = perms
    if req.policy:
        s["permission_policy"] = normalize_policy(req.policy)
    if req.sandbox_mode in ("required", "prefer", "off"):
        s["sandbox_mode"] = req.sandbox_mode
    if req.sandbox_image.strip():
        s["sandbox_image"] = req.sandbox_image.strip()[:200]
    _atomic_json_write(_settings_file(), s)
    return {"ok": True, "permissions": perms,
            "policy": normalize_policy(s.get("permission_policy")),
            "sandbox_mode": s.get("sandbox_mode", "required")}


@app.get("/api/sandbox/status")
def sandbox_status():
    status = get_sandbox_status()
    settings = load_settings()
    return {**status.__dict__, "configured_mode": settings.get("sandbox_mode", "required"),
            "image": settings.get("sandbox_image", "wenmo-agent:locked")}


@app.get("/api/sandbox/diagnostics")
def sandbox_diagnostics():
    return get_sandbox_diagnostics()


class MathRequest(BaseModel):
    latex: str = ""


@app.post("/api/math/omml")
def math_to_omml(req: MathRequest):
    """LaTeX 公式 → OMML（Office 数学格式，可粘贴进 Word / MathType 成为可编辑公式）"""
    try:
        import latex2mathml.converter  # 延迟导入（符号表 200KB+，仅公式转换时加载）
        mathml = latex2mathml.converter.convert(req.latex)
        xslt = etree.XSLT(etree.parse(os.path.join(BASE_DIR, "MML2OMML.XSL")))
        omml = str(xslt(etree.fromstring(mathml)))
        return {"ok": True, "omml": omml}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"公式转换失败: {e}")


# ============================================================================
# 智能体（子 agent）：把子任务委托给独立模型并行处理，主模型专注对话
# 智能体的模型在「设置 → 图像模型与上下文 → 智能体」里选择；
# 没选 → 默认用当前对话模型（前端传 provider/model 覆盖）。
# 上下文联系（对齐 opencode）：同一对话内的多次委托共享智能体历史（AGENT_SESSIONS），
# 不同对话互相隔离 —— 由 contextvar 记录当前 conversation_id。
# ============================================================================

import contextvars as _contextvars

_agent_session_var = _contextvars.ContextVar("agent_session_id", default="")
_agent_model_var = _contextvars.ContextVar("agent_model_ctx", default=None)   # (provider, model)：当前主对话模型
AGENT_SESSIONS = {}            # session_id -> messages 列表（智能体上下文历史）
AGENT_SESSION_MAX = 24         # 每个会话最多保留多少轮（防上下文爆炸）
AGENT_SESSION_LIMIT = 256      # process-wide bound; old inactive sessions are evicted
_AGENT_SESSION_TOUCHED = {}
_AGENT_SESSION_LOCKS = {}
_AGENT_SESSIONS_LOCK = threading.RLock()


def _set_agent_session(sid):
    """chat 生成器里调用：记录当前对话，供智能体插件关联上下文"""
    _agent_session_var.set(sid or "")


def _set_agent_model(provider, model):
    """chat 生成器里调用：记录当前主对话模型（智能体未配置时默认跟随）"""
    _agent_model_var.set((provider or "", model or ""))


class AgentDelegateRequest(BaseModel):
    task: str                                  # 要交给智能体处理的任务
    context: str = ""                          # 可选背景信息（如当前文件内容/对话摘要）
    provider: str = ""                         # 覆盖：不传用 settings.agent_provider
    model: str = ""                            # 覆盖：不传用 settings.agent_model
    max_tokens: int = 0
    session: str = ""                          # 会话 key（留空 = 当前对话）
    files: list[str] = []                      # 主模型授予智能体读取的本地文件（只读，限 workspace/files/项目内）


def _read_agent_files(files):
    """读取主模型授予智能体的文件内容（权限由主模型给；只读；限定安全目录）。"""
    workspace = current_workspace.get() or os.path.join(BASE_DIR, "workspace")
    allowed = [os.path.abspath(workspace), os.path.abspath(_files_dir())]
    out = []
    for f in (files or []):
        name = str(f).strip().strip('"')
        if not name:
            continue
        # 相对路径 → 优先 workspace
        path = os.path.abspath(name)
        if not os.path.isfile(path):
            alt = os.path.join(workspace, name)
            if os.path.isfile(alt):
                path = os.path.abspath(alt)
        norm = os.path.abspath(path)
        try:
            in_scope = any(os.path.commonpath([norm, root]) == root for root in allowed)
        except ValueError:
            in_scope = False
        if not in_scope:
            out.append(f"【文件 {name}：不在授权目录内，已拒绝读取】")
            continue
        try:
            if os.path.getsize(path) > 200_000:
                out.append(f"【文件 {name}：过大，仅截取前 200KB】")
            with open(path, encoding="utf-8", errors="replace") as fp:
                content = fp.read(200_000)
            out.append(f"【文件 {os.path.basename(path)}】\n{content}")
        except Exception as e:
            out.append(f"【文件 {name}：读取失败 {e}】")
    return "\n\n".join(out)


def _resolve_agent_cfg(provider, model):
    """解析智能体模型配置：显式 > settings > 当前主对话 > 本地已加载。"""
    s = load_settings()
    prov = provider or s.get("agent_provider", "")
    mdl = model or s.get("agent_model", "")
    if prov and mdl:
        return prov, mdl
    # 智能体默认跟随主模型（用户要求：配置默认和主模型一样）
    cur = _agent_model_var.get()
    if cur and cur[0] and cur[1]:
        return cur
    if LOCAL_STATE["status"] == "ready" and LOCAL_STATE["name"]:
        return "local", LOCAL_STATE["name"]
    return "", ""


@app.post("/api/agent/delegate")
def agent_delegate(req: AgentDelegateRequest):
    """把任务委托给智能体模型，返回智能体的回答（非流式，一次完整推理）。
    同一 session（默认当前对话）内的委托共享上下文历史，不同对话互不影响。"""
    if not req.task.strip():
        raise HTTPException(status_code=400, detail="任务不能为空")
    provider, model = _resolve_agent_cfg(req.provider, req.model)
    if not provider or not model:
        raise HTTPException(status_code=400, detail="未配置智能体模型，且当前没有可继承的主对话模型")
    if provider == "local":
        if LOCAL_STATE["status"] != "ready" or not LOCAL_STATE["port"]:
            raise HTTPException(status_code=400, detail="本地智能体模型未就绪（请先在供应商页加载本地模型）")
        cfg = {"base_url": f"http://127.0.0.1:{LOCAL_STATE['port']}/v1", "model": LOCAL_STATE["name"]}
        api_key = "local"
    else:
        providers = load_providers()
        if provider not in providers:
            raise HTTPException(status_code=400, detail=f"未知智能体供应商: {provider}")
        cfg = providers[provider]
        api_key = resolve_key(cfg)
    client = OpenAI(base_url=cfg["base_url"], api_key=api_key or "local")

    # 会话上下文（opencode 式：同对话连续、跨对话隔离）
    raw_sid = (req.session.strip() or _agent_session_var.get() or "default")[:200]
    sid = "%s:%s" % (current_tenant.get(), raw_sid)
    with _AGENT_SESSIONS_LOCK:
        call_lock = _AGENT_SESSION_LOCKS.setdefault(sid, threading.Lock())
    # Same-session calls are serialized; a failed call never leaves an unmatched user turn.
    with call_lock:
        with _AGENT_SESSIONS_LOCK:
            hist = [dict(message) for message in AGENT_SESSIONS.get(sid, [])]
        if not hist:
            hist.append({"role": "system", "content": "你是主助手的子智能体（agent）。专注完成交给你的子任务，"
                                                      "直接给出结果或结论，不要寒暄，不要自称其他产品身份。"
                                                      "如果后续任务与之前的任务相关，请结合前面已完成的子任务继续。"})
        user_content = f"任务：{req.task}"
        if req.context:
            user_content += f"\n\n背景信息：{req.context}"
        if req.files:
            user_content += "\n\n【主模型授权读取的本地文件】\n" + _read_agent_files(req.files)
        hist.append({"role": "user", "content": user_content})
        if len(hist) > AGENT_SESSION_MAX * 2 + 1:
            hist[:] = [hist[0]] + hist[-(AGENT_SESSION_MAX * 2):]

        kwargs = {} if not req.max_tokens else {"max_tokens": req.max_tokens}
        kwargs.setdefault("timeout", 120)
        try:
            resp = client.chat.completions.create(model=model, messages=hist, **kwargs)
            text = (resp.choices[0].message.content or "").strip()
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"智能体调用失败: {e}")
        hist.append({"role": "assistant", "content": text})

        usage_delta = from_openai_usage(getattr(resp, "usage", None))
        if usage_delta:
            if raw_sid != "default" and any(usage_delta.values()):
                history_store.add_conversation_usage(
                    raw_sid, usage_delta, provider=provider, model=model)

        with _AGENT_SESSIONS_LOCK:
            AGENT_SESSIONS[sid] = hist
            _AGENT_SESSION_TOUCHED[sid] = time.monotonic()
            while len(AGENT_SESSIONS) > AGENT_SESSION_LIMIT:
                oldest = min(_AGENT_SESSION_TOUCHED, key=_AGENT_SESSION_TOUCHED.get)
                if oldest == sid and len(AGENT_SESSIONS) > 1:
                    candidates = {key: value for key, value in _AGENT_SESSION_TOUCHED.items()
                                  if key != sid}
                    oldest = min(candidates, key=candidates.get)
                AGENT_SESSIONS.pop(oldest, None)
                _AGENT_SESSION_TOUCHED.pop(oldest, None)
                _AGENT_SESSION_LOCKS.pop(oldest, None)
        return {"ok": True, "provider": provider, "model": model, "result": text,
                "usage": usage_delta, "session": sid,
                "session_rounds": (len(hist) - 1) // 2}


def _agent_delegate_from_plugin(arguments):
    """Preserve tenant/conversation ContextVars by staying inside this process."""
    try:
        return agent_delegate(AgentDelegateRequest(**dict(arguments or {})))
    except HTTPException as exc:
        return {"error": str(exc.detail), "detail": str(exc.detail)}
    except Exception as exc:
        return {"error": str(exc)}


from agent_bridge import register_delegate as _register_agent_delegate
_register_agent_delegate(_agent_delegate_from_plugin)


# ============================================================================
# 小应用（小挂件 / 浮窗 / 桌宠）：模型可创建，图标放顶栏，点击打开
# ============================================================================
# 小应用目录：打包版数据目录 content/apps 优先（用户创建的应用不随更新覆盖），seed/apps 兜底
if CONTENT_DIR:
    _c_apps = os.path.join(CONTENT_DIR, "apps")
    if not os.path.isdir(_c_apps):
        try:
            _seed_apps = os.path.join(os.environ.get("WENMO_RES_DIR", ""), "seed", "apps")
            if os.path.isdir(_seed_apps):
                shutil.copytree(_seed_apps, _c_apps)
        except Exception:
            pass
    APPS_DIR = _c_apps
else:
    APPS_DIR = os.path.join(BASE_DIR, "apps")
os.makedirs(APPS_DIR, exist_ok=True)


class AppSaveRequest(BaseModel):
    name: str        # 应用标识（英文/数字/下划线）
    title: str = ""  # 显示名
    icon: str = ""   # emoji 图标
    html: str = ""   # 应用 HTML（模型生成的小挂件/浮窗/桌宠）


@app.get("/api/apps")
def list_apps():
    """列出全部小应用（目录里的 + 内置）"""
    apps = []
    seen = set()
    try:
        roots = [_user_apps_dir()]
        if os.path.normcase(os.path.abspath(APPS_DIR)) != os.path.normcase(os.path.abspath(roots[0])):
            roots.append(APPS_DIR)
        for root in roots:
            for entry in sorted(os.scandir(root), key=lambda e: e.name):
                if entry.name in seen:
                    continue
                seen.add(entry.name)
                if not entry.is_dir():
                    continue
                meta = {}
                meta_path = os.path.join(entry.path, "meta.json")
                if os.path.isfile(meta_path):
                    try:
                        with open(meta_path, encoding="utf-8") as f:
                            meta = json.load(f)
                    except Exception:
                        meta = {}
                if not os.path.isfile(os.path.join(entry.path, "index.html")):
                    continue
                apps.append({
                    "name": entry.name,
                    "title": meta.get("title") or entry.name,
                    "icon": meta.get("icon") or "🧩",
                    "builtin": root == APPS_DIR and bool(meta.get("builtin")),
                    "frameless": bool(meta.get("frameless")),
                    "width": meta.get("width") or 0,
                    "height": meta.get("height") or 0,
                })
    except OSError:
        pass
    return {"apps": apps}


class AppCreateRequest(BaseModel):
    name: str
    title: str = ""
    icon: str = ""
    html: str = ""


@app.post("/api/apps/create")
def create_app(req: AppCreateRequest):
    """模型创建小应用（小挂件/浮窗/桌宠）：写 apps/<name>/index.html + meta.json"""
    name = req.name.strip()
    if not name or not re.match(r"^[a-zA-Z0-9_-]{1,40}$", name):
        raise HTTPException(status_code=400, detail="应用名只能包含字母、数字、下划线、短横线（≤40）")
    if len(req.html) > 200_000:
        raise HTTPException(status_code=400, detail="应用 HTML 过大（≤200KB）")
    app_dir = os.path.join(_user_apps_dir(), name)
    os.makedirs(app_dir, exist_ok=True)
    with open(os.path.join(app_dir, "index.html"), "w", encoding="utf-8") as f:
        f.write(req.html)
    meta = {"name": name, "title": (req.title or name)[:50], "icon": (req.icon or "🧩")[:4]}
    with open(os.path.join(app_dir, "meta.json"), "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    return {"ok": True, **meta}


@app.delete("/api/apps/{name}")
def delete_app(name: str):
    """删除一个小应用（内置/模型创建的都可以删）"""
    if not re.match(r"^[a-zA-Z0-9_-]{1,40}$", name):
        raise HTTPException(status_code=400, detail="非法应用名")
    app_dir = os.path.join(_user_apps_dir(), name)
    if not os.path.isdir(app_dir):
        raise HTTPException(status_code=404, detail="应用不存在")
    import shutil
    shutil.rmtree(app_dir)
    return {"ok": True}



@app.get("/api/stock/kline")
def stock_kline(code: str = "sh600519", period: str = "day", count: int = 400):
    """股票 K线代理：转发腾讯 fqkline 接口（服务端无跨域限制）"""
    count = max(10, min(int(count), 800))
    url = "https://web.ifzq.gtimg.cn/appstock/app/fqkline/get?param={},{},,,{},qfq".format(code, period, count)
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0", "Referer": "https://gu.qq.com/"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json.loads(resp.read().decode("utf-8", "ignore"))
        node = (data.get("data") or {}).get(code) or {}
        rows = (node.get("qfq" + period) or node.get(period)
                or node.get("qfqday") or node.get("day") or [])
        if not rows:
            return {"ok": False, "error": "暂无K线数据"}
        return {"ok": True, "rows": rows}
    except Exception as e:
        return {"ok": False, "error": "K线获取失败: {}".format(e)}


@app.get("/api/stock/quote")
def stock_quote(codes: str = "sh600519"):
    """股票实时行情代理：转发腾讯 qt 接口（GBK 解码）"""
    codes = codes.strip()
    if not codes or len(codes) > 200:
        return {"ok": False, "error": "参数错误"}
    url = "https://qt.gtimg.cn/q=" + codes
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0", "Referer": "https://gu.qq.com/"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            raw = resp.read().decode("gbk", "ignore")
        out = {}
        for m in re.finditer(r'v_([a-z]{2}\d{6})="([^"]*)"', raw):
            f = m.group(2).split("~")
            if len(f) < 40:
                continue
            code = m.group(1)
            out[code] = {
                "name": f[1], "code": f[2], "price": float(f[3] or 0), "prevClose": float(f[4] or 0),
                "open": float(f[5] or 0), "vol": int(f[6] or 0), "time": f[30],
                "change": float(f[31] or 0), "changePct": float(f[32] or 0),
                "high": float(f[33] or 0), "low": float(f[34] or 0),
                "vol2": int(f[36] or 0), "amountWan": float(f[37] or 0),
                "turnover": float(f[38] or 0), "pe": float(f[39] or 0),
                "amp": float(f[43] or 0), "floatCap": float(f[44] or 0),
                "totalCap": float(f[45] or 0), "pb": float(f[46] or 0),
                "limitUp": float(f[47] or 0), "limitDown": float(f[48] or 0),
            }
        if not out:
            return {"ok": False, "error": "未获取到行情数据"}
        return {"ok": True, "quotes": out}
    except Exception as e:
        return {"ok": False, "error": "行情获取失败: {}".format(e)}


@app.get("/api/stock/search")
def stock_search(q: str = ""):
    """股票搜索联想代理：转发腾讯 smartbox 接口（GBK + unicode 反转义）"""
    import urllib.parse
    q = q.strip()
    if not q or len(q) > 20:
        return {"ok": True, "items": []}
    url = "https://smartbox.gtimg.cn/s3/?v=2&q=" + urllib.parse.quote(q) + "&t=all"
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0", "Referer": "https://gu.qq.com/"})
        with urllib.request.urlopen(req, timeout=8) as resp:
            raw = resp.read().decode("gbk", "ignore")
        m = re.search(r'v_hint="(.*?)"', raw, re.S)
        if not m:
            return {"ok": True, "items": []}
        content = m.group(1).encode("latin-1", "ignore").decode("unicode_escape")
        items = []
        for seg in content.split("^"):
            f = seg.split("~")
            if len(f) >= 3 and f[0] and f[1]:
                items.append({"code": f[0] + f[1], "name": f[2]})
        return {"ok": True, "items": items[:10]}
    except Exception:
        return {"ok": True, "items": []}



@app.get("/api/files/docx/{name}")
def docx_preview(name: str):
    """解析 docx → 结构化内容（标题/段落/列表/表格），供前端文档预览"""
    if not re.match(r"^[\w\u4e00-\u9fff\s（）()\-_]+\.(docx|doc)$", name):
        raise HTTPException(status_code=400, detail="非法文件名")
    path = os.path.join(_files_dir(), name)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="文件不存在")
    try:
        import docx as _docx
        from docx.document import Document as _Doc
        from docx.table import Table as _Table
        from docx.text.paragraph import Paragraph as _Paragraph
    except ImportError:
        raise HTTPException(status_code=500, detail="服务器缺少 python-docx，无法预览 Word")
    try:
        doc = _docx.Document(path)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法解析 Word: {e}")

    def _iter_blocks(parent):
        from docx.oxml.ns import qn
        body = parent.element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                p = _Paragraph(child, parent)
                style = (p.style.name or "").lower() if p.style else ""
                text = p.text.strip()
                # 提取 OMML 公式（m:oMath 内 m:t 文本）——预览展示公式内容
                formula = ""
                try:
                    from docx.oxml.ns import nsmap
                    ms = child.findall(".//" + qn("m:t"))
                    if ms:
                        formula = "".join(t.text or "" for t in ms).strip()
                except Exception:
                    pass
                if not text and not formula:
                    continue
                if formula:
                    yield {"type": "formula", "text": formula}
                    continue
                if style.startswith("heading") or style.startswith("标题"):
                    level = 1
                    m = re.search(r"(\d)", style)
                    if m:
                        level = min(4, int(m.group(1)))
                    yield {"type": "heading", "level": level, "text": text}
                elif style.startswith(("list", "列表")):
                    yield {"type": "list", "text": text}
                else:
                    yield {"type": "paragraph", "text": text}
            elif child.tag == qn("w:tbl"):
                t = _Table(child, parent)
                rows = []
                for row in t.rows[:200]:
                    rows.append([cell.text.strip() for cell in row.cells][:30])
                yield {"type": "table", "rows": rows}
    blocks = []
    for b in _iter_blocks(doc):
        blocks.append(b)
        if len(blocks) > 2000:
            blocks.append({"type": "paragraph", "text": "…（内容较多，已截断）"})
            break
    return {"ok": True, "name": name, "blocks": blocks}


@app.get("/api/files/preview/{name}")
def file_preview_pdf(name: str):
    """docx/pptx/xlsx → PDF（win32com 系统 Office 导出）：真正的 Office 渲染预览（复用 PDF 查看器）。
    转换缓存：files/_preview_cache/ 下同名 PDF（命中直接返回，秒开）。"""
    # 文件名允许：字母数字 中文 空格 全角括号 下划线连字符等（避免含空格/（）的文件名被拒）
    if not re.match(r"^[\w\u4e00-\u9fff\s（）()\-_]+\.(docx|doc|pptx|ppt|xlsx|xls)$", name):
        raise HTTPException(status_code=400, detail="不支持该格式预览")
    path = os.path.join(_files_dir(), name)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="文件不存在")
    ext = os.path.splitext(name)[1].lower().lstrip(".")
    cache_dir = os.path.join(_files_dir(), "_preview_cache")
    os.makedirs(cache_dir, exist_ok=True)
    base = os.path.splitext(name)[0]
    pdf_path = os.path.join(cache_dir, base + ".pdf")
    # 缓存命中：PDF 比源文件新（或相差 5 秒内，兼容文件系统时间戳微差）
    # → 直接用缓存秒开，避免每次打开都重新启动 Office 转换（10-40 秒）
    if os.path.isfile(pdf_path) and os.path.getsize(pdf_path) > 0 \
            and os.path.getmtime(pdf_path) >= os.path.getmtime(path) - 5:
        return FileResponse(pdf_path, media_type="application/pdf")
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()   # COM 初始化（线程池线程必需，否则 -2147221008 CoInitialize 未调用）
        try:
            if ext in ("docx", "doc"):
                app = win32com.client.Dispatch("Word.Application")
                app.Visible = False
                try:
                    doc = app.Documents.Open(path, ReadOnly=True)
                    doc.SaveAs(pdf_path, FileFormat=17)
                    doc.Close(False)
                finally:
                    app.Quit()
            elif ext in ("pptx", "ppt"):
                app = win32com.client.Dispatch("PowerPoint.Application")
                try:
                    prs = app.Presentations.Open(path, ReadOnly=True, WithWindow=False)
                    prs.SaveAs(pdf_path, 32)   # ppSaveAsPDF
                    prs.Close()
                finally:
                    app.Quit()
            elif ext in ("xlsx", "xls"):
                app = win32com.client.Dispatch("Excel.Application")
                app.Visible = False
                app.DisplayAlerts = False
                try:
                    wb = app.Workbooks.Open(path, ReadOnly=True)
                    wb.ExportAsFixedFormat(0, pdf_path)
                    wb.Close(False)
                finally:
                    app.Quit()
        finally:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
    except Exception as e:
        raise HTTPException(status_code=500, detail="Office 转换失败：%s" % str(e)[:150])
    if not os.path.isfile(pdf_path) or os.path.getsize(pdf_path) == 0:
        raise HTTPException(status_code=500, detail="Office 转换未产出 PDF")
    return FileResponse(pdf_path, media_type="application/pdf")


@app.get("/api/files/xlsx/{name}")
def xlsx_preview(name: str):
    """解析 xlsx → 工作表数据（表头+行），供前端表格预览"""
    if not re.match(r"^[\w\u4e00-\u9fff\s（）()\-_]+\.(xlsx|xls)$", name):
        raise HTTPException(status_code=400, detail="非法文件名")
    path = os.path.join(_files_dir(), name)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="文件不存在")
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except ImportError:
        raise HTTPException(status_code=500, detail="服务器缺少 openpyxl，无法预览 Excel")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法解析 Excel: {e}")
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                rows.append(vals)
            if len(rows) >= 300:
                rows.append(["(行数较多，已截断)"])
                break
        sheets.append({"name": ws.title, "rows": rows})
    return {"ok": True, "name": name, "sheets": sheets}


@app.get("/api/system/stats")
def system_stats():
    """系统状态（小挂件用）：CPU/内存/GPU/本地模型"""
    import psutil
    mem = _ram_info()
    vram = _vram_info()
    return {
        "cpu": psutil.cpu_percent(interval=0.4),
        "ram_used": mem["used"], "ram_total": mem["total"],
        "vram": vram,
        "local_model": LOCAL_STATE["name"] or None,
        "local_status": LOCAL_STATE["status"],
    }


@app.get("/api/files/pptx/{name}")
def pptx_preview(name: str):
    """解析 pptx → 幻灯片列表（PPT 界面预览用）：每页文本 + 图片（base64）"""
    if not re.match(r"^[\w\u4e00-\u9fff\s（）()\-_]+\.(pptx|ppt)$", name):
        raise HTTPException(status_code=400, detail="非法文件名")
    path = os.path.join(_files_dir(), name)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="文件不存在")
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(path)
    except ImportError:
        raise HTTPException(status_code=500, detail="服务器缺少 python-pptx，无法预览 PPT")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法解析 PPT: {e}")
    slides = []
    for slide in prs.slides:
        texts = []
        imgs = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.image:
                    img = shape.image
                    if len(img.blob) < 2_000_000:
                        import base64
                        imgs.append("data:" + img.content_type + ";base64," +
                                    base64.b64encode(img.blob).decode())
            except Exception:
                continue
        slides.append({"texts": texts[:12], "imgs": imgs[:4]})
    return {"ok": True, "count": len(slides), "slides": slides}


PENDING_ASKS = {}   # ask_user 弹窗询问：qid -> asyncio.Future（等前端回答）


class AskAnswerRequest(BaseModel):
    answer: str = ""


@app.post("/api/ask/{qid}")
async def ask_answer(qid: str, req: AskAnswerRequest):
    """前端回答 ask_user 弹窗后回调这里，唤醒等待中的模型"""
    fut = PENDING_ASKS.get(qid)
    if fut is None or fut.done():
        raise HTTPException(status_code=404, detail="问题不存在或已过期")
    fut.set_result(req.answer)
    return {"ok": True}


# ============================================================================
# 图片附件上传：前端把用户发的图先存到 files/，模型通过 see_image 工具读图
# （避免把 data URL 直接塞给不支持视觉的模型导致报错）
# ============================================================================

class UploadRequest(BaseModel):
    data_url: str   # data:image/png;base64,.....
    filename: str = ""   # 原始文件名（方案B：MIME 映射不到时用扩展名兜底）


@app.post("/api/upload")
def upload_attachment(req: UploadRequest):
    import re as _re
    m = _re.match(r"^data:([\w\-.+]+/[\w\-.+]+);base64,(.+)$", req.data_url.strip(), _re.S)
    if not m:
        raise HTTPException(status_code=400, detail="不支持的 base64 data URL（需 data:<类型>;base64,<数据>）")
    media, b64 = m.group(1), m.group(2)
    if len(b64) > 50_000_000:
        raise HTTPException(status_code=400, detail="文件过大（>50MB）")
    # MIME → 扩展名（图片 + 文档 + 压缩包 + 任意兜底）：支持任意文件上传
    # MIME → 扩展名（图片 + 文档 + 代码 + 文本 + 压缩 + 音视频 + 任意兜底）：支持任意文件上传
    ext = {
        # 图片
        "image/png": "png", "image/jpeg": "jpg", "image/jpg": "jpg", "image/gif": "gif",
        "image/webp": "webp", "image/bmp": "bmp", "image/svg+xml": "svg", "image/tiff": "tiff",
        # 文档
        "application/pdf": "pdf",
        "application/msword": "doc",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.ms-excel": "xls",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.ms-powerpoint": "ppt",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/rtf": "rtf",
        # 文本 / 代码 / 标记
        "text/plain": "txt", "text/markdown": "md", "text/html": "html", "text/csv": "csv",
        "text/css": "css", "text/xml": "xml", "application/xml": "xml", "application/json": "json",
        "text/javascript": "js", "application/javascript": "js",
        "text/x-python": "py", "text/x-script.python": "py", "text/x-shellscript": "sh",
        "application/x-sh": "sh", "application/x-bat": "bat",
        "text/yaml": "yaml", "application/yaml": "yaml", "text/x-yaml": "yaml",
        "text/x-ini": "ini", "text/x-config": "conf", "text/x-c": "c", "text/x-c++": "cpp",
        "text/x-java-source": "java", "text/x-ruby": "rb", "text/x-php": "php",
        "text/x-go": "go", "text/x-rust": "rs", "text/x-typescript": "ts",
        # 压缩 / 归档
        "application/zip": "zip", "application/x-zip-compressed": "zip",
        "application/x-7z-compressed": "7z", "application/x-rar-compressed": "rar",
        "application/gzip": "gz", "application/x-tar": "tar",
        # 可执行 / 其他
        "application/x-msdownload": "exe", "application/vnd.microsoft.portable-executable": "exe",
        "application/x-msdos-program": "exe", "application/x-iso9660-image": "iso",
        # 音视频
        "video/mp4": "mp4", "video/x-msvideo": "avi", "video/quicktime": "mov",
        "video/webm": "webm", "video/mpeg": "mpeg",
        "audio/mpeg": "mp3", "audio/wav": "wav", "audio/x-wav": "wav",
        "audio/ogg": "ogg", "audio/flac": "flac", "audio/aac": "aac", "audio/mp4": "m4a",
        # 兜底：任意未列出的类型也保留原始扩展名（从文件名提取，见下方）
    }
    # 先用 MIME 映射；映射不到再用原始文件名扩展名（方案B：前端带 filename）
    _ext = ext.get(media, "")
    if not _ext:
        _fn = (req.filename or "").strip() if hasattr(req, "filename") and req.filename else ""
        if _fn:
            _fe = os.path.splitext(_fn)[1].lstrip(".").lower()
            if _fe and len(_fe) <= 6:
                _ext = _fe
    if not _ext:
        _ext = "bin"
    try:
        raw = base64.b64decode(b64)
    except Exception:
        raise HTTPException(status_code=400, detail="文件数据损坏")
    name = f"att_{int(time.time())}_{os.urandom(3).hex()}.{_ext}"
    with open(os.path.join(_files_dir(), name), "wb") as f:
        f.write(raw)
    return {"ok": True, "name": name, "url": f"/files/{name}"}


# 核心保底工具：无论消息是什么都注入（常用/安全）
# ============ 免费综合供应商路由表（免费综合 = 聚合各平台免费档，非单一平台） ============
# 模型显示名 → 目标供应商 key（字符串）或 {provider, model}（显示名 ≠ 真实模型名时用 dict）
# 目标供应商需在 providers.json 中配置（用户去各平台注册拿免费 key 填到对应供应商）
FREE_MODEL_ROUTES = {
    # DeepSeek 免费档（已有 key 立即可用）
    "deepseek-v4-flash-free": "deepseek",
    "deepseek-chat": "deepseek",
    "deepseek-reasoner": "deepseek",
    # 智谱 GLM 官方免费模型（glm-4.7-flash 永久免费）
    "glm-4.7-flash": "zhipu",
    # 硅基流动限免模型（DeepSeek-OCR / 腾讯混元 MT-7B）
    "deepseek-ai/DeepSeek-OCR": "siliconflow",
    "tencent/Hunyuan-MT-7B": "siliconflow",
    # Kimi 免费档（kimi-latest 有免费额度）
    "kimi-latest": "kimi",
    "kimi-k3": "kimi",
    # 通义免费档（qwen-turbo 便宜/有免费额度）
    "qwen-turbo": "qianwen",
    "qwen-plus": "qianwen",
    # 豆包免费档
    "doubao-lite-32k": "doubao",
}

CORE_TOOL_NAMES = {
    "ask_user", "plugin_deliver_write_file_with_link", "plugin_workspace_read_file",
    "plugin_workspace_write_file", "plugin_system_info_get_system_info",
    "plugin_vision_see_image",
    "plugin_terminal_run_command",   # 终端命令：按项目常用，对标 opencode 终端工具
    "websearch_web_search",  # 联网搜索：MCP 工具（server_tool 前缀），核心保底
    "plugin_ppt_create_pptx", "plugin_word_create_docx", "plugin_excel_create_xlsx",  # 文档工具常用
    "plugin_skill_tools_skill_list", "plugin_skill_tools_skill_load", "plugin_skill_tools_skill_overview",  # 技能库按需加载
}

# 工具列表缓存：按对话锁定（缓存命中率优化——tools 稳定 → 前缀缓存最大化）
_TOOLS_CACHE = {}

# 前缀指纹（对标 Reasonix ImmutablePrefix 的 SHA-256 钉住）：对话 -> system+tools 前缀哈希。
# 同对话内前缀漂移（system 内容变、工具集变）→ 缓存必然全 miss；指纹用于检测并记录漂移原因。
_PREFIX_FINGERPRINTS = {}

# ============ 工具结果缓存（已按用户要求移除：与云端一致，所有工具调用真实执行） ============
# 背景：此前对读类工具做了 60 秒 LRU 本地缓存（命中直接复用结果、不发请求），
# 导致重复搜索/读文件不真实执行、云端 usage 与日志缺失，行为与云端不一致。
# 用户明确要求"问墨不应该自己搞一套，应该和云端的一致"→ 已禁用。
# 保留以下定义仅为兼容（无任何调用点），如需恢复可放开注释并还原 _exec_tool 逻辑。
# ============ 教训机制：模型从对话中学习，少犯错 ============
# 用户纠正词：检测到 → 记一条教训（模型以后避免同样错误）
CORRECTION_HINTS = ("不是", "不对", "错了", "错啦", "错误", "重新做", "重做", "别这样",
                    "别用", "不要这样", "应该改成", "应该这样", "理解错", "搞错", "反了",
                    "不对啊", "怎么又", "又错了", "改一下", "修正", "更正", "听我说")


def _load_lessons():
    data = load_tenant_json("lessons.json", [])
    return data if isinstance(data, list) else []


def _save_lesson(text, source="user_correction"):
    """保存一条教训（最多保留 50 条）"""
    try:
        lessons = _load_lessons()
        lessons.append({"text": text[:200], "source": source, "ts": time.time()})
        lessons = lessons[-50:]
        write_tenant_json("lessons.json", lessons)
    except Exception:
        pass


# ============ 技能使用统计（自演化 P2/P3：usage 驱动生命周期 + 反馈驱动自适应） ============
_USAGE_CACHE = {}          # tenant -> name -> {"uses": n, "last": ts}
_USAGE_FLUSH_COUNT = {}    # tenant -> count
_LAST_INJECTED = {}        # tenant -> last injected skill names


def _record_skill_usage(names):
    """记录技能被激活（注入到对话）——驱动试用期毕业/归档"""
    if not names:
        return
    import time as _t
    tenant = tenant_name()
    usage_cache = _USAGE_CACHE.setdefault(tenant, {})
    for n in names:
        u = usage_cache.setdefault(n, {"uses": 0, "last": 0})
        u["uses"] += 1
        u["last"] = _t.time()
    _LAST_INJECTED[tenant] = list(names)
    _USAGE_FLUSH_COUNT[tenant] = _USAGE_FLUSH_COUNT.get(tenant, 0) + 1
    # Persist every activation. This is small state and avoids tenant data being
    # stranded in process memory when the app exits outside a request context.
    _flush_usage_cache()
    if _USAGE_FLUSH_COUNT[tenant] % 10 == 0:   # 每 10 次激活写盘一次（降 IO）
        _flush_usage_cache()


def _record_skill_feedback(positive):
    """反馈闭环（P3）：用户纠正 → 上一轮注入的技能负反馈；任务正常 → 正反馈。
    score 用于注入优先级（低频技能淘汰）"""
    tenant = tenant_name()
    last_injected = _LAST_INJECTED.get(tenant, [])
    if not last_injected:
        return
    stats = load_tenant_json("usage_stats.json", {})
    stats = stats if isinstance(stats, dict) else {}
    for n in last_injected:
        s = stats.setdefault(n, {"uses": 0, "last": 0, "created": 0, "score": 0})
        s["score"] = s.get("score", 0) + (1 if positive else -1)
    try:
        write_tenant_json("usage_stats.json", stats)
    except Exception:
        pass
    _LAST_INJECTED[tenant] = []


def _skill_priority(names):
    """P3：按历史反馈分排序候选（高分优先注入；低分 < -3 淘汰）"""
    stats = load_tenant_json("usage_stats.json", {})
    if not isinstance(stats, dict):
        return names, 0
    scored = [(n, stats.get(n, {}).get("score", 0)) for n in names]
    kept = [n for n, s in scored if s >= -3]
    kept.sort(key=lambda n: -stats.get(n, {}).get("score", 0))
    return kept, len(kept) != len(names)


def _flush_usage_cache():
    try:
        tenant = tenant_name()
        stats = load_tenant_json("usage_stats.json", {})
        stats = stats if isinstance(stats, dict) else {}
        usage_cache = _USAGE_CACHE.get(tenant, {})
        for n, u in usage_cache.items():
            prev = stats.get(n, {"uses": 0, "last": 0, "created": 0})
            stats[n] = {"uses": prev.get("uses", 0) + u["uses"], "last": u["last"],
                        "created": prev.get("created", 0) or time.time()}
        write_tenant_json("usage_stats.json", stats)
        usage_cache.clear()
    except Exception:
        pass


def _lesson_system_prompt(query=""):
    """记忆图跨会话召回：把"最近教训" + "与当前问题语义最相关历史记忆"注入 system。
    借鉴 jcode Agent Memory——不再是简单取尾部 5 条，而是 TF-IDF 语义召回。"""
    parts = []
    # ① 最近的教训（用户纠正/工具失败，近 3 条——时间上最近的一般最相关）
    lessons = _load_lessons()
    if lessons:
        recent = lessons[-3:]
        _lines = []
        for l in recent:
            src = "用户纠正" if l.get("source") == "user_correction" else "工具失败"
            _lines.append(f"- [{src}] {l.get('text', '')[:150]}")
        if _lines:
            parts.append("【经验教训：以下是你之前犯过的错/用户纠正过的事，务必吸取，避免重犯】\n" + "\n".join(_lines))
    # ② 语义召回的历史记忆（跨对话，按当前问题相关度）
    if query:
        try:
            import memory_graph
            mem_prompt = memory_graph.memory_system_prompt(query, top_k=4)
            if mem_prompt:
                parts.append(mem_prompt)
        except Exception:
            pass   # 记忆图故障不影响主流程
    return "\n\n".join(parts)


def _match_tools_by_message(tools, message, core=CORE_TOOL_NAMES, max_tools=32):
    """按最后一条用户消息匹配工具描述（中文二元组 + 英文关键词）。
    对标 opencode 智能工具注入：只给相关工具，减少无关工具对模型的干扰（防降智）。"""
    if not tools:
        return tools
    message = message or ""
    grams = set()
    for i in range(len(message) - 1):
        pair = message[i:i + 2]
        if all("\u4e00" <= c <= "\u9fff" for c in pair):
            grams.add(pair)
    for w in re.findall(r"[a-zA-Z][a-zA-Z0-9_-]{2,}", message.lower()):
        grams.add("w:" + w)
    if not grams:
        fallback = [t for t in tools if t["function"]["name"] in core] or tools
        return fallback[:max_tools]

    scored = []
    for t in tools:
        name = t["function"]["name"]
        if name in core:
            scored.append((100, t))          # 核心保底永远在
            continue
        hay = (name + " " + t["function"].get("description", "")).lower()
        overlap = 0
        for g in grams:
            if g in hay or (g.startswith("w:") and g[2:] in hay):
                overlap += 1
        if overlap:
            scored.append((10 + overlap, t))
    scored.sort(key=lambda x: -x[0])
    out = [t for _, t in scored]
    # 冲突消解：deliver 的 write_file_with_link（给链接）在场时，
    # 移除 workspace 的 write_file（不给链接）——避免小模型选错、用户拿不到链接
    names = [t["function"]["name"] for t in out]
    if "plugin_deliver_write_file_with_link" in names and "plugin_workspace_write_file" in names:
        wants_download = any(hint in (message or "").lower() for hint in (
            "下载", "发给我", "给我文件", "导出", "download", "export"))
        remove_name = ("plugin_workspace_write_file" if wants_download
                       else "plugin_deliver_write_file_with_link")
        out = [t for t in out if t["function"]["name"] != remove_name]
    return out[:max_tools]


class ChatRequest(BaseModel):
    provider: str
    messages: list[dict]  # [{"role": "user", "content": "..."}, ...]
    conversation_id: str = ""  # 前端传来的会话 id（用于历史保存）
    project: str = "default"   # 所属项目
    model: str = ""            # 前端指定的模型（覆盖 providers.json 配置）
    online: bool = False       # 联网模式：优先保留搜索工具并提示 AI
    reasoning: str = ""        # 深度思考：low / medium / high（reasoning_effort）
    mode: str = ""             # plan（计划模式：只读规划+确认后执行）/ build（默认执行）
    persist: bool = True       # 对比/临时请求设为 False：不覆盖正文历史，但仍合并实际用量
    request_id: str = ""       # 前端生成；用于跨 HTTP 请求精确取消服务端流任务


_ACTIVE_CHAT_STREAMS = {}
_CANCELLED_CHAT_STREAMS = {}
_ACTIVE_CHAT_STREAMS_LOCK = threading.RLock()


@app.post("/api/chat/{request_id}/cancel")
async def cancel_chat_stream(request_id: str):
    """Cancel the actual server-side generator, not only the browser fetch."""
    if not re.fullmatch(r"[A-Za-z0-9._-]{8,128}", request_id or ""):
        raise HTTPException(status_code=400, detail="无效的流请求 id")
    stream_key = (current_tenant.get(), request_id)
    with _ACTIVE_CHAT_STREAMS_LOCK:
        task = _ACTIVE_CHAT_STREAMS.get(stream_key)
        _CANCELLED_CHAT_STREAMS[stream_key] = time.monotonic()
    if task is None or task.done():
        return {"ok": True, "cancelled": False}
    task.cancel()
    return {"ok": True, "cancelled": True}


# ============ 工具 Schema 展平（对标 Reasonix：DeepSeek 对 >2 层嵌套或 >10 叶子字段的
# schema 会静默丢参 → 展平为点号记法 a.b.c，调用时再还原嵌套结构）============
def _schema_leaf_count(p):
    """统计 schema 的叶子字段数（object 递归；array 算 items；非对象算 1）"""
    if not isinstance(p, dict):
        return 1
    t = p.get("type")
    if t == "object":
        sub = p.get("properties") or {}
        if not sub:
            return 1
        return sum(_schema_leaf_count(v) for v in sub.values())
    if t == "array":
        it = p.get("items")
        return _schema_leaf_count(it) if isinstance(it, dict) else 1
    return 1


def _schema_max_depth(p, d=1):
    """object 嵌套深度（数组不加深——数组内对象不展平）"""
    if not isinstance(p, dict):
        return d
    t = p.get("type")
    if t == "object":
        sub = p.get("properties") or {}
        return max((_schema_max_depth(v, d + 1) for v in sub.values()), default=d)
    return d


def _should_flatten(schema):
    """>2 层嵌套 或 >10 叶子字段 → 需要展平（Reasonix 阈值）"""
    if not isinstance(schema, dict):
        return False
    props = schema.get("properties") or {}
    if not isinstance(props, dict) or not props:
        return False
    if _schema_max_depth(schema) > 2:
        return True
    return sum(_schema_leaf_count(v) for v in props.values()) > 10


def _flatten_schema_params(schema):
    """把深嵌套 object 展平为点号键。返回 (新 schema, 映射{扁平键: 路径segments})。
    数组内的对象不展平（保持 items 结构），避免元素级还原的复杂度。"""
    props = schema.get("properties") or {}
    new_props = {}
    mapping = {}
    required = schema.get("required") or []

    def walk(prefix, p, path, parent_req):
        if not isinstance(p, dict):
            new_props[prefix] = p
            return
        t = p.get("type")
        if t == "object":
            sub = p.get("properties") or {}
            sub_req = p.get("required") or []
            if not sub:
                new_props[prefix] = {"type": "object"}
                return
            for k, v in sub.items():
                seg = path + [k]
                walk(f"{prefix}.{k}", v, seg, sub_req)
        else:
            new_props[prefix] = p
            mapping[prefix] = path

    for k, v in props.items():
        seg = [k]
        if isinstance(v, dict) and v.get("type") == "object" and (v.get("properties") or {}):
            # object 字段：递归展平其内部（无论深度，统一展平成点号）
            for k2, v2 in (v.get("properties") or {}).items():
                walk(f"{k}.{k2}", v2, seg + [k2], v.get("required") or [])
        else:
            new_props[k] = v
            mapping[k] = seg

    new_required = []
    for r in required:
        if r in mapping and mapping[r] == [r]:
            new_required.append(r)
        else:
            # 嵌套 required：展平后所有子键都必填（保守：不强制）
            pass
    out = {"type": "object", "properties": new_props}
    if new_required:
        out["required"] = new_required
    return out, mapping


def _renest_args(args, mapping):
    """把展平后的点号键还原为嵌套结构（对标 Reasonix dispatch 时 re-nest）"""
    if not mapping or not isinstance(args, dict):
        return args
    out = {}
    for k, v in args.items():
        path = mapping.get(k)
        if path is None or not path:
            out[k] = v
            continue
        cur = out
        for seg in path[:-1]:
            cur = cur.setdefault(seg, {})
        cur[path[-1]] = v
    return out


def _flatten_tools(tools):
    """Flatten schemas and return a request-local re-nesting map.

    A process-global map is unsafe: two concurrent conversations can otherwise
    overwrite each other's argument mapping before tool dispatch.
    """
    flatten_map = {}
    out = []
    for t in tools:
        try:
            fn = t["function"]
            params = fn.get("parameters") or {}
            if _should_flatten(params):
                new_params, mapping = _flatten_schema_params(params)
                if mapping:
                    flatten_map[fn.get("name", "")] = mapping
                    t = dict(t, function=dict(fn, parameters=new_params))
        except Exception:
            pass
        out.append(t)
    return out, flatten_map


# ============ Tool-call Repair（对标 Reasonix）：截断参数自动修复 + 重复调用风暴抑制 ============
def _repair_tool_args(raw):
    """修复被截断/畸形的工具参数 JSON（对标 Reasonix Truncation Repair）：
    覆盖截断（补引号/括号）+ 畸形（单引号/裸词键/尾逗号/缺逗号/注释）→ 预处理后解析。
    修复成功返回 dict，否则返回 None。"""
    if not isinstance(raw, str) or not raw.strip():
        return None
    s = raw.strip()
    try:
        return json.loads(s)
    except Exception:
        pass
    # ===== 预处理（畸形 JSON 修复）=====
    s2 = s
    # 1) 去掉注释 /* */ 和 // 行注释（值内 // 少见，谨慎：只去 /* */）
    s2 = re.sub(r"/\*[\s\S]*?\*/", "", s2)
    # 2) 单引号 → 双引号（JSON 只认双引号；值内单引号少见可接受）
    if "'" in s2 and '"' in s2:
        s2 = s2.replace("'", '"')
    elif s2.count("'") > s2.count('"'):
        s2 = s2.replace("'", '"')
    # 3) 裸词键 {title: → {"title":
    s2 = re.sub(r"([{,])\s*([A-Za-z_$][\w$-]*)\s*:", r'\1"\2":', s2)
    # 4) 尾逗号 ,} / ,] → } / ]
    s2 = re.sub(r",\s*([}\]])", r"\1", s2)
    # 5) 缺逗号：字符串后紧跟 " { [ （两个值之间）
    s2 = re.sub(r'(")\s*(?="|\{|\[)', r"\1, ", s2)
    # 6) 数字后缺逗号
    s2 = re.sub(r"(\d)\s*(?=\"|\{|\[)", r"\1, ", s2)
    # 7) 布尔/null 后缺逗号
    s2 = re.sub(r"\b(true|false|null)\s*(?=\"|\{|\[)", r"\1, ", s2)
    try:
        return json.loads(s2)
    except Exception:
        pass
    # ===== 原逻辑：补闭合引号/去尾逗号/补括号 =====
    s2 = s
    if s2.count('"') % 2 == 1:
        s2 += '"'
    s2 = re.sub(r",\s*$", "", s2)
    stack = []
    in_str = False
    esc = False
    for ch in s2:
        if esc:
            esc = False
            continue
        if ch == "\\":
            esc = True
            continue
        if ch == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if ch in "[{":
            stack.append(ch)
        elif ch in "]}":
            if stack and ((ch == "]" and stack[-1] == "[") or (ch == "}" and stack[-1] == "{")):
                stack.pop()
    for ch in reversed(stack):
        s2 += "]" if ch == "[" else "}"
    try:
        return json.loads(s2)
    except Exception:
        return None


# ============ 项目规则文件自动加载（对标 opencode 读 AGENTS.md）：项目文件夹下的规则进 system ============
_PROJECT_RULES_FILES = ("AGENTS.md", "CLAUDE.md")


def _load_project_rules(project_path):
    """读取项目规则文件（AGENTS.md / CLAUDE.md，对标 opencode），存在则返回 {file, content} 截断版。
    注意：文件内容字节稳定 → system 前缀稳定 → 不影响前缀缓存；规则文件被编辑时才导致一次缓存 miss（可接受）。"""
    if not project_path or not os.path.isdir(project_path):
        return None
    for fname in _PROJECT_RULES_FILES:
        fp = os.path.join(project_path, fname)
        if os.path.isfile(fp):
            try:
                with open(fp, encoding="utf-8") as f:
                    content = f.read()
                if content.strip():
                    return {"file": fname, "content": content[:4000]}
            except Exception:
                continue
    return None


def _fold_summary(folded):
    """把早期消息折叠成固定摘要（对标 Reasonix Auto-Compact）。
    关键：同一输入必须产出同一输出（字节稳定）——折叠后的摘要成为新的缓存前缀，
    后续请求才能继续命中；因此不做任何时间戳/随机拼接。"""
    parts = ["【对话早期内容摘要（已自动折叠以节省上下文）】"]
    n = 0
    for m in folded:
        c = m.get("content")
        role = "用户" if m.get("role") == "user" else ("AI" if m.get("role") == "assistant" else "工具")
        if isinstance(c, str) and c.strip():
            n += 1
            if n > 12:
                parts.append("（更早的内容不再展开）")
                break
            parts.append(f"- {role}: {c[:150].strip()}")
        elif isinstance(c, list):
            parts.append(f"- {role}: （含图片/附件）")
    parts.append("如需回顾早期细节请直接询问。")
    return "\n".join(parts)


def estimate_tokens(text):
    """估算 token 数；展示与预压缩使用估算，计费始终以供应商 usage 为准。

    cl100k_base 只对采用相同编码的模型精确；其他模型仍属于近似值。无 tiktoken
    时回退到字符/词估算，并用供应商返回的 prompt_tokens 温和校准。
    """
    if not text:
        return 0
    # ---- tiktoken BPE 估算（编码相同的模型可精确，其他模型仅近似）----
    try:
        import tiktoken
        _enc = getattr(estimate_tokens, "_enc", None)
        if _enc is None:
            try:
                _enc = tiktoken.get_encoding("cl100k_base")
            except Exception:
                _enc = tiktoken.encoding_for_model("gpt-4")
            estimate_tokens._enc = _enc
        return len(_enc.encode(text))
    except Exception:
        pass
    # ---- 回退估算（无 tiktoken 时）----
    cjk = len(re.findall(r"[\u4e00-\u9fff]", text))
    words = len(re.findall(r"[a-zA-Z0-9_]+", text))
    other = len(text) - cjk - sum(len(w) for w in re.findall(r"[a-zA-Z0-9_]+", text))
    k_cjk, k_word, k_other = _TOKEN_CAL
    return round(cjk * k_cjk + words * k_word + other * k_other)


def compact_messages_if_needed(msgs, tools, ctx_limit, observed_context=0):
    """Compact early model context once it reaches 85%, preserving recent turns."""
    before = estimate_tokens(json.dumps(msgs, ensure_ascii=False))
    before += estimate_tokens(json.dumps(tools or [], ensure_ascii=False))
    trigger_size = max(before, int(observed_context or 0))
    if trigger_size <= ctx_limit * 0.85 or len(msgs) <= 12:
        return msgs, False, trigger_size, before

    target = max(4096, int(ctx_limit * 0.35))
    keep_n = min(10, len(msgs) - 2)
    best = None
    while keep_n <= len(msgs) - 2:
        fold_end = max(2, len(msgs) - keep_n)
        recent = msgs[fold_end:]
        while recent and recent[0].get("role") == "tool":
            recent = recent[1:]
        candidate = msgs[:1] + [{"role": "user", "content": _fold_summary(msgs[1:fold_end])}] + recent
        size = estimate_tokens(json.dumps(candidate, ensure_ascii=False))
        size += estimate_tokens(json.dumps(tools or [], ensure_ascii=False))
        if size > target and best is not None:
            break
        best = (candidate, size)
        keep_n += 2
    compacted, after = best or (msgs, before)
    return compacted, compacted is not msgs, trigger_size, after


# ---- Token 估算自适应校准（对标原生 tokenizer：用真实 usage 回校系数） ----
# 系数格式 (中文/字, 英文/词, 其他/字符)；初始为近似值，收到真实 prompt_tokens 后修正
_TOKEN_CAL = [1.2, 1.3, 0.4]
_TOKEN_SAMPLES = 0
_MAX_SAMPLES = 30   # 采样 30 次后收敛

def _calibrate_token_estimate(real_tokens, text_len, cjk_len, word_len, other_len):
    """用真实 token 数回校估算系数（EMA 平滑，避免单次波动）。"""
    global _TOKEN_CAL, _TOKEN_SAMPLES
    if not real_tokens or real_tokens <= 0 or text_len <= 0:
        return
    _TOKEN_SAMPLES += 1
    if _TOKEN_SAMPLES > _MAX_SAMPLES:
        return
    # 单次真实 usage 不能可靠拆分到字符类型，只对本次出现的类型按同一比例做 EMA。
    # 限制比例可避免供应商特殊消息模板或异常 usage 把估算系数一次拉飞。
    est = cjk_len * _TOKEN_CAL[0] + word_len * _TOKEN_CAL[1] + other_len * _TOKEN_CAL[2]
    if est <= 0:
        return
    ratio = max(0.5, min(2.0, real_tokens / est))
    alpha = 0.3
    for idx, present in enumerate((cjk_len, word_len, other_len)):
        if present:
            _TOKEN_CAL[idx] *= (1.0 - alpha) + alpha * ratio


@app.get("/api/history")
def history_list(project: str = "", q: str = "", limit: int = 0, offset: int = 0):
    """对话历史列表（按更新时间倒序；置顶优先；可按项目/关键词过滤，支持分页）"""
    convs, total = history_store.list_conversations(
        project or None, q or None, limit or None, offset)
    return {"conversations": convs, "total": total}


class HistoryPinRequest(BaseModel):
    pin: bool = False


class HistoryRenameRequest(BaseModel):
    title: str = ""


@app.post("/api/history/{cid}/pin")
def history_pin(cid: str, req: HistoryPinRequest):
    return {"ok": history_store.set_pin(cid, req.pin)}


@app.post("/api/history/{cid}/rename")
def history_rename(cid: str, req: HistoryRenameRequest):
    if not req.title.strip():
        raise HTTPException(status_code=400, detail="对话名称不能为空")
    if not history_store.rename_conversation(cid, req.title):
        raise HTTPException(status_code=404, detail="对话不存在")
    return {"ok": True}


@app.get("/api/history/{cid}")
def history_get(cid: str):
    conv = history_store.get_conversation(cid)
    if conv is None:
        raise HTTPException(status_code=404, detail="对话不存在")
    # GET 不写盘，但返回前按权威累计 token 重算，旧版本脏 cost 不会继续显示。
    conv["usage"] = history_store.refresh_usage_prices(
        conv.get("usage"), conv.get("provider") or "", conv.get("model") or "")
    return conv


def _conversation_review_workspace(cid):
    conv = history_store.get_conversation(cid)
    if conv is None:
        raise HTTPException(status_code=404, detail="对话不存在")
    project_id = conv.get("project") or history_store.DEFAULT_PROJECT
    project = next((item for item in history_store.list_projects()
                    if item.get("id") == project_id), None)
    project_path = (project or {}).get("path", "")
    if not project_path:
        raise HTTPException(status_code=409, detail="该对话没有关联 Git 项目")
    try:
        workspace = resolve_conversation_workspace(
            cluster.WORKSPACES, cid, current_tenant.get(), project_path)
    except Exception as exc:
        raise HTTPException(status_code=409, detail="无法恢复会话 worktree: %s" % exc)
    if not workspace.isolated:
        raise HTTPException(status_code=409, detail="该项目不是 Git 仓库，没有独立 worktree")
    return workspace


@app.get("/api/history/{cid}/git/status")
def history_git_status(cid: str):
    workspace = _conversation_review_workspace(cid)
    return {"ok": True, "path": str(workspace.path), "branch": workspace.branch,
            "status": cluster.WORKSPACES.status(workspace.task_key)}


@app.get("/api/history/{cid}/git/diff")
def history_git_diff(cid: str):
    workspace = _conversation_review_workspace(cid)
    return {"ok": True, "path": str(workspace.path), "branch": workspace.branch,
            "diff": cluster.WORKSPACES.diff(workspace.task_key)}


@app.post("/api/projects/{pid}/open-folder")
def projects_open_folder(pid: str):
    """在资源管理器中打开项目文件夹（无路径则打开 workspace）"""
    for p in history_store.list_projects():
        if p.get("id") == pid:
            path = p.get("path") or os.path.join(BASE_DIR, "workspace")
            if not os.path.isdir(path):
                raise HTTPException(status_code=400, detail=f"文件夹不存在: {path}")
            try:
                os.startfile(path)   # Windows：用默认资源管理器打开
                return {"ok": True, "path": path}
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"无法打开文件夹: {e}")
    raise HTTPException(status_code=404, detail="项目不存在")


def _conv_to_markdown(conv):
    """把一条对话转成 Markdown 文本（导出用）"""
    lines = [f"# {conv.get('title') or '对话'}", ""]
    for m in conv.get("messages", []):
        role = m.get("role")
        content = m.get("content")
        if role == "user":
            if isinstance(content, list):
                content = " ".join(p.get("text", "") for p in content if p.get("type") == "text")
            if content:
                lines += ["**用户**", str(content), ""]
        elif role == "assistant" and isinstance(content, str) and content:
            lines += ["**AI**", content, ""]
    return "\n".join(lines)


@app.get("/api/projects/{pid}/export-zip")
def projects_export_zip(pid: str):
    """导出项目下所有对话为 ZIP（每个对话一个 Markdown 文件 + 清单）"""
    convs = history_store.list_conversations(project=pid)[0]
    if not convs:
        raise HTTPException(status_code=400, detail="该项目还没有对话")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        manifest = []
        for i, c in enumerate(convs, 1):
            conv = history_store.get_conversation(c["id"])
            md = _conv_to_markdown(conv) if conv else ""
            safe = re.sub(r'[\\/:*?"<>|]', "_", (c["title"] or f"对话{i}"))[:50]
            zf.writestr(f"{i:03d}_{safe}.md", md)
            manifest.append(f"{i:03d}_{safe}.md")
        zf.writestr("00_对话清单.txt", "\n".join(manifest))
    name = re.sub(r'[\\/:*?"<>|]', "_", (convs[0]["title"] or "对话"))[:40] or "对话"
    fname = urllib.parse.quote(name + "_项目对话.zip")
    return Response(
        content=buf.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{fname}"},
    )


@app.get("/api/history/{cid}/export")
def history_export(cid: str, format: str = "md"):
    """导出对话：md（Markdown）/ html（带样式，手机可看）/ json（完整含思考/工具记录）"""
    conv = history_store.get_conversation(cid)
    if conv is None:
        raise HTTPException(status_code=404, detail="对话不存在")
    title = conv.get("title") or "对话"
    if format == "html":
        try:
            import markdown as mdlib
            body_html = []
            for m in conv.get("messages", []):
                role = m.get("role")
                content = m.get("content")
                if role == "user":
                    if isinstance(content, list):
                        content = " ".join(str(p.get("text", "")) for p in content if p.get("type") == "text")
                    body_html.append(f"<div class='msg user'><div class='who'>用户</div><div class='body'>{mdlib.markdown(str(content))}</div></div>")
                elif role == "assistant" and isinstance(content, str) and content:
                    think_html = ""
                    if m.get("think"):
                        think_html = f"<details class='think'><summary>思考过程</summary>{mdlib.markdown(m['think'])}</details>"
                    tools_html = ""
                    if m.get("tools"):
                        tools_html = "<div class='tools'>" + "".join(
                            f"<div class='tool'>▸ {t.get('name','')}{' → 完成' if t.get('done') else ' → 失败'}</div>"
                            for t in m["tools"]) + "</div>"
                    body_html.append(f"<div class='msg ai'><div class='who'>AI</div>{think_html}{tools_html}<div class='body'>{mdlib.markdown(content)}</div></div>")
            html = f"""<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title} · 问墨·code导出</title>
<style>
body {{ max-width: 720px; margin: 0 auto; padding: 24px 16px 60px; font-family: system-ui, "Microsoft YaHei", sans-serif; color: #222; background: #fff; }}
h1 {{ font-size: 22px; border-bottom: 2px solid #6E8BFF; padding-bottom: 8px; }}
.msg {{ margin: 14px 0; padding: 12px 14px; border-radius: 10px; }}
.msg.user {{ background: #F0F4FF; }}
.msg.ai {{ background: #F7F7F9; }}
.who {{ font-size: 12px; color: #888; margin-bottom: 6px; font-weight: 600; }}
.tools {{ font-size: 12px; color: #B8860B; margin: 4px 0; }}
.think summary {{ color: #7C5CD6; font-size: 13px; cursor: pointer; }}
.think {{ margin: 4px 0; }}
code {{ background: #EEE; padding: 1px 5px; border-radius: 4px; }}
pre {{ background: #0D1117; color: #E7EAF0; padding: 12px; border-radius: 8px; overflow-x: auto; }}
pre code {{ background: none; }}
table {{ border-collapse: collapse; }}
td, th {{ border: 1px solid #CCC; padding: 4px 10px; }}
</style></head><body><h1>{title}</h1>{''.join(body_html)}</body></html>"""
            fname = re.sub(r'[\\/:*?"<>|]', "_", title)[:40] or "对话"
            return Response(content=html.encode("utf-8"), media_type="text/html; charset=utf-8",
                            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{urllib.parse.quote(fname + '_导出.html')}"})
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"HTML 导出失败: {e}")
    if format == "json":
        fname = re.sub(r'[\\/:*?"<>|]', "_", title)[:40] or "对话"
        return Response(content=json.dumps(conv, ensure_ascii=False, indent=2).encode("utf-8"),
                        media_type="application/json; charset=utf-8",
                        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{urllib.parse.quote(fname + '_导出.json')}"})
    return {"ok": True, "title": title, "markdown": _conv_to_markdown(conv)}


class FileOpenRequest(BaseModel):
    name: str


@app.get("/api/files/download/{name}")
def files_download(name: str):
    """下载文件（attachment 强制下载，WebView2/浏览器均可靠触发）。
    与 /files 静态直链不同：带 Content-Disposition attachment，客户端点击下载链接不会变成导航。"""
    import urllib.parse as _up
    if not re.match(r"^[\w\u4e00-\u9fff\s（）()\-_\.]+$", name):
        raise HTTPException(status_code=400, detail="非法文件名")
    path = os.path.join(_files_dir(), name)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="文件不存在")
    return FileResponse(
        path,
        media_type="application/octet-stream",
        headers={"Content-Disposition": "attachment; filename*=UTF-8''%s" % _up.quote(name)},
    )


@app.post("/api/files/open")
def files_open(req: FileOpenRequest):
    """用电脑上的默认软件打开下载区文件（PPT/PDF/Word/Excel/VS Code 等）。
    浏览器不能直接调起本地程序，所以由服务器代劳（os.startfile）。"""
    name = req.name.strip()
    if not name or "/" in name or "\\" in name or name in (".", ".."):
        raise HTTPException(status_code=400, detail="非法文件名")
    base = os.path.normpath(_files_dir())
    path = os.path.normpath(os.path.join(base, name))
    if not path.startswith(base) or not os.path.isfile(path):
        raise HTTPException(status_code=404, detail=f"文件不存在: {name}")
    try:
        os.startfile(path)  # Windows：按文件关联打开（Word/Excel/PPT/VS Code…）
        return {"ok": True, "path": path}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"无法打开文件: {e}")


class HistorySaveRequest(BaseModel):
    id: str = ""
    title: str = ""
    messages: list[dict] = []
    project: str = "default"
    provider: str = ""
    model: str = ""
    usage: dict = {}   # {input, output, cached} 对话累计用量（历史完整保存）


class HistoryAppendRequest(BaseModel):
    id: str = ""
    message: dict
    project: str = "default"
    provider: str = ""
    model: str = ""


@app.post("/api/history/message")
def history_append(req: HistoryAppendRequest):
    """Persist one user message before starting its model request."""
    try:
        cid, message = history_store.append_message(
            req.id or None,
            req.message,
            project=req.project,
            provider=req.provider,
            model=req.model,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    return {"ok": True, "id": cid, "message": message}


@app.post("/api/history")
def history_save(req: HistorySaveRequest):
    # 用量由 /api/chat 根据供应商返回的真实 usage 在服务端累计。
    # 浏览器只保存消息，避免旧页面/后台流用过时的累计值覆盖服务端账本。
    cid = history_store.save_conversation(req.messages, cid=req.id or None,
                                           title=req.title, project=req.project,
                                           provider=req.provider, model=req.model,
                                           usage=None)
    return {"ok": True, "id": cid}


@app.delete("/api/history/{cid}")
def history_delete(cid: str):
    if not history_store.delete_conversation(cid):
        raise HTTPException(status_code=404, detail="对话不存在或已删除")
    return {"ok": True}


# ==================== 用户登录（GitHub OAuth）====================
import auth as auth_module
from fastapi import Request as _AuthRequest

def _get_auth_token(request):
    """HttpOnly cookie is authoritative; header remains for legacy clients."""
    return (request.cookies.get("wenmo_session") or
            request.headers.get("X-Wenmo-Token", "")).strip()


@app.get("/api/auth/status")
async def auth_status(request: _AuthRequest):
    """登录状态：返回当前用户（无 token 时未登录）"""
    token = _get_auth_token(request)
    login = auth_module.validate_session(token) if token else None
    if login:
        u = auth_module.get_user(login)
        return {"logged_in": True, "user": u}
    return {"logged_in": False, "user": None,
            "oauth_configured": auth_module.is_configured()}


@app.get("/api/auth/login")
async def auth_login():
    """登录/配置入口：GitHub OAuth 已配置才跳转授权页；未配置返回 configured=False（前端引导配置模型 Key）"""
    if not auth_module.is_configured():
        return {"configured": False, "url": None,
                "hint": "GitHub 登录未配置（可选项）。直接用顶部模型下拉 + 设置页填 API Key 即可开始使用。"}
    return {"configured": True, "url": auth_module.github_oauth_url()}


@app.get("/api/auth/github/callback")
async def auth_github_callback(code: str = "", state: str = ""):
    """GitHub OAuth 回调：换 token → 取用户 → 建会话 → 重定向回前端"""
    if not code:
        return {"error": "缺少授权码"}
    if not auth_module.consume_oauth_state(state):
        raise HTTPException(status_code=400, detail="OAuth state 无效或已过期")
    tok = auth_module.github_exchange_code(code)
    if tok.get("error"):
        return {"error": f"OAuth 失败: {tok['error']}"}
    user = auth_module.github_get_user(tok.get("access_token", ""))
    if user.get("error"):
        return {"error": f"获取用户失败: {user['error']}"}
    login = user.get("login", "")
    if not login:
        return {"error": "GitHub 返回无 login"}
    auth_module.upsert_user(login, user)
    session_token = auth_module.create_session(login)
    response = RedirectResponse(url="/", status_code=303)
    response.set_cookie(
        "wenmo_session", session_token, httponly=True, samesite="strict",
        secure=auth_module.BASE_URL.lower().startswith("https://"),
        max_age=auth_module.SESSION_TTL, path="/",
    )
    return response


@app.post("/api/auth/logout")
async def auth_logout(request: _AuthRequest):
    """登出：清除会话 token"""
    token = _get_auth_token(request)
    if token:
        auth_module.logout_session(token)
    response = JSONResponse({"ok": True})
    response.delete_cookie("wenmo_session", path="/", samesite="strict")
    return response


@app.get("/api/auth/users")
async def auth_users(request: _AuthRequest):
    """Return only the signed-in user; there is no unauthenticated admin surface."""
    login = auth_module.validate_session(_get_auth_token(request))
    if not login:
        raise HTTPException(status_code=401, detail="未登录")
    return {"users": [auth_module.get_user(login)]}


# ==================== 自动更新（GitHub Releases）====================
import updater as updater_module
@app.get("/api/version")
async def api_version():
    """返回当前版本号（前端设置-通用页显示用）"""
    return {"version": updater_module.APP_VERSION, "name": "问墨·code"}


@app.get("/api/health")
async def api_health():
    """Small local readiness endpoint for startup checks and release smoke tests."""
    history_ready = os.path.isdir(history_store.BASE)
    config_ready = os.path.isdir(os.path.dirname(os.path.abspath(_settings_file())))
    ok = history_ready and config_ready
    return {
        "status": "ok" if ok else "degraded",
        "version": updater_module.APP_VERSION,
        "storage": {"history": history_ready, "config": config_ready},
    }


@app.get("/api/update/check")
async def update_check():
    """检查 GitHub Releases 是否有新版本。返回 {has_update, version, notes, url}
    放线程池执行（check_update 内部是同步 urllib，避免阻塞事件循环/其他请求）"""
    try:
        import asyncio as _asyncio
        import concurrent.futures as _cf
        _loop = _asyncio.get_running_loop()
        info = await _loop.run_in_executor(None, updater_module.check_update, 5)   # 5秒超时，快速失败
        if info:
            # 更新模式：相邻版本(1.0.4→1.0.5)且有 manifest → 增量 delta；否则全量
            _mode = "full"
            _url = info.get("url", "")
            _asset = info.get("asset_name", "")
            try:
                _cur = [int(x) for x in updater_module.APP_VERSION.split(".")]
                _new = [int(x) for x in info.get("version", "").split(".")]
                if len(_cur) >= 3 and len(_new) >= 3 and _new[0] == _cur[0] and _new[1] == _cur[1] and _new[2] == _cur[2] + 1:
                    _mode = "delta"
                    # delta 模式：优先用增量包地址（小几十 MB），前端下载它而非全量 190MB
                    if info.get("delta_url"):
                        _url = info["delta_url"]
                        _asset = info.get("delta_name", "")
            except Exception:
                pass
            return {
                "has_update": True,
                "version": info.get("version"),
                "notes": (info.get("notes") or "")[:500],
                "url": _url,
                "asset_name": _asset,
                "current": updater_module.APP_VERSION,
                "update_mode": _mode,
            }
        return {"has_update": False, "current": updater_module.APP_VERSION}
    except Exception as e:
        return {"has_update": False, "error": str(e), "current": updater_module.APP_VERSION}


# 更新下载进度（全局状态，供 /api/update/progress 轮询）
_UPDATE_DL = {"running": False, "downloaded": 0, "total": 0, "percent": 0, "done": False, "error": "", "exe_path": ""}

@app.get("/api/update/progress")
def update_progress():
    """返回下载进度 {running, downloaded, total, percent, done, error}"""
    return dict(_UPDATE_DL)

@app.post("/api/update/download")
def update_download():
    """下载新版本安装包到临时目录。返回 {ok, message, exe_path}"""
    try:
        if _UPDATE_DL.get("running"):
            return {"ok": True, "message": "已在下载中", "downloading": True}
        info = updater_module.check_update()
        if not info or not info.get("url"):
            return {"ok": False, "message": "没有可下载的更新"}
        # delta 模式（相邻版本且远程有 delta 包）→ 直接走增量下载（小几十 MB），避免下全量 190MB
        if info.get("delta") or info.get("delta_url"):
            _UPDATE_DL.update(running=True, downloaded=0, total=0, percent=0, done=False, error="", exe_path="")
            def _dworker():
                try:
                    def _cb(dl, total):
                        # 真实字节进度（total 未知 → -1，前端显示 MB，根治 0% 卡死）
                        _UPDATE_DL["downloaded"] = dl
                        _UPDATE_DL["total"] = total
                        _UPDATE_DL["percent"] = int(dl * 100 / total) if total else -1
                    _ok, _msg = updater_module.try_delta_update(progress_cb=_cb)
                    if _ok:
                        _UPDATE_DL["done"] = True
                        _UPDATE_DL["message"] = _msg
                    else:
                        # 增量失败 → 自动降级全量下载（方案：不报错卡死）
                        _UPDATE_DL["message"] = "增量不可用（" + _msg + "），自动切换全量下载…"
                        try:
                            _i2 = updater_module.check_update()
                            if not _i2 or not _i2.get("url"):
                                _UPDATE_DL["error"] = "全量下载不可用：无更新地址"
                                return
                            _ep, _er = updater_module.download_update(_i2.get("url", ""), _i2.get("sha256", ""), progress_cb=_cb)
                            if _er:
                                _UPDATE_DL["error"] = _er
                            else:
                                _UPDATE_DL["exe_path"] = _ep
                                updater_module.set_downloaded_path(_ep)
                                _UPDATE_DL["message"] = "全量包下载完成，可点击应用更新"
                                _UPDATE_DL["done"] = True
                        except Exception as _e2:
                            _UPDATE_DL["error"] = "全量下载失败: " + str(_e2)
                except Exception as e:
                    _UPDATE_DL["error"] = str(e)
                    _UPDATE_DL["done"] = True
                finally:
                    _UPDATE_DL["running"] = False
            threading.Thread(target=_dworker, daemon=True).start()
            return {"ok": True, "message": "开始增量下载…", "downloading": True, "delta": True}
        url = info.get("url", ""); sha = info.get("sha256", "")

        def _cb(downloaded, total):
            # 真实字节进度：total 未知（国内 CDN 无 Content-Length）→ percent=-1，
            # 前端显示"已下载 X MB"而不是卡 0%（根治 0% 卡死）
            _UPDATE_DL["downloaded"] = downloaded
            _UPDATE_DL["total"] = total
            _UPDATE_DL["percent"] = int(downloaded * 100 / total) if total else -1

        _UPDATE_DL.update(running=True, downloaded=0, total=0, percent=0, done=False, error="", exe_path="")

        def _worker():
            try:
                exe_path, err = updater_module.download_update(url, sha, progress_cb=_cb)
                if err:
                    _UPDATE_DL["error"] = err
                else:
                    _UPDATE_DL["exe_path"] = exe_path
                    updater_module.set_downloaded_path(exe_path)
            except Exception as e:
                _UPDATE_DL["error"] = str(e)
            finally:
                _UPDATE_DL["running"] = False
                _UPDATE_DL["done"] = True

        threading.Thread(target=_worker, daemon=True).start()
        return {"ok": True, "message": "开始下载", "downloading": True}
    except Exception as e:
        return {"ok": False, "message": "下载异常: " + str(e)}


@app.post("/api/update/apply")
def update_apply():
    """应用更新：zip → 文件级替换；exe → 启动安装包。返回 {ok, message}"""
    try:
        # 全量包已经下载完成时，必须应用该包；不能再次启动增量下载。
        exe_path = updater_module.get_downloaded_path()
        if exe_path and os.path.exists(exe_path):
            if exe_path.lower().endswith(".zip"):
                err = updater_module.apply_zip_update(exe_path)
                if err:
                    return {"ok": False, "message": "应用更新失败: " + err}
                return {"ok": True, "message": "正在替换程序文件，问墨即将自动重启..."}
            err = updater_module.apply_update(exe_path)
            if err:
                return {"ok": False, "message": "启动安装失败: " + err}
            return {"ok": True, "message": "正在安装更新，问墨即将关闭..."}

        # 优先增量更新（delta）：相邻版本 + manifest + delta 包；失败自动回退全量
        # 线程化执行：delta 下载过程更新 _UPDATE_DL 进度（前端轮询显示真实进度条）
        def _delta_worker():
            try:
                _UPDATE_DL.update(running=True, downloaded=0, total=0, percent=0, done=False, error="", exe_path="")
                def _cb(dl, total):
                    # 真实字节进度（total 未知 → -1 显示 MB，根治 0% 卡死）
                    _UPDATE_DL["downloaded"] = dl
                    _UPDATE_DL["total"] = total
                    _UPDATE_DL["percent"] = int(dl * 100 / total) if total else -1
                _ok, _msg = updater_module.try_delta_update(progress_cb=_cb)
                if _ok:
                    _UPDATE_DL["done"] = True
                    _UPDATE_DL["message"] = _msg
                    return
                # 增量失败 → 自动静默降级全量下载（方案：国内 CDN 增量不可用时不报错卡死）
                _UPDATE_DL["message"] = "增量不可用（" + _msg + "），自动切换全量下载…"
                try:
                    _info2 = updater_module.check_update()
                    if not _info2 or not _info2.get("url"):
                        _UPDATE_DL["error"] = "全量下载不可用：无更新地址"
                        return
                    exe_path, err = updater_module.download_update(_info2.get("url", ""), _info2.get("sha256", ""), progress_cb=_cb)
                    if err:
                        _UPDATE_DL["error"] = err
                    else:
                        _UPDATE_DL["exe_path"] = exe_path
                        updater_module.set_downloaded_path(exe_path)
                        _UPDATE_DL["message"] = "全量包下载完成，可点击应用更新"
                        _UPDATE_DL["done"] = True          # 完成标记：前端停止轮询并提示
                except Exception as e2:
                    _UPDATE_DL["error"] = "全量下载失败: " + str(e2)
            except Exception as e:
                _UPDATE_DL["error"] = str(e)
            finally:
                _UPDATE_DL["running"] = False
        threading.Thread(target=_delta_worker, daemon=True).start()
        return {"ok": True, "message": "正在增量更新…", "downloading": True, "delta": True}
    except Exception as e:
        return {"ok": False, "message": "应用更新异常: " + str(e)}


@app.get("/api/projects")
def projects_list():
    """项目列表（每个项目有独立的对话历史）"""
    return {"projects": history_store.list_projects()}


class ProjectRequest(BaseModel):
    name: str
    path: str = ""


@app.post("/api/projects")
def projects_add(req: ProjectRequest):
    pid = history_store.add_project(req.name, req.path)
    return {"ok": True, "id": pid}


@app.post("/api/projects/{pid}/rename")
def projects_rename(pid: str, req: ProjectRequest):
    ok = history_store.rename_project(pid, req.name)
    if not ok:
        raise HTTPException(status_code=404, detail="项目不存在")
    return {"ok": True}


@app.post("/api/projects/{pid}/update")
def projects_update(pid: str, req: dict = None):
    """更新项目自定义属性（图标颜色/显示文字/启动脚本/文件夹路径）"""
    body = req or {}
    ok = history_store.update_project(
        pid,
        icon_color=body.get("icon_color"),
        icon_text=body.get("icon_text"),
        launch_cmd=body.get("launch_cmd"),
        path=body.get("path"),
    )
    if not ok:
        raise HTTPException(status_code=404, detail="项目不存在")
    return {"ok": True}


@app.post("/api/launch")
def launch_cmd(req: dict = None):
    """Run the stored launch command for one tenant-owned project."""
    body = req or {}
    project_id = str(body.get("project_id") or "").strip()
    project = next((item for item in history_store.list_projects()
                    if item.get("id") == project_id), None)
    if project is None:
        raise HTTPException(status_code=404, detail="项目不存在")
    cmd = str(project.get("launch_cmd") or "").strip()
    cwd = os.path.realpath(str(project.get("path") or "").strip())
    if not cmd:
        raise HTTPException(status_code=400, detail="该项目未配置启动命令")
    if not cwd or not os.path.isdir(cwd):
        raise HTTPException(status_code=400, detail="项目工作目录不存在")
    decision = evaluate_permission(
        "project_launch", {"command": cmd, "path": cwd},
        load_runtime_policy(_settings_file()))
    if decision.effect == "deny":
        raise HTTPException(status_code=403, detail={
            "message": "权限策略拒绝项目启动命令",
            "permission": decision.__dict__,
        })
    if decision.effect == "ask" and not body.get("_confirmed"):
        raise HTTPException(status_code=409, detail={
            "message": "项目启动命令需要确认",
            "need_confirmation": True,
            "permission": decision.__dict__,
        })
    try:
        # shell=False prevents shell metacharacters from turning one configured
        # executable into an unrelated command chain. Windows CreateProcess
        # accepts a command line string; POSIX receives an explicitly split argv.
        args = cmd if os.name == "nt" else shlex.split(cmd)
        subprocess.Popen(args, shell=False, cwd=cwd,
                         creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
        return {"ok": True, "project_id": project_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"启动失败: {e}")

@app.post("/api/emergency-stop")
def emergency_stop():
    """紧急停止：终止所有正在生成的流 + 清空待处理询问（对标硬件急停，防 AI 持续下发指令）"""
    n = 0
    for qid in list(PENDING_ASKS.keys()):
        try:
            fut = PENDING_ASKS[qid]
            if not fut.done():
                fut.set_result("（紧急停止：任务已终止）")
                n += 1
        except Exception:
            pass
    return {"ok": True, "stopped_asks": n}


@app.delete("/api/projects/{pid}")
def projects_delete(pid: str):
    return {"ok": history_store.delete_project(pid)}


@app.post("/api/pick-folder")
async def pick_folder():
    """在服务器所在机器上弹出原生文件夹选择器（浏览器沙箱拿不到路径，
    本机场景由服务器代劳：用 tkinter 弹 Windows 资源管理器文件夹选择）。"""
    code = (
        "import tkinter as tk;"
        "from tkinter import filedialog;"
        "root = tk.Tk(); root.withdraw(); root.attributes('-topmost', True);"
        "r = filedialog.askdirectory(title='选择项目文件夹');"
        "print(r or '')"
    )
    try:
        proc = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True, text=True, timeout=180,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        path = (proc.stdout or "").strip()
        return {"ok": True, "path": path}
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="文件夹选择超时")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"无法打开文件夹选择器: {e}")


@app.get("/api/plugins")
def list_plugins():
    """列出已加载的插件及其工具"""
    plugins = plugins_loader.load_plugins()
    return {"plugins": plugins}


@app.get("/api/mcp/servers")
async def mcp_servers():
    """列出 MCP 服务器与它们的工具（供前端展示）"""
    await mcp_refresh()
    return {"servers": await mcp_snapshot()}


@app.post("/api/mcp/refresh")
async def mcp_refresh_endpoint():
    """重新连接所有 MCP 服务器"""
    await mcp_refresh()
    return {"servers": await mcp_snapshot()}


class MCPCallRequest(BaseModel):
    server: str
    tool: str
    arguments: dict = {}


@app.post("/api/mcp/call")
async def mcp_call_endpoint(req: MCPCallRequest):
    """手动调用一个 MCP 工具（测试/调试用）"""
    if os.environ.get("WENMO_ENABLE_DEBUG_API") != "1":
        raise HTTPException(status_code=404, detail="调试接口未启用")
    try:
        await mcp_refresh()  # 确保已连接
        text = await mcp_call(f"{req.server}_{req.tool}", req.arguments)
        return {"ok": True, "result": text[:2000]}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


class KeySettings(BaseModel):
    provider: str
    api_key: str = ""  # 传空字符串 = 清除该供应商的 key
    model: str = ""    # 传非空 = 同时切换该供应商的模型


class PricingSyncRequest(BaseModel):
    provider: str


@app.get("/api/pricing/status")
def pricing_status():
    path = tenant_file("pricing_catalogs.json")
    try:
        with open(path, encoding="utf-8") as handle:
            data = json.load(handle)
    except Exception:
        data = {"version": 1, "providers": {}}
    return data


@app.post("/api/pricing/sync")
async def pricing_sync(req: PricingSyncRequest):
    """Sync provider model metadata; use embedded prices only when the API supplies them."""
    providers = load_providers()
    cfg = providers.get(req.provider)
    if not cfg:
        raise HTTPException(status_code=404, detail="未知供应商")
    catalog_url = str(cfg.get("pricing_catalog_url") or "").strip()
    if not catalog_url and req.provider == "zen":
        catalog_url = "https://opencode.ai/zen/v1/models"
    if not catalog_url:
        raise HTTPException(status_code=409, detail="该供应商未配置官方模型/价格目录 API")
    try:
        from pricing_sync import fetch_catalog, save_provider_catalog
        catalog = await asyncio.to_thread(
            fetch_catalog,
            catalog_url,
            resolve_key(cfg) or "",
            20,
            str(cfg.get("pricing_catalog_currency") or ""),
            str(cfg.get("pricing_catalog_unit") or ""),
        )
        path = tenant_file("pricing_catalogs.json")
        await asyncio.to_thread(save_provider_catalog, path, req.provider, catalog)
    except Exception as exc:
        raise HTTPException(status_code=502, detail="目录同步失败: %s" % exc)
    applied = len(catalog.get("models") or {})
    if applied:
        note = "已应用 %d 个单位明确的目录价；金额仍是本地估算，不是供应商账单" % applied
    elif catalog.get("price_fields_ignored"):
        note = "目录含价格字段，但币种或单位未明确，因此拒绝计费，只同步可用模型"
    else:
        note = "供应商未返回可用价格，只同步模型目录；费用继续使用本地估算"
    return {"ok": True, "provider": req.provider, **catalog,
            "applied_price_models": applied, "billing_note": note,
            "invoice_reconciled": False}


@app.post("/api/settings")
def save_key(req: KeySettings):
    """保存/清除某个供应商的 api_key（可选同时切换 model）到 providers.json。
    教学点：key 写进服务器本地文件，浏览器永远看不到它。"""
    providers = load_providers()
    if req.provider not in providers:
        raise HTTPException(status_code=400, detail=f"未知供应商: {req.provider}")
    providers[req.provider]["api_key"] = req.api_key.strip()
    if req.model.strip():
        providers[req.provider]["model"] = req.model.strip()
    save_providers(providers)
    return {"ok": True, "provider": req.provider, "has_key": bool(req.api_key.strip())}


@app.post("/api/chat")
def chat(req: ChatRequest):
    if req.provider == "local":
        # 本地 GGUF：直接路由到当前运行的推理服务
        if LOCAL_STATE["status"] != "ready" or not LOCAL_STATE["port"]:
            raise HTTPException(status_code=400, detail=(
                f"本地模型未就绪（状态: {LOCAL_STATE['status']}）。"
                "请打开 设置 → 供应商 → 本地 GGUF，扫描并加载一个模型后再试。"))
        cfg = {"base_url": f"http://127.0.0.1:{LOCAL_STATE['port']}/v1", "model": LOCAL_STATE["name"]}
        api_key = "local"
    else:
        providers = load_providers()
        if req.provider not in providers:
            raise HTTPException(status_code=400, detail=f"未知供应商: {req.provider}")
        cfg = providers[req.provider]
        # ---- 免费综合供应商：按所选模型路由到对应国产平台（各平台免费档）----
        # free 不是单一平台，而是聚合入口：模型名 → 目标供应商 key → 用该家的 base_url + key
        if req.provider == "free" and req.model.strip():
            _route = FREE_MODEL_ROUTES.get(req.model.strip())
            if _route and _route in providers:
                cfg = providers[_route]        # 用目标供应商的 base_url + key（用户需在设置里配好）
        api_key = resolve_key(cfg)
        if not api_key and not cfg.get("base_url", "").startswith(("http://localhost", "http://127.0.0.1")):
            raise HTTPException(status_code=400, detail="该远程供应商尚未配置 API Key")
    # 本地模型（含 Ollama/llama.cpp）不需要 key，但 SDK 要求非空，填个占位符
    client = OpenAI(base_url=cfg["base_url"], api_key=api_key or "local")   # 同步（辅助函数用）
    async_client = AsyncOpenAI(base_url=cfg["base_url"], api_key=api_key or "local")  # 异步（主对话流式用，不阻塞事件循环）
    # 前端指定的模型覆盖（"模型设置"里的默认文本模型）
    if req.model.strip():
        # 免费综合：路由到目标供应商后，模型名也用目标家的真实模型名（free 显示名 → 真实名）
        _real_model = req.model.strip()
        if req.provider == "free":
            _real_model = FREE_MODEL_ROUTES.get(_real_model, _real_model)
            # 若路由后是 dict 形式（{provider, model}），取 model
            if isinstance(_real_model, dict):
                _real_model = _real_model.get("model", req.model.strip())
        cfg = {**cfg, "model": _real_model}
    # 记录主对话模型：智能体未配置时默认跟随它
    _set_agent_model(req.provider, cfg["model"])
    request_id = (req.request_id or uuid.uuid4().hex).strip()
    if not re.fullmatch(r"[A-Za-z0-9._-]{8,128}", request_id):
        raise HTTPException(status_code=400, detail="无效的流请求 id")

    request_tenant = current_tenant.get()

    async def gen():
        """SSE 生成器（函数调用版）：把模型的流式输出转成 data: {...} 事件。

        事件类型：
          {"delta": "..."}   正文流
          {"think": "..."}   思考过程流（思考型模型）
          {"tool": {...}}    工具调用开始/结束（MCP 工具）
          {"done": true}     结束
          {"error": "..."}   错误

        核心循环：模型可请求调用 MCP 工具 → 我们执行并把结果喂回 →
        模型继续回答，最多 8 轮工具调用。这就是"agent 循环"：
        思考 → 行动 → 观察 → 再思考。
        """
        MAX_LOOPS = 50
        MAX_CONTINUE = 8
        CONTINUE_MARKER = "（你的回答被截断了，请紧接上文继续，不要重复已说过的内容）"
        # StreamingResponse may outlive the HTTP middleware call. Capture and
        # restore the authenticated tenant inside the stream task itself.
        tenant_token = current_tenant.set(request_tenant)
        workspace_token = None
        stream_task = asyncio.current_task()
        stream_key = (request_tenant, request_id)
        with _ACTIVE_CHAT_STREAMS_LOCK:
            cutoff = time.monotonic() - 60
            for stale_id, cancelled_at in list(_CANCELLED_CHAT_STREAMS.items()):
                if cancelled_at < cutoff:
                    _CANCELLED_CHAT_STREAMS.pop(stale_id, None)
            if _CANCELLED_CHAT_STREAMS.pop(stream_key, None) is not None:
                return
            previous_task = _ACTIVE_CHAT_STREAMS.get(stream_key)
            if previous_task is not None and not previous_task.done() and previous_task is not stream_task:
                previous_task.cancel()
            _ACTIVE_CHAT_STREAMS[stream_key] = stream_task
        # 记录当前对话：智能体委托时用它关联会话上下文（不同对话互不影响）
        _set_agent_session(req.conversation_id or req.project or "default")
        def _ends_cleanly(text):
            """判断回复是否以正常标点收尾（否则视为被截断，需要续写）"""
            t = (text or "").rstrip()
            if not t:
                return False
            if t.endswith("```"):
                return True
            # 剥离末尾 emoji/符号（如 ✅）再判断标点，避免误判
            t = re.sub(r"[\U0001F000-\U0001FAFF\u2600-\u27BF\uFE0F\u2B00-\u2BFF]+$", "", t)
            if not t:
                return True
            return t[-1] in "。！？；：.!?;:…\"'”’）】》〉,，、"

        def _merge_continuations(msgs):
            """把"截断续写"产生的 assistant(部分)+user(继续)+assistant(续) 合并成一条完整回复"""
            out = []
            i = 0
            while i < len(msgs):
                m = msgs[i]
                if (m.get("role") == "user" and m.get("content") == CONTINUE_MARKER
                        and out and out[-1].get("role") == "assistant"
                        and i + 1 < len(msgs) and msgs[i + 1].get("role") == "assistant"):
                    out[-1]["content"] = (out[-1].get("content") or "") + (msgs[i + 1].get("content") or "")
                    i += 2
                    continue
                out.append(m)
                i += 1
            return out

        try:
            project_record = next((p for p in history_store.list_projects()
                                   if p.get("id") == req.project), None)
            project_path = (project_record or {}).get("path", "")
            if project_path:
                conversation_workspace = resolve_conversation_workspace(
                    cluster.WORKSPACES, req.conversation_id, request_tenant, project_path)
                workspace_token = current_workspace.set(str(conversation_workspace.path))
            msgs = list(req.messages)
            # ---- 用户纠正检测：发现纠正词 → 记教训（模型以后避免重犯）----
            for _m in msgs:
                _c = _m.get("content")
                if _m.get("role") == "user" and isinstance(_c, str):
                    if any(h in _c for h in CORRECTION_HINTS):
                        _save_lesson(f"用户纠正/指出问题：{_c[:120]}", source="user_correction")
                        break
            # 上下文只记录"用户输入 + 模型正式输出"：剥离思考过程/工具记录等展示性字段
            for _m in msgs:
                _m.pop("think", None)
                _m.pop("tools", None)
                _m.pop("tool_calls_log", None)
            # 防御：历史/旧格式消息里的 image_url 图片 → 剥离为占位提示（模型走 see_image 读图）
            for _i, _m in enumerate(msgs):
                _c = _m.get("content")
                if isinstance(_c, list):
                    _parts = []
                    for _p in _c:
                        if isinstance(_p, dict) and _p.get("type") == "image_url":
                            _parts.append({"type": "text", "text": "[图片附件（旧格式），请用 see_image 工具查看]"})
                        else:
                            _parts.append(_p)
                    _m["content"] = _parts if _parts else _c
            # ---- 自适应预算：按模型上下文大小调节工具/技能用量 ----
            ctx_limit = load_settings()["local_ctx"] if req.provider == "local" else (
                8192 if req.provider == "ollama" else load_settings()["remote_ctx"])
            max_tools = min(32, max(12, ctx_limit // 1024))
            skill_count = 3
            # ---- 系统提示：基础规则 + 项目上下文 + 通用技能 + 联网 + 个性技能 ----
            # 注意：最终必须合并成【一条】system 消息——llama.cpp / 部分网关模板
            # 要求 system 位于消息最开头，两条连续 system 会报
            # "System message must be at the beginning" 500 错误。
            # 缓存优化：system 只放【静态内容】（基础规则+项目+通用技能），
            # 保证前缀稳定 → prompt cache 命中率最大化；
            # 动态内容（个性技能、联网提示）追加到【最后一条 user 消息】，不破坏前缀。
            sys_parts = ["你是本软件的内置 AI 助手，没有其他身份设定。不要自称或扮演其他产品"
                         "（如 Claude、ChatGPT、DeepSeek、OpenAI 等）的身份；如需介绍自己，"
                         "说明你是本软件（问墨·code）的助手即可。保持真实、直接、实事求是。",
                         "【了解自己的构成】你运行在『问墨·code』里——一个本地部署的多供应商 AI 聊天软件"
                         "（Python FastAPI 后端 + Web 前端）。你了解自己的内部构成："
                         "① 模型：可切换云端模型（zen/DeepSeek/通义/智谱等）或本地 GGUF 模型（llama.cpp 加载）；"
                         "② 工具：插件工具（写文件并给链接、终端指令 run_command、画图 draw_chart、视觉 see_image、"
                         "智能体委托 delegate_to_agent、系统信息等）与 MCP 工具（网页搜索 web_search、文件系统、GitHub）；"
                         "③ 技能库：199 个技能按 14 类组织（frontend/backend/data/lang/testing/security/debug/git/devops/"
                         "doc/research/agent/plan/perspective）。技能【不注入 prompt】——当任务涉及某领域需要规范指导时，"
                         "先调 skill_list 查看该类技能，再调 skill_load 按名加载正文遵循执行；skill_overview 可看全貌。"
                         "通用技能（编码规范/错误处理/API 设计等）时刻生效。"
                         "④ 功能：对话历史按项目分组、图片识别、公式转 Word、文件下载/预览、发送队列、深度思考档位。"
                         "用户询问软件功能、模型、架构、操作方式时，如实清晰地介绍。",
                         "当你调用工具时，最终回答里要用一两句话向用户说明你做了什么、得到了什么结果。",
                         "当你需要澄清用户意图、或在多个方案中选择时，使用 ask_user 工具弹窗询问用户，不要擅自假设。",
                         "【文件创建规则】当用户要求写文件/创建文件/保存到文件/生成文档、报告、代码文件时，"
                         "你必须调用 plugin_deliver_write_file_with_link 工具把完整内容写入文件并给出下载/预览链接，"
                         "再在对话里简要说明——绝不能只在对话里输出文字而不给文件。"
                         "【长内容写入】写文件工具支持长内容，可完整写入整份报告/长文，无需拆分截断；"
                         "若内容确实很长，为保证完整与质量，宁可完整写入一个文件也不为省事拆成碎片。"
                         "（注：工具参数若有解析问题可重试或分批，但不要因担心截断而主动裁剪内容。）",
                         "【文档排版规则】用户要求 PPT/Word/Excel/表格等办公文档时，必须调用专门工具"
                         "（plugin_ppt_create_pptx / plugin_word_create_docx / plugin_excel_create_xlsx），注重排版与设计感：结构清晰"
                         "（封面/标题/分点/分节）、要点化而不是大段文字、配色协调、善用工具的主题色/层级/表格参数；"
                         "绝不要用纯文本文件替代（没有格式也没有审美）。",
                         "【禁止脚本生成文档】docx/pptx/xlsx/pdf 等办公文档【必须】用专门工具"
                         "（create_docx / create_pptx / create_xlsx / create_pdf），"
                         "【严禁】自己写脚本（python/批处理）生成办公文档——脚本生成的文件路径不可控，"
                         "预览/下载链接会 404，且脚本失败时对话却可能谎报成功。"
                         "公式用 create_docx 的 formula 块（LaTeX/Unicode 数学符号），不要为此写脚本。"
                         "文件读取/写入/打包/压缩用 file_mcp 工具（file_read / file_write / file_zip）"
                         "或 workspace 工具（写入自动落到下载区保证链接可用）。",
"【文档美学】生成文档/PPT/表格/Markdown 时注重设计感（对标 ChatGPT 的质量）：\n① 排版：层级清晰（封面→目录→分节→要点）、段落留白合理、不堆大段文字；\n② 视觉：统一配色（选 1 个主色 + 中性色）、标题与正文对比、善用表格/列表/图标分层；\n③ 内容密度：要点化（每页/每节 3-6 个要点），不要整页文字墙；\n④ PPT 尤其：每页一个主题，标题醒目，正文精简，善用配图/色块。",
                         "【PPT 配图规则】用户要求 PPT 带图片/截图时：优先用 playwright-mcp 的浏览器工具"
                         "（playwright-mcp_browser_navigate 打开网页 → playwright-mcp_browser_take_screenshot 截图），"
                         "把截图通过 plugin_ppt_create_pptx 的每页 image 参数嵌入；也可用图片 URL。"
                         "搜索工具只返回文字，无法获得图片——需要图片必须用浏览器截图或图片 URL。",
                         "【权限与自主性】当 write_files / run_command 权限为『允许』时，文件写入、保存、"
                         "命令执行等操作你【可以直接做决定并执行】，不需要先征求用户同意——用户已授权你自主工作。"
                         "用 ask_user 要克制：需求明确或可合理推断时直接做，只有真正拿不定（多方案重大分歧、"
                         "涉及删除/覆盖重要文件、不可逆或危险操作）才询问。",
                          "【步骤规划规范】区分『简单任务』与『复杂任务』（最重要的判断）：\n"
                          "一、简单任务（问候/闲聊/简单解释/简单计算/翻译/常识问答/单步操作）→ 直接回答，不规划、不调用工具。\n"
                          "二、复杂任务（调研分析/代码修改/科研实验/仿真计算/软件操作/多文件项目/多步流程）→ 严格走：\n"
                          "① 先给【方案设计】（对标 opencode Plan）：列出 2-3 个可行方案，各含思路、关键步骤、"
                          "优劣（成本/效果/风险），并给出推荐；\n"
                          "② 用 ask_user 让用户选择方案（或确认推荐）——除非用户已明确指定做法；\n"
                          "③ 确认后输出【步骤规划】块——把全部步骤一气呵成列出（循序渐进、直达目标，不是一步步挤）：\n"
                          "【步骤规划】\n1. <第一步>\n2. <第二步>\n…\nN. 验证（最后一步一定是验证："
                          "重读文件/自查结果是否符合要求）\n"
                          "④ 然后按规划逐步执行；每完成一步输出一行【步骤完成：N】；\n"
                          "⑤ 全部完成输出【全部完成】。\n"
                          "单步即可完成的任务无需规划，直接执行即可。\n"
                          "判断标准：需要『方案抉择』的（调研方向、修改思路、实验设计、工具选型）→ 必出方案先选；"
                          "执行路径唯一清晰的 → 直接规划步骤执行，不必问。",
                         "【交付反馈】任务完成后，在回答末尾给出简短反馈（对标 Sol 的 verifier 自查机制）："
                         "① 完成了什么 ② 如何验证的（重读/测试/对比）③ 是否完全满足要求、有无未完成或遗留事项。"
                         "交付前先自查：有没有遗漏用户的核心要求？结果是否符合预期？如有偏差如实说明并给出下一步。",
                         "【简单问题直接回答】问候、闲聊、简单解释、简单计算、翻译、常识问答等无需工具的请求："
                         "直接回答，不要调用任何工具、不要搜索、不要规划、不要走工具循环——一个回复搞定。"
                         "只在确实需要文件读写/联网搜索/执行命令/生成文档/处理图片时才调用工具。"
                         "【技能注入是辅助参考】你收到的技能说明只是参考规范，不代表用户要求执行该技能："
                         "用户的主意图明确（如『做 PPT/写文档/查资料』）时按主意图执行，"
                         "绝不要因为注入了多个技能而向用户确认『要执行哪个技能』。",
                         "【果断行动】默认直接执行、自己拿主意，不犹豫（对标 opencode 的行动风格）：\n① 需求可合理推断 → 按最合理理解直接做，做完说明你的假设；\n② 不确定的细节影响不大 → 用合理默认值（主题/风格/格式等选常见方案），不要打断用户问；\n③ 只有真正重大时才用 ask_user：用户意图完全无法判断、多方案成本差异巨大、涉及删除/覆盖重要文件或不可逆/危险操作。\n你是行动的 agent，不是犹豫的：宁可多做事后修正，也不要频繁弹窗打断用户。",
"【按原话执行】按用户字面要求执行，不要联想扩展需求：用户说『写两个公式』就只写公式，不要问/不要主动做成 Word/PPT/文件；用户没要求格式就不要调用文档工具。只有用户明确要求（『做成文档』『保存』『生成 PPT』『给个文件』）时才调用对应工具。不要替用户决定他们没要求的下一步。",
"【无法完成时的处理】用户提出需求而当前能力不足/工具缺失/技术受限时，"
                         "绝不要一句『做不到/不支持』了事，要往能完成的方向走："
                         "① 分析需求，找出可行路径：拆解任务、换一种实现方式、分步完成、或做简化版先交付；"
                         "② 与用户交流方案：说明限制和可行的替代方案，让用户选择方向；"
                         "③ 缺失能力时先使用已安装工具或给出可执行替代方案；需要新增插件、技能、MCP 或依赖时，"
                         "说明用途、来源、权限和费用，交由用户在设置/项目代码中显式安装，不得自行改写运行时扩展；"
                         "④ 持续沟通：每一步进展告诉用户，方向不对及时调整，目标是把需求真正落地。",
                         "【自我修复】你具备自我修复能力：如果发现自身组件有问题（插件加载失败、MCP 连接异常、"
                         "工具缺失/调用失败、技能异常、代码报错等），应主动诊断并修复，不要坐等问题影响使用。"
                         "修复原则：小问题（参数错误、配置缺失、简单代码修复、重启单个组件）直接自己修复，"
                         "并在回答中告知用户你发现了什么问题、做了什么修复；大问题（需要重启整个服务、"
                         "修改核心文件、安装依赖、涉及数据安全或不可逆操作）必须先告知用户问题与修复方案，"
                         "用 plugin_vision_see_image 无关的 ask_user 弹窗征得用户同意后再动手。"
                         "诊断工具：plugin_management_list_extensions 查插件、/api/mcp/servers 查 MCP、"
                         "/api/skills 查技能、plugin_workspace_write_file 修改当前任务工作区内代码、"
                         "plugin_terminal_run_command 跑诊断命令（须确认）。",
                         "【看图规则】see_image 工具只对【无法识别图像的模型】管用（借助外部视觉模型看图）。"
                         "若你的模型本身支持视觉输入（多模态，能直接读取图片附件），就不要调用 see_image，直接描述图片内容。"
                         "当消息里出现「图片附件：xxx.png」之类的提示、或用户发送/引用了图片（包括报错提到图片文件）"
                         "而你无法直接看图时，必须调用 plugin_vision_see_image 工具（传文件路径/文件名）查看后再回答，"
                         "绝不要直接说看不到。如果用户提到『这张图/这个截图/图片里/识别图像/看图』等字眼但没附图片，"
                         "应询问图片位置（文件路径），并用 see_image 主动查看。"]
            # 注入历史教训（跨会话学习：之前犯过的错避免重犯）+ 记忆图跨会话语义召回
            # 先取最后用户消息作为记忆图查询（提前计算，供 lessons 注入使用）
            _mem_query = ""
            for _m in reversed(msgs):
                if _m.get("role") == "user" and isinstance(_m.get("content"), str):
                    _mem_query = _m["content"]
                    break
            _lessons_prompt = _lesson_system_prompt(_mem_query)
            if _lessons_prompt:
                sys_parts.append(_lessons_prompt)
            for p in history_store.list_projects():
                if p.get("id") == req.project:
                    if p.get("path"):
                        effective_workspace = current_workspace.get() or p["path"]
                        isolation_note = "（该会话的独立 Git worktree）" if effective_workspace != p["path"] else ""
                        sys_parts.append(f"当前项目：{p.get('name', '')}，任务工作区：{effective_workspace}{isolation_note}（文件和终端工具只能在此任务工作区操作）")
                        # 项目规则文件（对标 opencode AGENTS.md）：存在则自动加载，让 AI 遵循项目约定
                        _rules = _load_project_rules(effective_workspace)
                        if _rules:
                            sys_parts.append(f"【项目规则（{_rules['file']}）】\n{_rules['content']}")
                    else:
                        sys_parts.append(f"当前项目：{p.get('name', '')}（未设置项目文件夹）")
                    break
            # Only advertise generic skill names. Bodies are loaded on demand by skill_load;
            # injecting every body/summary each turn wastes prompt tokens and blurs priorities.
            generic = skills_loader.generic_skills()
            if generic:
                sys_parts.append("可按需加载的通用技能：" + ", ".join(g["name"] for g in generic)
                                 + "。仅在当前任务相关时调用 skill_load。")
            # 计划模式（对标 opencode Plan/Build）：只读规划 → 输出计划 → ask_user 确认 → 才执行
            if req.mode == "plan":
                sys_parts.append(
                    "【计划模式】你当前处于计划模式（Plan）：\n"
                    "① 只做分析和规划：绝不修改文件、绝不执行有副作用的操作（不写文件、不删改、不提交）；"
                    "可读取文件、可联网搜索（搜集信息辅助规划）；\n"
                    "② 理解需求后，输出清晰的实施计划：目标 → 步骤（含涉及文件/要调用的工具）→ 风险/注意点；\n"
                    "③ 计划输出完毕后，必须调用 ask_user 询问用户『是否按此计划执行』；\n"
                    "④ 用户确认后（回答是/可以/执行），你再切换执行模式按计划逐步实施；"
                    "用户修改计划则按新指示调整。")
            msgs = [{"role": "system", "content": "\n\n".join(sys_parts)}] + msgs
            # 动态内容追加到末尾 user 消息（不破坏 system 前缀缓存）
            dynamic_parts = []
            if req.online:
                dynamic_parts.append("【联网模式已开启】你可以使用搜索工具（工具名 websearch_web_search，"
                                     "一次调用即多源聚合：必应+百度+搜狗）。"
                                     "【搜索规范——语义第一，绝不拆词】"
                                     "① 先理解再搜：搜索前先想清楚用户【真正想问什么】，把整个意图浓缩成"
                                     "【一个】语义完整、精准的搜索查询（完整问句或完整概念），"
                                     "绝不把需求拆成零散关键词（如用户问『光学工程专业的概念』→ 搜『什么是光学工程专业』，"
                                     "而不是『光学』『工程』；用户问前景 → 搜『光学工程专业的前景是什么』；"
                                     "用户问区别 → 搜『A与B的区别』）。"
                                     "② 一次搜索搞定：query 要一步到位，一次调用返回多源结果；"
                                     "不要反复调用同一工具搜近似词。"
                                     "③ 反馈机制：搜索后说明『搜了什么、结果是否相关』；"
                                     "结果不满意就换一个更精准的完整查询重搜（最多 2-3 次），"
                                     "若仍无相关结果必须如实告知用户『未能找到相关信息』，绝不编造。"
                                     "涉及新闻/时效/事实核查的问题，务必先搜索再回答，并注明信息来源。")
            last_user = ""
            for m in reversed(msgs):
                if m.get("role") == "user" and isinstance(m.get("content"), str):
                    last_user = m["content"]
                    break
            # 技能激活 = 字面候选召回 + LLM 语义判断（意思相同才激活，不是字面相同）
            # 反馈闭环（P3）：用户纠正 → 上一轮注入的技能负反馈
            if any(_h in last_user for _h in CORRECTION_HINTS):
                _record_skill_feedback(False)
            candidate_skills = [sk for sk in skills_loader.match_skills(last_user, max_n=3)
                                if sk["name"] not in skills_loader.GENERIC_SKILL_NAMES]
            # P3：按历史反馈分排序候选（高分优先，低分淘汰）
            if candidate_skills:
                _names = [sk["name"] for sk in candidate_skills]
                _ordered, _dropped = _skill_priority(_names)
                if _dropped:
                    _by_name = {sk["name"]: sk for sk in candidate_skills}
                    candidate_skills = [_by_name[n] for n in _ordered if n in _by_name]
            # Local deterministic recall avoids a second paid model request on every user turn.
            matched_skills = candidate_skills[:skill_count]
            if matched_skills:
                _record_skill_usage([sk["name"] for sk in matched_skills])   # 自演化 P2：记录激活，驱动生命周期
                # 轻量注入（缓存友好）：只注入技能【名称清单】，不注入正文。
                # 正文由 skill_load 工具按需获取 → user 消息稳定 → 前缀缓存命中；
                # 技能正文不进 prompt → 每轮省数万 tokens。
                names = ", ".join(sk["name"] for sk in matched_skills)
                dynamic_parts.append("本次任务涉及以下技能（来自技能库）：" + names
                                     + "。如需遵循其中规范，调用 skill_load 工具按名加载对应技能正文再执行。")
            if dynamic_parts and msgs and msgs[-1].get("role") == "user" and isinstance(msgs[-1].get("content"), str):
                msgs[-1]["content"] += "\n\n" + "\n\n".join(dynamic_parts)
            tools = await mcp_openai_tools()  # MCP 工具（无服务器时返回 None）
            plugin_tools = plugins_loader.openai_tools()  # 插件工具
            flatten_map = {}
            # ---- 缓存优化：工具列表按对话锁定（同一对话内稳定）----
            # OpenAI/zen 前缀缓存要求 tools 参数一致；工具随消息变化会破坏前缀缓存。
            # 首次请求确定工具集（按消息匹配防降智），同对话后续轮次复用 → 前缀稳定、命中率最大化。
            # key 含联网状态：开启/关闭联网会重新匹配（联网时搜索工具必须可用）。
            # The selected set follows the current intent. Locking it only by conversation made
            # first-turn tools persist forever, so later tasks could not access their tools.
            intent_key = hashlib.sha256(last_user.strip().lower().encode("utf-8")).hexdigest()[:16]
            tools_key = req.provider + "|" + str(req.online) + "|" + intent_key
            if tools_key in _TOOLS_CACHE:
                entry = _TOOLS_CACHE[tools_key]
                tools = entry["tools"]
                flatten_map = dict(entry.get("flatten") or {})
            else:
                if tools or plugin_tools:
                    # 插件（软件自身能力，含自我改进）优先于 MCP 外部工具
                    tools = (plugin_tools or []) + (tools or [])
                    # 工具 Schema 展平（对标 Reasonix：DeepSeek 对深/宽 schema 静默丢参）
                    tools, flatten_map = _flatten_tools(tools)
                    # 按消息智能匹配工具（对标 opencode：只注入相关工具，防无关工具干扰降智）
                    matched = _match_tools_by_message(tools, last_user, max_tools=max_tools)
                    if matched:
                        tools = matched
                    # 联网模式无需重排工具：websearch_web_search 已在 CORE_TOOL_NAMES 核心保底
                    #（永远注入），重排只会改变 tools 字节顺序、破坏前缀缓存（对标 Reasonix：
                    # 工具定义顺序必须字节稳定，缓存才命中）。
                    # 工具数量不限：按消息匹配后全量保留（模型需要哪个都能用）
                _TOOLS_CACHE[tools_key] = {"tools": tools, "flatten": flatten_map}
                if len(_TOOLS_CACHE) > 64:
                    _TOOLS_CACHE.clear()
            # ask_user：弹窗询问用户（对标 opencode 的 AskUserQuestion）
            tools = (tools or []) + [{
                "type": "function",
                "function": {
                    "name": "ask_user",
                    "description": "弹窗询问用户。⚠️ 使用原则：不清楚就问！当用户意图不明确、要求含糊、"
                                   "存在多个合理解释、或缺少关键信息（主题/风格/参数/用途等）时，必须调用本工具确认，"
                                   "绝不要擅自假设。仅当任务完全无歧义时才直接执行。"
                                   "options 是建议选项（最多 8 个），用户可点选或自由输入。",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "question": {"type": "string", "description": "要问用户的问题"},
                            "options": {"type": "array", "items": {"type": "string"}, "description": "建议选项（最多5个）"}
                        },
                        "required": ["question"],
                    },
                },
            }]

            # 工具集解析完成后才能准确计算上下文；旧实现提前读取 tools，异常被吞后永不压缩。
            _existing_conv = history_store.get_conversation(
                req.conversation_id, req.project) if req.conversation_id else None
            _observed_context = ((_existing_conv or {}).get("usage") or {}).get("context", 0)
            msgs, _did_compact, _compact_before, _compact_after = compact_messages_if_needed(
                msgs, tools, ctx_limit, _observed_context)

            called_tools = {}   # 工具去重计数：(name|args) -> 次数（风暴抑制：同参数反复调用 ≥3 次强停）
            tool_fail_count = {}   # 工具失败计数（教训机制：连续失败 → 记教训）
            continues = 0
            _retried_compress = False   # 上下文超限已压缩重试（只允许一次，防死循环）
            _drift_flag = False
            # ---- 前缀指纹：检测前缀漂移（对标 Reasonix ImmutablePrefix 的 SHA-256 钉住）----
            # system（msgs[0]）+ tools 字节变化 → 缓存必然全量 miss；记录变化以便诊断
            try:
                _sys_text = msgs[0]["content"] if msgs and msgs[0].get("role") == "system" else ""
                _fp = hashlib.sha256((_sys_text + "\x00" + json.dumps(tools, ensure_ascii=False, sort_keys=True)).encode("utf-8")).hexdigest()[:16]
                _cid_k = req.conversation_id or "new"
                _prev = _PREFIX_FINGERPRINTS.get(_cid_k)
                if _prev and _prev != _fp:
                    _drift_flag = True
                    print(f"[cache-drift] 对话 {_cid_k} 前缀指纹变化 {_prev} -> {_fp}（本请求缓存将全量 miss）", flush=True)
                _PREFIX_FINGERPRINTS[_cid_k] = _fp
            except Exception:
                pass
            # 远端模型不设输出上限：超长文档/代码/工具参数完整输出，不被截断。
            # 防失控靠：MAX_LOOPS（工具/续写轮上限）+ 同参数工具去重提示。
            create_kwargs = {}
            # 深度思考（reasoning_effort，轻/中/深）；不支持的供应商报错时自动去掉重试
            if req.reasoning in ("low", "medium", "high"):
                create_kwargs["reasoning_effort"] = req.reasoning
            # 路由亲和（对标 opencode promptCacheKey=sessionID）：同会话请求附加缓存键，
            # 支持路由亲和的网关/代理可显著提升前缀缓存命中率；不支持的端点会忽略该字段（已验证无害）。
            # 本地模型（llama.cpp/ollama）无前缀缓存意义，跳过。
            if req.provider not in ("local", "ollama"):
                create_kwargs["extra_body"] = {"promptCacheKey": (req.conversation_id or "new")[:64]}
            # 累计各轮真实用量（含缓存命中；拿不到时兜底估算）
            # context = 当前上下文实际占用（最后一次 prompt_tokens，覆盖式：每轮 prompt 都是完整上下文）
            usage_totals = {"input": 0, "output": 0, "cached": 0, "reasoning": 0, "context": 0}
            _secondary_usage = {}

            def _priced_usage(delta_usage):
                delta_usage = dict(delta_usage or {})
                try:
                    from pricing import calc_cost, get_price_source
                    _cost, _est = calc_cost(
                        req.provider, cfg["model"], delta_usage.get("input", 0),
                        delta_usage.get("output", 0), delta_usage.get("cached", 0))
                    if _cost is None:
                        delta_usage.pop("cost", None)
                        delta_usage["cost_unavailable"] = True
                    else:
                        delta_usage["cost"] = _cost
                        delta_usage.pop("cost_unavailable", None)
                    delta_usage["cost_est"] = _est
                    delta_usage["cost_source"] = get_price_source(req.provider, cfg["model"])
                except Exception:
                    pass
                delta_usage["last_input"] = delta_usage.get("input", 0)
                delta_usage["last_cached"] = delta_usage.get("cached", 0)
                return delta_usage

            def _commit_secondary_usage(report):
                if not report:
                    return
                provider = report["provider"] or req.provider
                model = report["model"] or cfg["model"]
                priced = history_store.price_usage_delta(report["usage"], provider, model)
                merged = history_store.merge_cumulative_usage(_secondary_usage, priced)
                _secondary_usage.clear()
                _secondary_usage.update(merged)
                if req.conversation_id and not report.get("usage_committed"):
                    history_store.add_conversation_usage(
                        req.conversation_id, report["usage"], project=req.project,
                        provider=provider, model=model)

            for _round in range(MAX_LOOPS):
                try:
                    stream = await async_client.chat.completions.create(
                        model=cfg["model"], messages=msgs, stream=True, tools=tools,
                        **create_kwargs,
                    )
                except Exception as e:
                    # 供应商不支持 reasoning_effort → 去掉该参数重试一次
                    if "reasoning" in str(e).lower() and "reasoning_effort" in create_kwargs:
                        create_kwargs.pop("reasoning_effort", None)
                        stream = await async_client.chat.completions.create(
                            model=cfg["model"], messages=msgs, stream=True, tools=tools,
                            **create_kwargs,
                        )
                    else:
                        # 上下文超限（配置窗口 > 模型实际窗口）→ 自动压缩后重试一次，不打断用户
                        err_l = str(e).lower()
                        if (any(k in err_l for k in ("context length", "maximum context",
                                                     "context_length", "too many tokens",
                                                     "max context", "exceeds", "超出上下文"))
                                and len(msgs) > 8 and not _retried_compress):
                            _retried_compress = True
                            half = max(1, (len(msgs) - 1) // 2)
                            msgs = msgs[:1] + msgs[-half:]
                            last = msgs[-1]
                            if isinstance(last.get("content"), str):
                                last["content"] += ("\n\n（提示：上下文超限，已自动压缩保留最近一半消息，"
                                                    "如需回顾早期内容请直接询问。）")
                            continue
                        raise
                acc = ""
                tool_calls = {}  # index -> {"id", "name", "args"}
                finish_reason = None
                _last_save_t = time.time()
                async for chunk in stream:
                    # 流末尾的 usage chunk（choices 为空）→ 收集真实用量
                    u = getattr(chunk, "usage", None)
                    if u is not None:
                        det = getattr(u, "prompt_tokens_details", None) or {}
                        cached = 0
                        if isinstance(det, dict):
                            cached = det.get("cached_tokens", 0) or 0
                        else:
                            cached = getattr(det, "cached_tokens", 0) or 0
                        # 思考过程 token（统计时不计入正式输出）
                        cdet = getattr(u, "completion_tokens_details", None) or {}
                        reasoning = 0
                        if isinstance(cdet, dict):
                            reasoning = cdet.get("reasoning_tokens", 0) or 0
                        else:
                            reasoning = getattr(cdet, "reasoning_tokens", 0) or 0
                        if not reasoning:
                            reasoning = getattr(u, "reasoning_tokens", 0) or 0
                        usage_totals["input"] += getattr(u, "prompt_tokens", 0) or 0
                        usage_totals["output"] += getattr(u, "completion_tokens", 0) or 0
                        usage_totals["cached"] += cached
                        usage_totals["reasoning"] += reasoning
                        # 当前上下文占用 = 该轮完整 prompt；取全程最大值（agent 循环内
                        # 每轮 prompt 因工具消息追加而增长；取峰值避免"5万↔3万"锯齿）
                        _pt = getattr(u, "prompt_tokens", 0) or 0
                        usage_totals["context"] = max(usage_totals["context"], _pt)
                        # Token 估算自适应校准：用真实 prompt_tokens 回校本地估算系数
                        try:
                            _ptext = json.dumps(msgs, ensure_ascii=False)
                            _cj = len(re.findall(r"[\u4e00-\u9fff]", _ptext))
                            _wd = len(re.findall(r"[a-zA-Z0-9_]+", _ptext))
                            _ot = len(_ptext) - _cj - sum(len(w) for w in re.findall(r"[a-zA-Z0-9_]+", _ptext))
                            _calibrate_token_estimate(_pt, len(_ptext), _cj, _wd, _ot)
                        except Exception:
                            pass
                    if not chunk.choices:
                        continue
                    delta = chunk.choices[0].delta
                    if chunk.choices[0].finish_reason:
                        finish_reason = chunk.choices[0].finish_reason
                    content = getattr(delta, "content", None)
                    if content:
                        acc += content
                        yield f"data: {json.dumps({'delta': content}, ensure_ascii=False)}\n\n"
                    extra = getattr(delta, "model_extra", None) or {}
                    thinking = extra.get("reasoning_content")
                    if thinking:
                        yield f"data: {json.dumps({'think': thinking}, ensure_ascii=False)}\n\n"
                    for tc in getattr(delta, "tool_calls", None) or []:
                        slot = tool_calls.setdefault(tc.index, {"id": "", "name": "", "args": ""})
                        if tc.id:
                            slot["id"] = tc.id
                        if tc.function and tc.function.name:
                            slot["name"] = tc.function.name
                        if tc.function and tc.function.arguments:
                            slot["args"] += tc.function.arguments

                if not tool_calls:
                    # 回答被截断 → 自动续写，不打断用户（对标 opencode 自动衔接）
                    # 仅在提供方明确返回 length 时续写，避免因代码块/列表结尾被误判而重复生成。
                    truncated = finish_reason == "length"
                    if truncated and continues < MAX_CONTINUE:
                        continues += 1
                        msgs.append({"role": "assistant", "content": acc})
                        msgs.append({"role": "user", "content": CONTINUE_MARKER})
                        continue
                    if acc.strip():
                        msgs.append({"role": "assistant", "content": acc})
                    break  # 正常结束

                # 有工具调用：把 assistant 消息（含 tool_calls）记入历史
                ordered = [tool_calls[i] for i in sorted(tool_calls)]
                msgs.append({
                    "role": "assistant",
                    "content": acc or None,
                    "tool_calls": [
                        {"id": s["id"], "type": "function",
                         "function": {"name": s["name"], "arguments": s["args"]}}
                        for s in ordered
                    ],
                })
                for s in ordered:
                    yield f"data: {json.dumps({'tool': {'name': s['name'], 'phase': 'start', 'args': (s['args'] or '')[:500]}}, ensure_ascii=False)}\n\n"
                # ---- 并行执行独立工具（提高效能）：非 ask_user 工具并行跑，
                # ask_user 需要用户交互，单独串行处理 ----
                async def _exec_tool(s):
                    # 在任何解析、缓存或分发分支之前初始化；所有工具共用此路径。
                    dup_key = s["name"] + "|" + (s.get("args") or "")
                    cnt = called_tools.get(dup_key, 0) + 1
                    called_tools[dup_key] = cnt
                    if cnt >= 3:
                        return (s, None,
                                "已阻止第 3 次相同工具与参数调用：疑似死循环，请换工具、换参数或结束任务。",
                                False)
                    try:
                        try:
                            args = json.loads(s["args"] or "{}")
                        except Exception as je:
                            # 参数 JSON 解析失败：多半是内容过长被截断 → 先自动修复（对标 Reasonix
                            # Truncation Repair：补引号/括号），修复不了则触发"续写补齐"（让模型
                            # 紧接不完整参数继续输出剩余部分，而不是重新调用——重新调用同样会截断）
                            args = _repair_tool_args(s["args"] or "")
                            if args is None:
                                return (s, None,
                                        "工具参数 JSON 不完整；请缩短参数或拆成多次调用后重试。",
                                        False)
                        # 还原展平后的点号参数（对标 Reasonix dispatch 时 re-nest）
                        _map = flatten_map.get(s["name"])
                        if _map:
                            args = _renest_args(args, _map)
                        decision = evaluate_permission(
                            s["name"], args, load_runtime_policy(_settings_file()))
                        if decision.effect == "deny":
                            denied = json.dumps({
                                "error": "权限策略拒绝了该工具调用",
                                "permission": decision.__dict__,
                            }, ensure_ascii=False)
                            return (s, denied, None, False)
                        if decision.effect == "ask" and not args.get("_confirmed"):
                            confirmation = json.dumps({
                                "error": "该工具调用需要用户确认",
                                "need_confirmation": True,
                                "safety": "权限规则 %s 要求确认：%s" % (decision.rule, s["name"]),
                            }, ensure_ascii=False)
                            return (s, confirmation, None, True)
                        # 工具结果缓存已移除（用户要求与云端一致：每次工具调用都真实执行，
                        # 云端 usage/计费/日志完整记录，不做本地复用）
                        if s["name"].startswith("plugin_"):
                            result_text = await asyncio.to_thread(plugins_loader.call, s["name"], args)
                        else:
                            result_text = await mcp_call(s["name"], args)
                        # Some explicit tools call a second model (vision/model delegation).
                        # Persist provider-reported usage immediately so cancellation of the
                        # outer chat stream cannot make billed tokens disappear. Tools that
                        # already committed usage set usage_committed=true: they are included
                        # in the live display but not written a second time.
                        _tool_usage = reported_tool_usage(result_text)
                        _commit_secondary_usage(_tool_usage)
                        # ---- 系统级权限确认标记（对标 opencode：权限弹窗是系统行为）----
                        # 插件返回 need_confirmation（高危/敏感操作）→ 标记由外层串行确认
                        if isinstance(result_text, str) and "need_confirmation" in result_text:
                            try:
                                _nc = json.loads(result_text)
                                if _nc.get("need_confirmation"):
                                    return (s, result_text, None, True)   # 第4位=需权限确认
                            except Exception:
                                pass
                        # 风暴抑制（对标 Reasonix Storm Detection）：同 (name,args) 反复调用 ≥3 次 → 强停提示
                        if cnt == 2:
                            result_text += "\n（提示：该工具与之前相同参数已调用过，如结果无变化请勿重复调用）"
                        return (s, result_text, None, False)
                    except Exception as e:
                        # 教训机制：同工具连续失败 → 记教训（模型下次避免）
                        tool_fail_count[s["name"]] = tool_fail_count.get(s["name"], 0) + 1
                        _fail_cnt = tool_fail_count[s["name"]]
                        if _fail_cnt >= 2:
                            _save_lesson(f"工具 {s['name']} 连续失败 {_fail_cnt} 次（错误：{str(e)[:80]}），注意参数格式与内容长度",
                                         source="tool_fail")
                        # 失败升级（对标 Reasonix：3+ 错误强制换思路）——连续失败说明方案不可行
                        _err_msg = f"工具错误: {e}"
                        if _fail_cnt >= 3:
                            _err_msg += ("\n（⚠️ 该工具已连续失败 3 次，说明当前方案不可行。"
                                         "请换一种实现思路：换工具 / 换参数 / 拆解问题，不要重复同一失败方案。）")
                        return (s, None, _err_msg, False)
                par_calls = [s for s in ordered if s["name"] != "ask_user"]
                ask_calls = [s for s in ordered if s["name"] == "ask_user"]
                if par_calls:
                    par_results = await asyncio.gather(*[_exec_tool(s) for s in par_calls])
                else:
                    par_results = []
                for item in par_results:
                    s = item[0]; result_text = item[1]; err = item[2]
                    need_confirm = len(item) > 3 and item[3] == True
                    # ---- 系统级权限确认（对标 opencode：权限弹窗是系统行为，不依赖模型自觉）----
                    # 工具返回 need_confirmation（高危/敏感操作）→ 自动 ask_user 弹窗，用户确认后才放行
                    if need_confirm:
                        _perm_question = f"⚠️ 需要确认：工具 {s['name']} 检测到高危/敏感操作，是否允许执行？"
                        try:
                            _nc = json.loads(result_text)
                            if _nc.get("safety"):
                                _perm_question = f"⚠️ 需要确认：{_nc.get('safety')}\n是否允许执行？"
                        except Exception:
                            pass
                        qid = uuid.uuid4().hex[:8]
                        fut = asyncio.get_running_loop().create_future()
                        PENDING_ASKS[qid] = fut
                        yield f"data: {json.dumps({'ask_user': {'id': qid, 'question': _perm_question,
                                                                   'options': ['允许执行', '拒绝']}}, ensure_ascii=False)}\n\n"
                        try:
                            answer = await asyncio.wait_for(fut, timeout=180)
                        except asyncio.TimeoutError:
                            answer = "拒绝"
                        finally:
                            PENDING_ASKS.pop(qid, None)
                        if answer and ("允许" in str(answer) or "同意" in str(answer) or "是" in str(answer)):
                            # 用户允许 → 带 _confirmed 标记重调（插件确认后执行）
                            try:
                                _args2 = dict(json.loads(s["args"] or "{}"))
                                _args2["_confirmed"] = True
                                if s["name"].startswith("plugin_"):
                                    result_text = await asyncio.to_thread(plugins_loader.call, s["name"], _args2)
                                else:
                                    _args2.pop("_confirmed", None)
                                    result_text = await mcp_call(s["name"], _args2)
                                _confirmed_usage = reported_tool_usage(result_text)
                                _commit_secondary_usage(_confirmed_usage)
                                err = None
                            except Exception as e2:
                                err = str(e2)
                        else:
                            result_text = f"（权限确认被拒绝：用户不同意执行 {s['name']}，请勿重试，向用户说明或换方案）"
                    _rt = result_text if result_text is not None else (err or "")
                    # headroom 轻量压缩（Rust SmartCrusher）：仅压缩大 JSON 工具结果（数组/对象，
                    # 压缩后更小才采用）。文档全文/代码等非 JSON 内容原样保留（不截断）。
                    try:
                        from deps.headroom_compressor import compress_json
                        _rt = compress_json(_rt, min_len=600)
                    except Exception:
                        pass
                    msgs.append({"role": "tool", "tool_call_id": s["id"], "content": _rt})
                    # 提取文件变更 diff（前端渲染可展开的更改视图）
                    _diff = None
                    if "【文件变更】" in _rt:
                        _i = _rt.index("【文件变更】")
                        _diff = _rt[_i + len("【文件变更】"):][:4000]
                    yield f"data: {json.dumps({'tool': {'name': s['name'], 'phase': 'done', 'ok': err is None, 'result': _rt[:100], 'diff': _diff}}, ensure_ascii=False)}\n\n"
                for s in ask_calls:
                    try:
                        args = json.loads(s["args"] or "{}")
                        # 弹窗询问用户：发事件给前端，等回答
                        qid = uuid.uuid4().hex[:8]
                        fut = asyncio.get_running_loop().create_future()
                        PENDING_ASKS[qid] = fut
                        question = str(args.get("question", "请确认一下"))
                        options = (args.get("options") or [])[:8]
                        yield f"data: {json.dumps({'ask_user': {'id': qid, 'question': question, 'options': options}}, ensure_ascii=False)}\n\n"
                        try:
                            answer = await asyncio.wait_for(fut, timeout=180)
                            result_text = f"用户的回答是：{answer}"
                        except asyncio.TimeoutError:
                            result_text = "用户没有在 3 分钟内回答（可能已离开），请根据已有信息继续或做合理假设。"
                        finally:
                            PENDING_ASKS.pop(qid, None)
                        msgs.append({"role": "tool", "tool_call_id": s["id"], "content": result_text})
                        yield f"data: {json.dumps({'tool': {'name': s['name'], 'phase': 'done', 'ok': True, 'result': result_text[:100]}}, ensure_ascii=False)}\n\n"
                    except Exception as e:
                        msgs.append({"role": "tool", "tool_call_id": s["id"], "content": f"工具错误: {e}"})
                        yield f"data: {json.dumps({'tool': {'name': s['name'], 'phase': 'done', 'ok': False, 'result': str(e)[:100]}}, ensure_ascii=False)}\n\n"
            else:
                yield f"data: {json.dumps({'error': '工具调用轮数超限'}, ensure_ascii=False)}\n\n"
                return
            # ---- 合并截断续写产生的消息（历史干净） ----
            msgs = _merge_continuations(msgs)
            # ---- 用量统计：真实数据优先（含缓存命中）；思考 token 单独记录，不计入正式输出 ----
            # context = 当前对话在模型里的实际占用（不含思考），不是累计消耗
            real = usage_totals.get("input") or usage_totals.get("output")
            if real:
                ctx_used = usage_totals.get("context") or usage_totals["input"]
                usage = {"input": usage_totals["input"], "output": usage_totals["output"],
                         "cached": usage_totals["cached"], "reasoning": usage_totals["reasoning"],
                         "context": ctx_used}
            else:
                prompt_text = json.dumps(msgs, ensure_ascii=False) + json.dumps(tools or [], ensure_ascii=False)
                est_in = estimate_tokens(prompt_text)
                usage = {"input": est_in, "output": estimate_tokens(acc), "cached": 0,
                         "reasoning": 0, "context": est_in}
            # 缓存告警：前缀漂移 或 大请求几乎零命中（缓存被破坏/冷启动）→ 前端提示（对标 opencode #40796）
            usage["cache_drift"] = bool(_drift_flag)
            usage["cache_bust"] = bool(usage["input"] > 20000 and usage["cached"] < usage["input"] * 0.1)
            usage["compacted"] = bool(_did_compact)
            if _did_compact:
                usage["context_before_compact"] = _compact_before
                usage["context_after_compact"] = _compact_after
            usage = _priced_usage(usage)
            display_usage = history_store.merge_cumulative_usage(
                history_store.price_usage_delta(usage, req.provider, cfg["model"]),
                _secondary_usage)
            if _secondary_usage:
                display_usage["cost_source"] = {
                    "kind": "mixed_model_ledger", "invoice_reconciled": False,
                    "entries": len(_secondary_usage.get("by_model") or {}),
                }
            yield f"data: {json.dumps({'usage': display_usage}, ensure_ascii=False)}\n\n"
            # ---- 权威提交：只追加完整 assistant，绝不把压缩后的模型 prompt 覆盖到历史。----
            if not req.persist:
                if req.conversation_id:
                    history_store.add_conversation_usage(
                        req.conversation_id, usage, req.project, req.provider, cfg["model"])
                yield 'data: {"done": true}\n\n'
                return
            final_text = ""
            for final_message in reversed(msgs):
                if (final_message.get("role") == "assistant"
                        and isinstance(final_message.get("content"), str)
                        and final_message.get("content")):
                    final_text = final_message["content"]
                    break
            if not final_text.strip():
                raise RuntimeError("模型未产生可提交的完整回复")
            history_store.append_message(
                req.conversation_id or None,
                {"id": "m-" + request_id, "role": "assistant", "content": final_text,
                 "ts": int(time.time() * 1000)},
                project=req.project, provider=req.provider, model=cfg["model"],
                usage_delta=usage,
            )
            yield 'data: {"done": true}\n\n'
        except asyncio.CancelledError:
            # 会话切换/停止：用户消息已单独持久化；未完成 assistant 永不落库。
            return
        except Exception as e:
            import traceback
            traceback.print_exc()   # 落到日志文件，方便排查
            yield f"data: {json.dumps({'error': str(e)}, ensure_ascii=False)}\n\n"
        finally:
            with _ACTIVE_CHAT_STREAMS_LOCK:
                if _ACTIVE_CHAT_STREAMS.get(stream_key) is stream_task:
                    _ACTIVE_CHAT_STREAMS.pop(stream_key, None)
                _CANCELLED_CHAT_STREAMS.pop(stream_key, None)
            if workspace_token is not None:
                current_workspace.reset(workspace_token)
            current_tenant.reset(tenant_token)

    return StreamingResponse(
        gen(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# AI 交付文件的下载区（deliver 插件写入，/files 直接可下载）
# 注意：必须挂在 "/" 通配挂载之前，否则会被前端静态路由吞掉
# 下载区：打包版用数据目录 files（可写、持久；开发版用项目目录）
@app.get("/files/{name}")
def tenant_file_response(name: str):
    """Serve only the current tenant's attachment; never bypass tenant lookup."""
    if not name or name != os.path.basename(name) or name in (".", ".."):
        raise HTTPException(status_code=400, detail="非法文件名")
    path = os.path.join(_files_dir(), name)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="文件不存在")
    response = FileResponse(path)
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["Content-Security-Policy"] = (
        "sandbox; default-src 'none'; img-src 'self' data:; style-src 'unsafe-inline'")
    return response
@app.get("/apps/{name}/{resource:path}")
def tenant_app_resource(name: str, resource: str):
    """Serve tenant apps in an opaque-origin sandbox with a restrictive CSP."""
    if not re.match(r"^[a-zA-Z0-9_-]{1,40}$", name):
        raise HTTPException(status_code=400, detail="非法应用名")
    relative = (resource or "index.html").replace("\\", "/").lstrip("/")
    if not relative or any(part in ("", ".", "..") for part in relative.split("/")):
        raise HTTPException(status_code=400, detail="非法应用资源路径")
    roots = [os.path.join(_user_apps_dir(), name), os.path.join(APPS_DIR, name)]
    path = None
    for root in roots:
        candidate = os.path.abspath(os.path.join(root, *relative.split("/")))
        if os.path.commonpath([candidate, os.path.abspath(root)]) == os.path.abspath(root) \
                and os.path.isfile(candidate):
            path = candidate
            break
    if not path:
        raise HTTPException(status_code=404, detail="应用资源不存在")
    response = FileResponse(path)
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["Content-Security-Policy"] = (
        "sandbox allow-scripts; default-src 'self' data: blob:; "
        "script-src 'self' 'unsafe-inline'; style-src 'self' 'unsafe-inline'; "
        "img-src 'self' data: blob:; connect-src 'none'; frame-src 'none'")
    return response

# 前端页面挂载在根路径：访问 / 就是打开 gui/static/index.html
# ============ 分布式任务集群 P1：/api/tasks（任务队列 + worker 池） ============

class TaskSubmitRequest(BaseModel):
    provider: str
    messages: list[dict]
    model: str = ""
    online: bool = False
    reasoning: str = ""
    project: str = "default"
    repo_path: str = ""
    isolate_git: bool = True


@app.get("/api/tasks")
def tasks_list(limit: int = 30):
    """任务列表（按创建时间倒序）"""
    return {"tasks": cluster._store.list(
        limit=min(max(limit, 1), 100), tenant=current_tenant.get())}


@app.get("/api/tasks/{tid}")
def task_get(tid: str):
    t = cluster._store.get(tid, tenant=current_tenant.get())
    if not t:
        raise HTTPException(status_code=404, detail="任务不存在")
    return t


@app.post("/api/tasks")
def task_submit(req: TaskSubmitRequest):
    """提交任务 → 入队 → worker 池自动执行"""
    if not req.messages:
        raise HTTPException(status_code=400, detail="messages 不能为空")
    tid = uuid.uuid4().hex[:12]
    worktree_path = ""
    branch = ""
    repo_path = req.repo_path.strip()
    if req.isolate_git and repo_path:
        try:
            isolated = cluster.WORKSPACES.create(tid, repo_path)
            repo_path = str(isolated.repo)
            worktree_path = str(isolated.path)
            branch = isolated.branch
        except Exception as exc:
            raise HTTPException(status_code=400, detail="无法创建任务 Git worktree: %s" % exc)
    cluster._store.submit({
        "id": tid, "provider": req.provider, "model": req.model,
        "messages": req.messages, "online": req.online, "reasoning": req.reasoning,
        "project": req.project, "repo_path": repo_path,
        "worktree_path": worktree_path, "branch": branch,
        "sandbox_mode": load_settings().get("sandbox_mode", "required"),
        "tenant": current_tenant.get(),
    })
    return {"id": tid, "status": "pending", "worktree_path": worktree_path, "branch": branch}


def _task_workspace(tid):
    task = cluster._store.get(tid, tenant=current_tenant.get())
    if not task:
        raise HTTPException(status_code=404, detail="任务不存在")
    if not task.get("worktree_path"):
        raise HTTPException(status_code=409, detail="该任务不是 Git 仓库任务，没有独立 worktree")
    item = cluster.WORKSPACES.get(tid)
    if item is None:
        try:
            item = cluster.WORKSPACES.attach(
                tid, task.get("repo_path"), task.get("worktree_path"), task.get("branch"))
        except Exception as exc:
            raise HTTPException(status_code=409, detail=str(exc))
    return item


@app.get("/api/tasks/{tid}/git/status")
def task_git_status(tid: str):
    _task_workspace(tid)
    return {"ok": True, "status": cluster.WORKSPACES.status(tid)}


@app.get("/api/tasks/{tid}/git/diff")
def task_git_diff(tid: str):
    _task_workspace(tid)
    return {"ok": True, "diff": cluster.WORKSPACES.diff(tid)}


@app.post("/api/tasks/{tid}/retry")
def task_retry(tid: str):
    t = cluster._store.get(tid, tenant=current_tenant.get())
    if not t:
        raise HTTPException(status_code=404, detail="任务不存在")
    cluster._store.retry(tid)
    return {"id": tid, "status": "pending"}


@app.post("/api/tasks/{tid}/cancel")
def task_cancel(tid: str):
    t = cluster._store.get(tid, tenant=current_tenant.get())
    if not t:
        raise HTTPException(status_code=404, detail="任务不存在")
    if t["status"] in ("pending", "running"):
        cluster._store.finish(tid, error="已取消")
    return {"id": tid, "status": "cancelled"}


@app.on_event("startup")
async def _start_cluster():
    """启动任务调度（worker 池后台泵）+ 自演化生命周期初始扫描"""
    # 历史跟随项目文件夹：把旧集中目录里的历史搬进各项目文件夹的 .wenmo/（幂等，失败跳过）
    try:
        moved = await asyncio.to_thread(history_store.migrate_history_to_projects)
        if moved:
            print("[history] 已迁移 %d 条历史到项目文件夹 .wenmo/" % moved, flush=True)
    except Exception as e:
        print("[history] 历史迁移跳过: %s" % e, flush=True)

    # 启动时清理残留 llama-server 进程（上次异常退出/停止失败遗留的孤儿进程）
    # 性能：to_thread 后台执行，不阻塞 startup 完成（tasklist 是同步 subprocess）
    try:
        killed = await asyncio.to_thread(_kill_llama_processes)
        if killed:
            print("[local] 启动清理残留 llama-server: %s" % killed, flush=True)
    except Exception as e:
        print("[local] 启动清理异常: %s" % e, flush=True)
    try:
        cluster.start_pump()
        print("[cluster] 任务集群已启动（worker 池 + 节点注册表）", flush=True)
    except Exception as e:
        print("[cluster] 启动失败: %s" % e, flush=True)
@app.on_event("shutdown")
async def _graceful_shutdown():
    """受控退出：flush 内存统计 + 停止本地推理 + 断开全部 MCP 子进程。
    （打包版窗口关闭走 os._exit 会跳过 atexit，但 uvicorn 的 shutdown 事件
    在正常退出路径（Ctrl+C / webview 关闭前先停 uvicorn）仍会触发。）"""
    # 1) 内存中未落盘的 usage_stats 增量 → 写盘（避免不足 10 次激活的统计丢失）
    try:
        _flush_usage_cache()
        print("[shutdown] usage_stats 已 flush", flush=True)
    except Exception as e:
        print("[shutdown] flush usage_stats 失败: %s" % e, flush=True)
    # 2) 停止本地 llama-server（受管进程 + 残留扫描双保险）
    try:
        await asyncio.to_thread(_local_stop)
    except Exception as e:
        print("[shutdown] 停止本地模型失败: %s" % e, flush=True)
    # 3) 断开全部 MCP 服务器（终止 stdio 子进程，避免残留孤儿进程）
    try:
        import mcp_client as _mc
        handles = list(getattr(_mc.manager, "_handles", {}).keys())
        for _name in handles:
            try:
                await _mc.run_on_loop(_mc.manager._disconnect(_name))
            except Exception:
                pass
        if handles:
            print("[shutdown] 已断开 %d 个 MCP 服务器" % len(handles), flush=True)
    except Exception as e:
        print("[shutdown] 断开 MCP 失败: %s" % e, flush=True)


# ============ 分布式 P2/P3：Worker API（远程 worker 拉任务/回传/注册/心跳） ============
_CLUSTER_TOKEN = os.environ.get("CLUSTER_TOKEN", "").strip()


def _check_worker_auth(request):
    if not _CLUSTER_TOKEN:
        raise HTTPException(status_code=503, detail="未配置 CLUSTER_TOKEN，远程 worker API 已关闭")
    tok = request.headers.get("X-Cluster-Token", "")
    if tok != _CLUSTER_TOKEN:
        raise HTTPException(status_code=401, detail="集群令牌无效")


@app.get("/api/workers/poll")
async def worker_poll(request: Request):
    """远程 worker 拉取一个 pending 任务（拉模式天然负载均衡）"""
    _check_worker_auth(request)
    wid = request.headers.get("X-Worker-Id", "remote-" + uuid.uuid4().hex[:6])
    cluster.register_node(wid, addr=request.client.host if request.client else "")
    cluster.heartbeat_node(wid)
    cluster._store.reset_stale()
    tid = cluster._store.next_pending()
    if not tid:
        return {"task": None}
    task = cluster._store.claim(tid, worker_id=wid)
    if not task:
        return {"task": None}
    return {"task": task}


@app.post("/api/workers/result")
async def worker_result(request: Request):
    """远程 worker 回传任务结果"""
    _check_worker_auth(request)
    body = await request.json()
    tid = body.get("id", "")
    result = body.get("result")
    error = body.get("error")
    wid = request.headers.get("X-Worker-Id", "")
    t = cluster._store.get(tid)
    if not t:
        raise HTTPException(status_code=404, detail="任务不存在")
    if error:
        cluster._store.finish(tid, error=str(error)[:500], worker_id=wid)
    else:
        cluster._store.finish(tid, result=result, worker_id=wid)
    cluster.heartbeat_node(wid)
    return {"ok": True}


@app.post("/api/workers/register")
async def worker_register(request: Request):
    _check_worker_auth(request)
    body = await request.json()
    wid = request.headers.get("X-Worker-Id", "node-" + uuid.uuid4().hex[:6])
    cluster.register_node(wid, name=body.get("name", ""), addr=request.client.host if request.client else "")
    return {"ok": True, "worker_id": wid}


@app.post("/api/workers/heartbeat")
async def worker_heartbeat(request: Request):
    _check_worker_auth(request)
    wid = request.headers.get("X-Worker-Id", "")
    cluster.heartbeat_node(wid)
    return {"ok": True}


@app.get("/api/workers/nodes")
def workers_list():
    """节点列表（集群状态）"""
    cluster.prune_nodes()
    return {"nodes": [{"worker_id": k, **v} for k, v in cluster.NODES.items()]}


app.mount("/", StaticFiles(directory=STATIC_DIR, html=True), name="static")

def _find_server_port(preferred=8000):
    """端口自适应：8000 被占用时自动 +1 找空闲端口，支持新旧实例共存"""
    import socket
    for port in range(preferred, preferred + 50):
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            try:
                s.bind(("127.0.0.1", port))
                return port
            except OSError:
                continue
    return preferred


def _wait_server(port, timeout=20):
    """等待服务器就绪"""
    url = f"http://127.0.0.1:{port}/"
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            with urllib.request.urlopen(url, timeout=1) as r:
                if r.status == 200:
                    return True
        except Exception:
            time.sleep(0.3)
    return False


if __name__ == "__main__":
    # ===== 单实例锁：按安装位置隔离（不同目录的问墨可共存，同一份只开一个） =====
    import ctypes as _ct
    import hashlib as _hl
    _lock_seed = os.path.abspath(getattr(sys, "_MEIPASS", os.path.dirname(os.path.abspath(__file__))))
    _lock_hash = _hl.md5(_lock_seed.encode("utf-8")).hexdigest()[:16]
    _lock_name = "WenMo_Single_Instance_Mutex_" + _lock_hash
    _mutex = _ct.windll.kernel32.CreateMutexW(None, False, _lock_name)
    # PID 文件按锁 seed 隔离：开发版与客户端（不同安装目录）各写各的，
    # 避免互相覆盖 lock.pid → 下次启动读到对方 PID 误杀进程
    _lock_dir = os.path.join(os.environ.get("APPDATA", os.path.expanduser("~")), "问墨")
    _pid_file = os.path.join(_lock_dir, "lock_%s.pid" % _lock_hash)
    if _ct.windll.kernel32.GetLastError() == 183:  # ERROR_ALREADY_EXISTS
        _still_alive = False
        _old_pid = 0
        try:
            if os.path.exists(_pid_file):
                _old_pid = int(open(_pid_file, "r").read().strip() or "0")
                if _old_pid > 0:
                    _h = _ct.windll.kernel32.OpenProcess(0x1000, False, _old_pid)  # PROCESS_QUERY_LIMITED_INFORMATION
                    if _h:
                        _ec = _ct.c_ulong(0)
                        _ct.windll.kernel32.GetExitCodeProcess(_h, _ct.byref(_ec))
                        _ct.windll.kernel32.CloseHandle(_h)
                        _still_alive = (_ec.value == 259)  # STILL_ACTIVE
        except Exception:
            _still_alive = False
        if _still_alive:
            _ret = _ct.windll.user32.MessageBoxW(None,
                "问墨正在运行中（可能窗口已关闭但进程残留在后台）！\n\n是否结束现有问墨进程并重新打开？\n\n选“是”=结束旧进程并启动新实例\n选“否”=退出本次启动",
                "问墨·code", 0x34)  # MB_YESNO | MB_ICONQUESTION
            if _ret == 6:  # IDYES
                try:
                    _h = _ct.windll.kernel32.OpenProcess(0x0001, False, _old_pid)  # PROCESS_TERMINATE
                    if _h:
                        _ct.windll.kernel32.TerminateProcess(_h, 0)
                        _ct.windll.kernel32.CloseHandle(_h)
                        time.sleep(0.8)
                except Exception:
                    pass
            else:
                sys.exit(0)
        # 旧进程已死（残留 PID/僵尸锁）→ 接管继续启动
    # 记录当前 PID（供下次启动校验残留进程）
    try:
        os.makedirs(os.path.dirname(_pid_file), exist_ok=True)
        with open(_pid_file, "w") as _f:
            _f.write(str(os.getpid()))
    except Exception:
        pass

    # ===== 端口自适应：8000 被占用自动顺延（支持不同目录的问墨共存） =====
    port = _find_server_port(8000)

    # uvicorn 服务器实例（可控制优雅退出：置 should_exit=True 触发 shutdown 事件）
    _server = None

    def _run_server():
        global _server
        try:
            import uvicorn
            cfg = uvicorn.Config(app, host="127.0.0.1", port=port, log_level="warning")
            _server = uvicorn.Server(cfg)
            _server.run()
        except Exception:
            pass

    t = threading.Thread(target=_run_server, daemon=True)
    t.start()
    _wait_server(port)

    def _graceful_exit():
        """受控退出：先触发 uvicorn shutdown（执行 @app.on_event shutdown 钩子：
        flush 统计 + 停本地模型 + 断开 MCP），再强制退出（pythonnet 可能阻止正常返回）。"""
        try:
            if _server is not None:
                _server.should_exit = True        # 触发 uvicorn shutdown → 执行 _graceful_shutdown
        except Exception:
            pass
        # 给 shutdown 钩子一点时间跑完（flush/断开），再兜底退出
        try:
            import time as _t
            _t.sleep(1.2)
        except Exception:
            pass
        try:
            os._exit(0)
        except Exception:
            pass

    # ===== 桌面版：pywebview 原生窗口（像 Codex 桌面应用一样） =====
    try:
        # WebView2 在窗口失焦/被遮挡时会暂停后台计时器与渲染，导致切屏后流式对话内容不更新。
        # 通过环境变量禁用后台节流，保证切屏后对话仍持续渲染。
        _wv2_args = os.environ.get("WEBVIEW2_ADDITIONAL_BROWSER_ARGUMENTS", "")
        for _flag in ("--disable-background-timer-throttling",
                      "--disable-backgrounding-occluded-windows",
                      "--disable-renderer-backgrounding"):
            if _flag not in _wv2_args:
                _wv2_args = (_wv2_args + " " + _flag).strip()
        os.environ["WEBVIEW2_ADDITIONAL_BROWSER_ARGUMENTS"] = _wv2_args
        import webview
        _win = webview.create_window("问墨·code", "http://127.0.0.1:%d" % port,
                                     width=1280, height=820, min_size=(960, 600))
        def _on_closed():
            """窗口关闭 → 受控退出（触发 shutdown 钩子清理，再彻底退出进程）"""
            _graceful_exit()
        try:
            _win.events.closed += _on_closed
        except Exception:
            pass
        webview.start()
        # start() 返回（窗口已全部关闭）：兜底受控退出（pythonnet 可能阻止正常退出）
        _graceful_exit()
    except Exception as _e:
        # 桌面窗口失败：记录日志 + 弹窗提示，最后才回退浏览器
        try:
            _log_dir = os.path.join(os.environ.get("APPDATA", os.path.expanduser("~")), "问墨")
            os.makedirs(_log_dir, exist_ok=True)
            with open(os.path.join(_log_dir, "webview_error.log"), "a", encoding="utf-8") as _f:
                _f.write("[" + time.strftime("%Y-%m-%d %H:%M:%S") + "] " + repr(_e) + "\n")
        except Exception:
            pass
        try:
            _ct.windll.user32.MessageBoxW(None,
                "桌面窗口启动失败，将改用浏览器打开。\n错误：%s\n\n详情已写入 %%APPDATA%%\\问墨\\webview_error.log" % _e,
                "问墨·code", 0x30)
        except Exception:
            pass
        import webbrowser
        try:
            webbrowser.open("http://127.0.0.1:%d" % port)
        except Exception:
            pass
        t.join()
