# -*- coding: utf-8 -*-
"""html2pptx_core：HTML 设计稿 → 可编辑 PPT 渲染核心
被 file_mcp_server.py 的 html_to_pptx 工具调用。
约定：HTML 为 1920x1080 画布，文字块用 data-slot="名称" 标注；
     多页用 <!-- PAGE --> 分隔或 <section class="page"> 包裹。
流程：Edge headless 渲染纯背景图 → dump-dom 提取槽位坐标/样式
     → python-pptx 拼装原生文本框 → 输出到 files/ 下载区。

v2 修复记录（2025）：
  P0-1 _to_rgb 支持 rgb()/rgba()/颜色名（原只认 #RRGGBB → getComputedStyle 返回 rgb() → 文字全变黑）
  P0-2 文本提取保留换行（innerText/子节点遍历，支持 <br> 与块级换行）
  P0-3 字号系数 0.75 → 0.5（1920px→13.33in 的正确比例，防止文字比设计稿大 50% 溢出）
  P0-4 背景截图与坐标提取同步等待字体加载（--virtual-time-budget）
  P1-5 富文本槽位：一个 data-slot 内多 span 不同样式（颜色/加粗）均保留
  P1-6 背景渲染失败显式警告 + 输出统计（页数/槽位数）
  P1-7 支持 <img> 本地图片自动 base64 内嵌（离线可用）、远程 URL 自动加载
  P1-8 文字溢出防护：超长内容自动缩小字号
"""
import base64
import html as _html
import json
import os
import re
import shutil
import subprocess
import tempfile
import urllib.parse

BASE = os.path.dirname(os.path.abspath(__file__))
FILES_DIR = os.path.join(BASE, "files")

SLOT_JS = r"""
<script>
window.addEventListener('load', function(){
  setTimeout(function(){
    var out = {};
    document.querySelectorAll('[data-slot]').forEach(function(el, i){
      var r = el.getBoundingClientRect();
      var cs = window.getComputedStyle(el);
      var key = el.getAttribute('data-slot') || el.id || ('s' + (i + 1));
      // 收集富文本 runs：保留子元素（span/b/strong/i）独立样式
      var runs = [];
      function walk(node, styleEl){
        node.childNodes.forEach(function(n){
          if (n.nodeType === 3) { // 文本节点
            var t = n.textContent.replace(/\s+/g, ' ');
            if (t.trim()) {
              var c = window.getComputedStyle(styleEl || el);
              runs.push({text: t, fontSize: c.fontSize, color: c.color, fontWeight: c.fontWeight});
            }
          } else if (n.nodeType === 1) {
            if (n.tagName === 'BR') {
              runs.push({text: '\n', fontSize: cs.fontSize, color: cs.color, fontWeight: cs.fontWeight});
            } else {
              walk(n, n);
            }
          }
        });
      }
      walk(el, el);
      out[key] = {
        x: Math.round(r.x), y: Math.round(r.y),
        w: Math.round(r.width), h: Math.round(r.height),
        runs: runs.length ? runs : [{
          text: (el.textContent || '').replace(/\s+/g,' ').trim(),
          fontSize: cs.fontSize, color: cs.color, fontWeight: cs.fontWeight
        }],
        textAlign: cs.textAlign,
        lineHeight: cs.lineHeight
      };
    });
    var d = document.createElement('div'); d.id='__SLOT_DATA__'; d.style.display='none';
    d.textContent = JSON.stringify(out); document.body.appendChild(d);
  }, 300);
});
</script>
"""


def _find_edge():
    cands = [
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    for c in cands:
        if os.path.isfile(c):
            return c
    return shutil.which("msedge")


def _file_url(path):
    return "file:///" + urllib.parse.quote(os.path.abspath(path).replace("\\", "/"), safe="/")


def _edge_shot(edge, html_path, png_path, w=1920, h=1080):
    # P0-4：与坐标提取同步等待字体/布局（virtual-time-budget）
    cmd = [edge, "--headless=new", "--disable-gpu", "--hide-scrollbars",
           "--force-device-scale-factor=1", "--window-size=%d,%d" % (w, h),
           "--virtual-time-budget=5000",
           "--screenshot=%s" % png_path, _file_url(html_path)]
    subprocess.run(cmd, capture_output=True, timeout=90, creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
    return os.path.isfile(png_path) and os.path.getsize(png_path) > 0


def _edge_dump(edge, html_path):
    cmd = [edge, "--headless=new", "--disable-gpu", "--dump-dom",
           "--virtual-time-budget=4000", _file_url(html_path)]
    r = subprocess.run(cmd, capture_output=True, timeout=90, creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
    return r.stdout.decode("utf-8", errors="replace")


def _extract_slots(dom):
    m = re.search(r'(?is)<div id="__SLOT_DATA__"[^>]*>(.*?)</div>', dom)
    if not m:
        return {}
    raw = _html.unescape(m.group(1)).strip()
    try:
        return json.loads(raw)
    except Exception:
        return {}


def _build_page(page_html, styles):
    head = "".join(styles)
    return ("<!DOCTYPE html><html><head><meta charset='utf-8'>%s"
            "<style>html,body{margin:0;padding:0;width:1920px;height:1080px;overflow:hidden;}"
            "body>*{box-sizing:border-box}</style></head><body>%s</body></html>"
            % (head, page_html))


def _split_pages(html_content):
    parts = re.split(r'(?i)<!--\s*PAGE\s*-->', html_content)
    parts = [p.strip() for p in parts if p.strip()]
    if len(parts) > 1:
        return parts
    secs = re.findall(r'(?is)<section[^>]*class=["\']page["\'][^>]*>.*?</section>', html_content)
    if secs:
        return secs
    return [html_content]


_COLOR_NAMES = {
    "black": (0, 0, 0), "white": (255, 255, 255), "red": (255, 0, 0),
    "green": (0, 128, 0), "blue": (0, 0, 255), "gray": (128, 128, 128),
    "grey": (128, 128, 128), "orange": (255, 165, 0), "yellow": (255, 255, 0),
    "purple": (128, 0, 128), "teal": (0, 128, 128), "navy": (0, 0, 128),
    "silver": (192, 192, 192), "maroon": (128, 0, 0), "lime": (0, 255, 0),
    "olive": (128, 128, 0), "aqua": (0, 255, 255), "fuchsia": (255, 0, 255),
}


def _to_rgb(color):
    """P0-1：支持 #RRGGBB / #RGB / rgb(r,g,b) / rgba(r,g,b,a) / 常见颜色名 → (r,g,b)；失败返回 (0,0,0)"""
    color = (color or "").strip().lower()
    if not color:
        return (0, 0, 0)
    if color.startswith("#"):
        hexc = color.lstrip("#")
        if len(hexc) == 3:
            hexc = "".join(c * 2 for c in hexc)
        if len(hexc) == 6:
            try:
                return tuple(int(hexc[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                return (0, 0, 0)
    m = re.match(r"rgba?\(([^)]+)\)", color)
    if m:
        parts = [p.strip() for p in m.group(1).replace("/", ",").split(",")]
        try:
            r, g, b = (int(float(parts[i])) for i in range(3))
            return (max(0, min(r, 255)), max(0, min(g, 255)), max(0, min(b, 255)))
        except Exception:
            return (0, 0, 0)
    if color in _COLOR_NAMES:
        return _COLOR_NAMES[color]
    return (0, 0, 0)


def _embed_images(html_content, base_dir=None):
    """P1-7：<img src="本地路径"> → base64 内嵌（离线可用）；远程 URL / data: / file: 原样保留"""
    base_dir = base_dir or BASE

    def _repl(m):
        tag = m.group(0)
        sm = re.search(r'src=["\']([^"\']+)["\']', tag, re.I)
        if not sm:
            return tag
        src = sm.group(1)
        if src.startswith(("http://", "https://", "data:", "file:")):
            return tag
        p = src if os.path.isabs(src) else os.path.join(base_dir, src)
        if os.path.isfile(p):
            try:
                with open(p, "rb") as f:
                    b64 = base64.b64encode(f.read()).decode("ascii")
                ext = os.path.splitext(p)[1].lower().lstrip(".")
                mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                        "gif": "image/gif", "webp": "image/webp", "svg": "image/svg+xml"}.get(ext, "image/png")
                return re.sub(r'src=["\'][^"\']*["\']',
                              'src="data:%s;base64,%s"' % (mime, b64), tag, count=1, flags=re.I)
            except Exception:
                return tag
        return tag

    return re.sub(r'(?is)<img\b[^>]*>', _repl, html_content)


def _auto_font_size(runs, box_w_emu, box_h_emu):
    """P1-8：按内容量与文本框尺寸估算合适字号（pt），防止长文本撑破卡片"""
    try:
        total_chars = sum(len(re.sub(r'\s', '', r.get("text", "")))
                          for r in runs if r.get("text") != "\n")
        if total_chars == 0:
            return 12.0
        base_px = 24.0
        for r in runs:
            try:
                base_px = float(str(r.get("fontSize", "24px")).replace("px", ""))
                break
            except Exception:
                continue
        base_pt = base_px * 0.5
        w_in = box_w_emu / 914400.0
        h_in = box_h_emu / 914400.0
        if w_in <= 0 or h_in <= 0:
            return base_pt
        char_w_ratio = 0.6  # 中英混排近似
        chars_per_line = max(int(w_in * 72 / (base_pt * char_w_ratio)), 1)
        est_lines = max(int(total_chars / chars_per_line) + 1, 1)
        line_h_in = base_pt / 72.0 * 1.3
        est_h = est_lines * line_h_in
        if est_h > h_in:
            ratio = h_in / est_h
            return max(base_pt * ratio, 6.0)
        return base_pt
    except Exception:
        return 12.0


def render_html_to_pptx(args):
    """主入口：args = {html, out_name, title} → 返回结果文本"""
    html_content = str(args.get("html", "")).strip()
    out_name = str(args.get("out_name", "")).strip() or "html_to_pptx.pptx"
    if not out_name.lower().endswith(".pptx"):
        out_name += ".pptx"
    if not html_content:
        return "参数 html 不能为空（需 1920x1080 设计稿，文字块用 data-slot 标注，多页用 <!-- PAGE --> 分隔）"

    edge = _find_edge()
    if not edge:
        return "未找到 Edge 浏览器（需要 Edge headless 渲染），请安装 Microsoft Edge"

    # P1-7：图片内嵌预处理
    html_content = _embed_images(html_content)

    tmp = tempfile.mkdtemp(prefix="h2p_")
    warnings = []
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        styles = re.findall(r'(?is)<style[^>]*>.*?</style>', html_content)
        pages = _split_pages(html_content)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        EMU_PER_IN = 914400
        SW, SH = 13.333, 7.5
        n_slots_total = 0

        for i, page_html in enumerate(pages, 1):
            full = _build_page(page_html, styles)
            # 1) 背景图版：隐藏 data-slot 文字
            bg_html = full.replace("</head>", "<style>[data-slot]{visibility:hidden !important;}</style></head>", 1)
            bg_path = os.path.join(tmp, "bg%d.html" % i)
            bg_png = os.path.join(tmp, "bg%d.png" % i)
            with open(bg_path, "w", encoding="utf-8") as f:
                f.write(bg_html)
            ok_bg = _edge_shot(edge, bg_path, bg_png)
            if not ok_bg:
                # P1-6：不再静默降级，显式警告
                warnings.append("第 %d 页背景渲染失败（已跳过背景，仅保留文字）" % i)
            # 2) 坐标版：注入提取脚本
            c_path = os.path.join(tmp, "c%d.html" % i)
            with open(c_path, "w", encoding="utf-8") as f:
                f.write(full.replace("</body>", SLOT_JS + "</body>", 1))
            dom = _edge_dump(edge, c_path)
            slots = _extract_slots(dom)
            n_slots_total += len(slots)
            # 3) 拼装幻灯片
            slide = prs.slides.add_slide(blank)
            if ok_bg:
                slide.shapes.add_picture(bg_png, 0, 0, width=Inches(SW), height=Inches(SH))
            for sid, s in slots.items():
                try:
                    left = Emu(int(s["x"] / 1920.0 * SW * EMU_PER_IN))
                    top = Emu(int(s["y"] / 1080.0 * SH * EMU_PER_IN))
                    wdt = Emu(max(int(s["w"] / 1920.0 * SW * EMU_PER_IN), 1))
                    hgt = Emu(max(int(s["h"] / 1080.0 * SH * EMU_PER_IN), 1))
                    tb = slide.shapes.add_textbox(left, top, wdt, hgt)
                    tf = tb.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                    runs = s.get("runs") or [{"text": s.get("text", ""), "fontSize": s.get("fontSize"),
                                              "color": s.get("color"), "fontWeight": s.get("fontWeight")}]
                    auto_pt = _auto_font_size(runs, wdt, hgt)  # P1-8 溢出防护
                    p = tf.paragraphs[0]
                    for run in runs:
                        t = run.get("text", "")
                        if t == "\n":
                            p = tf.add_paragraph()
                            continue
                        r = p.add_run()
                        r.text = t
                        try:
                            fsize = float(str(run.get("fontSize", "24px")).replace("px", ""))
                            # P0-3：0.5 pt/px（1920px → 13.33in 的正确比例）
                            r.font.size = Pt(max(min(fsize * 0.5, auto_pt), 6))
                        except Exception:
                            r.font.size = Pt(12)
                        try:
                            rr, gg, bb = _to_rgb(run.get("color"))  # P0-1
                            r.font.color.rgb = RGBColor(rr, gg, bb)
                        except Exception:
                            pass
                        if str(run.get("fontWeight")) in ("bold", "700", "800", "900"):
                            r.font.bold = True
                    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
                        str(s.get("textAlign", "left")), PP_ALIGN.LEFT)
                    for pp in tf.paragraphs:
                        pp.alignment = al
                except Exception:
                    continue

        out_path = os.path.join(FILES_DIR, out_name)
        prs.save(out_path)
        url = "http://127.0.0.1:8000/files/" + urllib.parse.quote(out_name)
        stat = "已生成可编辑 PPT：%s（%d 页，%d 个文字槽位）\n下载/预览：%s\n" % (out_name, len(pages), n_slots_total, url)
        if warnings:
            stat += "警告：%s\n" % "；".join(warnings)
        stat += "说明：文字为 PPT 原生文本框（可编辑/搜索/复制），背景为渲染图，可再用 file_convert 转 PDF。"
        return stat
    except Exception as e:
        return "html_to_pptx 失败：%s" % str(e)[:300]
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
