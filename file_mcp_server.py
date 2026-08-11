# -*- coding: utf-8 -*-
"""文件 MCP 服务器：统一文件操作（读取/写入/打包/压缩/列举）。
对标用户需求：文件生成、打包、压缩、读取等操作统一由本 MCP 管理——
所有生成的文件强制落到 files/ 下载区（保证预览/下载链接可用），
杜绝"脚本生成到临时路径 → 链接 404"的问题。

工具：
  file_read   读取文件（纯文本 + Word/Excel/PPT/PDF 转文本）
  file_write  写文本文件到下载区，返回预览/下载链接
  file_zip    把多个文件打包为 zip（下载区）
  file_list   列举工作区/下载区文件

启动：由 gui_server.py 的 MCP 管理器自动拉起（注册于 mcp.json）。
"""

import html
import json
import os
import re
import sys
import urllib.parse

BASE = os.path.dirname(os.path.abspath(__file__))
FILES_DIR = os.path.join(BASE, "files")
WORKSPACE = os.path.join(BASE, "workspace")
for _d in (FILES_DIR, WORKSPACE):
    os.makedirs(_d, exist_ok=True)

from mcp.server import Server  # noqa: E402
from mcp.server.stdio import stdio_server  # noqa: E402
from mcp.types import CallToolRequestParams, CallToolResult, ListToolsRequest, ListToolsResult, TextContent, Tool  # noqa: E402

server = Server("file-mcp")


def _safe_dir(path):
    """路径必须落在 workspace 或 files 内（沙箱）"""
    p = os.path.abspath(path)
    for root in (WORKSPACE, FILES_DIR):
        try:
            if os.path.commonpath([p, root]) == root:
                return p
        except Exception:
            continue
    return None


def _read_any(path):
    """读取任意文件为文本（文本直接读；办公文档转文本）"""
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    if ext in ("txt", "md", "markdown", "json", "csv", "yaml", "py", "js", "ts", "html", "htm", "css"):
        with open(path, encoding="utf-8", errors="replace") as f:
            return f.read()
    try:
        if ext in ("docx",):
            import docx
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs if p.text.strip())
        if ext in ("pptx", "ppt"):
            import pptx
            prs = pptx.Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append("【第%d页】" % i)
                for sh in slide.shapes:
                    if sh.has_text_frame:
                        for para in sh.text_frame.paragraphs:
                            t = "".join(r.text for r in para.runs).strip()
                            if t:
                                parts.append(t)
            return "\n".join(parts)
        if ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append("【工作表:%s】" % ws.title)
                for row in ws.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    if any(vals):
                        parts.append(" | ".join(vals))
            return "\n".join(parts)
        if ext in ("pdf",):
            import fitz
            doc = fitz.open(path)
            # 完整读取全部页（不限制页数——用户要求文档读取不截断）
            return "\n".join(doc[i].get_text() for i in range(len(doc)))
    except Exception as e:
        return "（无法解析 %s 文档：%s）" % (ext, e)
    return "（不支持的格式：.%s）" % ext


def _tool_result(text):
    return CallToolResult(content=[TextContent(type="text", text=str(text))])


def handle_file_read(args):
    path = str(args.get("path", "")).strip()
    if not path:
        return "错误：需要 path 参数"
    p = os.path.abspath(path) if os.path.isabs(path) else os.path.join(WORKSPACE, path)
    if not os.path.isfile(p):
        alt = os.path.join(FILES_DIR, os.path.basename(path))
        if os.path.isfile(alt):
            p = alt
        else:
            return "错误：文件不存在 %s" % path
    content = _read_any(p)
    # 完整返回（不截断——用户要求文档读取不设长度限制，避免影响逐词校对等场景）
    return "文件: %s\n%s" % (os.path.basename(p), content)


def handle_file_write(args):
    path = str(args.get("path", "")).strip()
    content = str(args.get("content", ""))
    if not path:
        return "错误：需要 path 参数"
    name = os.path.basename(path) or ("file_%d.txt" % int(__import__("time").time()))
    fpath = os.path.join(FILES_DIR, name)
    with open(fpath, "w", encoding="utf-8") as f:
        f.write(content)
    url = "/files/" + urllib.parse.quote(name)
    size = os.path.getsize(fpath)
    return "已写入下载区：%s（%d 字节）\n预览/下载：%s\n（模型请在回答中引用该链接）" % (name, size, url)


def handle_file_zip(args):
    zip_name = str(args.get("zip_name", "打包文件")).strip() or "打包文件"
    if not zip_name.lower().endswith(".zip"):
        zip_name += ".zip"
    sources = args.get("files") or []
    if not sources or not isinstance(sources, list):
        return "错误：需要 files 参数（要打包的文件路径列表）"
    import zipfile
    zpath = os.path.join(FILES_DIR, zip_name)
    with zipfile.ZipFile(zpath, "w", zipfile.ZIP_DEFLATED) as zf:
        for src in sources:
            p = os.path.abspath(src) if os.path.isabs(str(src)) else os.path.join(WORKSPACE, str(src))
            if not os.path.isfile(p):
                alt = os.path.join(FILES_DIR, os.path.basename(str(src)))
                if os.path.isfile(alt):
                    p = alt
                else:
                    continue
            zf.write(p, arcname=os.path.basename(p))
    url = "/files/" + urllib.parse.quote(zip_name)
    return "已打包 %d 个文件 → %s（%d 字节）\n下载：%s" % (len(sources), zip_name, os.path.getsize(zpath), url)


def handle_file_list(args):
    path = str(args.get("path", "")).strip() or WORKSPACE
    p = os.path.abspath(path) if os.path.isabs(path) else os.path.join(WORKSPACE, path)
    if not os.path.isdir(p):
        return "错误：目录不存在 %s" % path
    entries = sorted(os.listdir(p))
    lines = []
    for e in entries:
        full = os.path.join(p, e)
        if os.path.isdir(full):
            lines.append("📁 %s/" % e)
        else:
            lines.append("📄 %s (%d B)" % (e, os.path.getsize(full)))
    return "目录 %s | %d 项：\n%s" % (p, len(entries), "\n".join(lines) if lines else "（空）")


def _safe_path(path, default_dir):
    """解析为工作区/下载区内路径（沙箱）"""
    p = os.path.abspath(path) if os.path.isabs(str(path)) else os.path.join(default_dir, str(path))
    for root in (WORKSPACE, FILES_DIR):
        try:
            if os.path.commonpath([p, root]) == root:
                return p
        except Exception:
            continue
    raise ValueError("路径不在允许的工作区内（仅限项目目录与 files 下载区）")


def handle_file_move(args):
    src = str(args.get("source", "")).strip()
    dst = str(args.get("target", "")).strip()
    if not src or not dst:
        return "错误：需要 source 和 target 参数"
    try:
        sp = _safe_path(src, WORKSPACE)
        dp = os.path.abspath(dst) if os.path.isabs(dst) else os.path.join(FILES_DIR, os.path.basename(dst))
        if not os.path.exists(sp):
            return "错误：源不存在 %s" % src
        if os.path.isdir(dp):
            dp = os.path.join(dp, os.path.basename(sp))
        os.makedirs(os.path.dirname(dp), exist_ok=True)
        os.rename(sp, dp)
        return "已移动/重命名：%s → %s" % (sp, dp)
    except ValueError as ve:
        return "错误：%s" % ve
    except Exception as e:
        return "移动失败：%s" % str(e)[:150]


def handle_file_delete(args):
    path = str(args.get("path", "")).strip()
    if not path:
        return "错误：需要 path 参数"
    try:
        p = _safe_path(path, WORKSPACE)
        if os.path.isdir(p):
            import shutil
            shutil.rmtree(p)
        elif os.path.isfile(p):
            os.remove(p)
        else:
            return "错误：路径不存在 %s" % path
        return "已删除：%s" % p
    except ValueError as ve:
        return "错误：%s" % ve
    except Exception as e:
        return "删除失败：%s" % str(e)[:150]


def handle_file_convert(args):
    """Office 文档 → PDF（系统 Office 渲染）：docx/pptx/xlsx 转 PDF，返回下载链接。
    用于：PPT/Word/Excel 查看器预览、导出为 PDF。"""
    path = str(args.get("path", "")).strip()
    if not path:
        return "错误：需要 path 参数（文件路径）"
    name = os.path.basename(path)
    ext = os.path.splitext(name)[1].lower().lstrip(".")
    if ext not in ("docx", "doc", "pptx", "ppt", "xlsx", "xls"):
        return "错误：仅支持 Office 文档（docx/pptx/xlsx）转 PDF"
    p = os.path.join(FILES_DIR, name) if os.path.isfile(os.path.join(FILES_DIR, name)) else (
        os.path.abspath(path) if os.path.isfile(os.path.abspath(path)) else None)
    if not p:
        return "错误：文件不存在 %s" % path
    import urllib.parse
    try:
        # 复用服务器预览转换端点（Office COM 渲染）
        import urllib.request
        url = "/api/files/preview/" + urllib.parse.quote(name)
        req = urllib.request.Request(url)
        with urllib.request.urlopen(req, timeout=120) as resp:
            if resp.status != 200:
                return "转换失败：HTTP %d" % resp.status
            pdf_name = os.path.splitext(name)[0] + ".pdf"
            pdf_path = os.path.join(FILES_DIR, pdf_name)
            with open(pdf_path, "wb") as f:
                f.write(resp.read())
        link = "/files/" + urllib.parse.quote(pdf_name)
        return "已转换 %s → PDF（%s 字节）\n预览/下载：%s" % (name, os.path.getsize(pdf_path), link)
    except ValueError as ve:
        return "错误：%s" % ve
    except Exception as e:
        return "转换失败：%s" % str(e)[:200]


TOOLS = [
    Tool(
        name="file_read",
        description="读取文件内容（纯文本 + Word/Excel/PPT/PDF 自动转文本）。参数 path=文件路径（工作区相对或绝对）。",
        input_schema={"type": "object", "properties": {"path": {"type": "string"}}, "required": ["path"]},
    ),
    Tool(
        name="file_write",
        description="把文本内容写入文件（保存到下载区 files/，返回预览/下载链接，回答中引用）。"
                    "参数 path=文件名；content=文本内容。办公文档（docx/pptx/xlsx/pdf）请用专门的生成工具。",
        input_schema={"type": "object", "properties": {
            "path": {"type": "string"}, "content": {"type": "string"}}, "required": ["path", "content"]},
    ),
    Tool(
        name="file_zip",
        description="把多个文件打包为 zip（保存到下载区，返回下载链接）。参数 zip_name=压缩包名；files=文件路径列表。",
        input_schema={"type": "object", "properties": {
            "zip_name": {"type": "string"}, "files": {"type": "array", "items": {"type": "string"}}},
            "required": ["zip_name", "files"]},
    ),
    Tool(
        name="file_list",
        description="列举工作区/下载区目录下的文件。参数 path=目录路径（可选，默认工作区）。",
        input_schema={"type": "object", "properties": {"path": {"type": "string"}}},
    ),
    Tool(
        name="file_move",
        description="移动/重命名文件（仅限工作区与下载区内，沙箱）。参数 source=源路径；target=目标路径（文件名）。",
        input_schema={"type": "object", "properties": {
            "source": {"type": "string"}, "target": {"type": "string"}}, "required": ["source", "target"]},
    ),
    Tool(
        name="file_delete",
        description="删除文件/目录（仅限工作区与下载区内，沙箱；拒绝工作区外路径）。参数 path=路径。",
        input_schema={"type": "object", "properties": {"path": {"type": "string"}}, "required": ["path"]},
    ),
    Tool(
        name="file_convert",
        description="把 Office 文档转成 PDF（系统 Office 渲染，排版/配色/公式原样）：docx/pptx/xlsx → PDF，"
                    "返回 PDF 预览/下载链接。用于：PPT/Word/Excel 查看、导出 PDF。参数 path=文件路径。",
        input_schema={"type": "object", "properties": {"path": {"type": "string"}}, "required": ["path"]},
    ),
    Tool(
        name="html_to_pptx",
        description="把 HTML 设计稿（1920x1080，文字块用 data-slot=\"名称\" 标注）渲染为可编辑 PPT。"
                    "多页用 <!-- PAGE --> 分隔。内部：Edge headless 渲染纯背景图 + 提取槽位坐标/样式 "
                    "+ python-pptx 拼装原生文本框（文字可编辑/搜索/复制）→ 输出到 files/ 下载区返回链接。"
                    "参数：html=完整 HTML 设计稿；out_name=输出文件名（可选，默认 html_to_pptx.pptx）。"
                    "生成后可再用 file_convert 转 PDF。",
        input_schema={"type": "object", "properties": {
            "html": {"type": "string", "description": "完整 HTML 设计稿，1920x1080，文字块用 data-slot 标注，多页用 <!-- PAGE --> 分隔"},
            "out_name": {"type": "string", "description": "输出 pptx 文件名（可选）"}},
            "required": ["html"]},
    ),]


async def handle_list_tools(*args, **kwargs) -> ListToolsResult:
    return ListToolsResult(tools=TOOLS)


async def handle_call_tool(*args, **kwargs) -> CallToolResult:
    params = kwargs.get("params")
    if params is None:
        for a in args:
            if isinstance(a, CallToolRequestParams):
                params = a
                break
    if params is None or params.name not in ("file_read", "file_write", "file_zip", "file_list", "file_move", "file_delete", "file_convert", "html_to_pptx"):
        return _tool_result("未知工具")
    args_dict = params.arguments or {}
    try:
        if params.name == "file_read":
            return _tool_result(handle_file_read(args_dict))
        if params.name == "file_write":
            return _tool_result(handle_file_write(args_dict))
        if params.name == "file_zip":
            return _tool_result(handle_file_zip(args_dict))
        if params.name == "file_move":
            return _tool_result(handle_file_move(args_dict))
        if params.name == "file_delete":
            return _tool_result(handle_file_delete(args_dict))
        if params.name == "file_convert":
            return _tool_result(handle_file_convert(args_dict))
        if params.name == "html_to_pptx":
            return _tool_result(handle_html_to_pptx(args_dict))
        return _tool_result(handle_file_list(args_dict))
    except Exception as e:
        return _tool_result("文件操作失败：%s" % str(e)[:200])



def handle_html_to_pptx(args):
    """HTML 设计稿 → 可编辑 PPT（复用 html2pptx_core 渲染管线）"""
    from html2pptx_core import render_html_to_pptx
    return render_html_to_pptx(args)


server.add_request_handler("tools/list", ListToolsRequest, handle_list_tools)
server.add_request_handler("tools/call", CallToolRequestParams, handle_call_tool)


async def main():
    async with stdio_server() as (read, write):
        await server.run(read, write, server.create_initialization_options())


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
